import base64
import ipaddress
import json
import logging
import os
import random
import re
import shutil
import socket
import subprocess
import threading
import time
import uuid
import webbrowser
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Optional
from urllib.parse import parse_qs, quote, quote_plus, unquote_plus, urlparse

import requests

from automation import browser_driver, desktop_uia, interaction_intelligence, screen_ocr
from automation.app_profiles.action_merge import (
    normalize_actions_with_merge_rules,
)
from automation.app_profiles.prompt_fragments import load_planner_fragment
from config import MODEL_POOL
from llm.volcengine_llm import VolcengineLLM, _normalize_reasoning_effort
from memory.auto_memory import AutoMemoryManager
from memory.memory_system import LongTermMemory, MidTermMemory, ShortTermMemory
from runtime.permissions import PermissionModel, SAFE_ACTION_TYPES
from runtime.shell_danger import shell_command_blocked_reason
from runtime.timing_breakdown import compute_timing_breakdown

logger = logging.getLogger(__name__)

_OFFICE_MAX_CONTENT_CHARS = 400_000
_OFFICE_MAX_ROWS = 2000
_OFFICE_MAX_COLS = 64
_OFFICE_MAX_PPT_BULLETS = 100

# desktop_open_app：Windows 下扫描快捷方式/可执行文件的上限，避免开始菜单过大拖慢执行
_WIN_APP_SCAN_CAP = 8000

# 聊天窗通常不渲染 LaTeX：先写人话与分步，再视需要补充严谨式子
# 这些 Agent 主要产出 JSON 或内部结构化结果，流式推送到聊天窗观感差，默认不参与 SSE 文本流
_LLM_STREAM_BLOCK_AGENTS = frozenset(
    {
        "TaskParser",
        "MethodSaver",
        "SolutionLearner",
        "TaskSplitter",
        "QualityChecker",
        "ReActAgent",
        "MediaSummarizer",
    }
)

_MATH_NOTATION_FOR_CHAT = (
    "数学与公式：默认在普通聊天里展示。先用中文分步、换行、普通数字与单位把思路和结论写清楚；"
    "需要时用常见 Unicode（如 ² ³ × ÷）或「a/b」「根号下…」等口语化写法，让读者不看符号约定也能读懂。"
    "若仍需严谨记法，在直观说明之后单独用一行「参考记法：」写出；避免正文大段只有 $...$、\\boxed、\\mathrm 等 LaTeX 源码（除非用户明确要求交 LaTeX 源码）。"
)


def _windows_desktop_path_bonus(p: Path) -> int:
    try:
        s = str(p.resolve()).lower()
    except OSError:
        s = str(p).lower()
    if "desktop" in s or "桌面" in s:
        return 120
    return 0


def _windows_open_app_keywords(query: str) -> list[str]:
    """从用户/规划器给出的名称扩展出用于匹配快捷方式、exe 的关键词。"""
    q0 = (query or "").strip().strip('"').strip("'")
    if not q0:
        return []
    ql = q0.lower()
    parts: list[str] = [q0]
    if ql != q0:
        parts.append(ql)
    if "微" in q0 or "wechat" in ql or "weixin" in ql or "tencent" in ql:
        parts.extend(["wechat", "微信", "WeChat", "Weixin"])
    if "wps" in ql or "kingsoft" in ql or "金山" in q0:
        parts.extend(["wps", "WPS", "kingsoft", "金山", "WPS Office"])
    seen: set[str] = set()
    out: list[str] = []
    for x in parts:
        t = x.strip()
        if len(t) >= 1 and t not in seen:
            seen.add(t)
            out.append(t)
    return out


def _windows_score_app_match(keywords: list[str], path: Path) -> int:
    stem = path.stem
    name = path.name
    stem_l = stem.lower()
    name_l = name.lower()
    best = 0
    for kw in keywords:
        k = kw.strip()
        if len(k) < 1:
            continue
        kl = k.lower()
        if k == stem or k == name or kl == stem_l or kl == name_l:
            best = max(best, 1000)
        elif len(k) >= 2 and (k in stem or k in name or kl in stem_l or kl in name_l):
            best = max(best, 880)
        elif len(kl) >= 3 and (kl in stem_l.replace(" ", "") or kl in name_l.replace(" ", "")):
            best = max(best, 720)
    return best


def _windows_known_office_exes() -> list[Path]:
    env = os.environ
    pf = env.get("ProgramFiles", r"C:\Program Files")
    pf86 = env.get("ProgramFiles(x86)", r"C:\Program Files (x86)")
    la = env.get("LOCALAPPDATA", "")
    return [
        Path(pf) / "Tencent" / "WeChat" / "WeChat.exe",
        Path(pf86) / "Tencent" / "WeChat" / "WeChat.exe",
        Path(la) / "Kingsoft" / "WPS Office" / "ksolaunch.exe",
        Path(pf) / "Kingsoft" / "WPS Office" / "office6" / "wps.exe",
        Path(pf86) / "Kingsoft" / "WPS Office" / "office6" / "wps.exe",
    ]


def _windows_collect_shortcuts_and_exes() -> list[Path]:
    home = Path.home()
    env = os.environ
    shallow_roots: list[Path] = []
    for d in (
        home / "Desktop",
        home / "OneDrive" / "Desktop",
        home / "OneDrive" / "桌面",
        Path(env.get("PUBLIC", "")) / "Desktop",
    ):
        if d.is_dir():
            shallow_roots.append(d)

    found: list[Path] = []
    seen: set[str] = set()

    def push(p: Path) -> bool:
        if not p.is_file():
            return True
        try:
            key = str(p.resolve())
        except OSError:
            key = str(p)
        if key in seen:
            return True
        seen.add(key)
        found.append(p)
        return len(found) < _WIN_APP_SCAN_CAP

    for base in shallow_roots:
        if len(found) >= _WIN_APP_SCAN_CAP:
            break
        try:
            for pattern in ("*.lnk", "*.exe"):
                for p in base.glob(pattern):
                    if not push(p):
                        return found
            for sub in base.iterdir():
                if len(found) >= _WIN_APP_SCAN_CAP:
                    break
                if not sub.is_dir():
                    continue
                try:
                    for pattern in ("*.lnk", "*.exe"):
                        for p in sub.glob(pattern):
                            if not push(p):
                                return found
                except OSError:
                    continue
        except OSError:
            continue

    for base in (
        Path(env.get("APPDATA", "")) / "Microsoft" / "Windows" / "Start Menu" / "Programs",
        Path(env.get("PROGRAMDATA", "")) / "Microsoft" / "Windows" / "Start Menu" / "Programs",
    ):
        if len(found) >= _WIN_APP_SCAN_CAP:
            break
        if not base.is_dir():
            continue
        try:
            for p in base.rglob("*.lnk"):
                if not push(p):
                    return found
            for p in base.rglob("*.exe"):
                if not push(p):
                    return found
        except OSError:
            continue

    return found


def _windows_resolve_app_executable(app: str) -> tuple[Path | None, dict]:
    """
    在常见安装目录、桌面、开始菜单中解析可启动文件（.lnk / .exe）。
    返回：(可执行文件路径，信息字典)
    信息字典包含:
    - found: bool 是否找到
    - web_alternative: str|None 网页版 URL 建议
    - scan_locations: list 扫描过的位置
    """
    # 常见应用的网页版映射
    WEB_ALTERNATIVES = {
        "wechat": "https://web.wechat.com",
        "微信": "https://web.wechat.com",
        "weixin": "https://web.wechat.com",
        "qq": "https://im.qq.com",
        "钉钉": "https://www.dingtalk.com",
        "dingtalk": "https://www.dingtalk.com",
        "企业微信": "https://work.weixin.qq.com",
        "wps": "https://www.kdocs.cn",
        "word": "https://www.office.com",
        "excel": "https://www.office.com",
        "powerpoint": "https://www.office.com",
        "ppt": "https://www.office.com",
        "photoshop": "https://www.adobe.com",
        "ps": "https://www.adobe.com",
        "adobe": "https://www.adobe.com",
    }

    raw = (app or "").strip().strip('"').strip("'")
    if not raw:
        return None, {"found": False, "web_alternative": None}

    # 查找网页版替代方案
    web_alt = None
    raw_lower = raw.lower()
    for key, url in WEB_ALTERNATIVES.items():
        if key in raw_lower or raw_lower in key:
            web_alt = url
            break

    trial = Path(raw)
    if trial.is_file():
        return trial, {"found": True, "web_alternative": web_alt}
    if not trial.is_absolute():
        t2 = Path.cwd() / raw
        if t2.is_file():
            return t2, {"found": True, "web_alternative": web_alt}

    kws = _windows_open_app_keywords(raw)
    if not kws:
        return None, {"found": False, "web_alternative": web_alt}

    for cand in _windows_known_office_exes():
        if cand.is_file() and _windows_score_app_match(kws, cand) >= 650:
            return cand, {"found": True, "web_alternative": web_alt}

    scored: list[tuple[int, Path]] = []
    for p in _windows_collect_shortcuts_and_exes():
        sc = _windows_score_app_match(kws, p)
        if sc <= 0:
            continue
        scored.append((sc + _windows_desktop_path_bonus(p), p))

    if not scored:
        return None, {"found": False, "web_alternative": web_alt}

    scored.sort(key=lambda x: (-x[0], len(str(x[1]))))
    best, path = scored[0]
    if best < 650:
        return None, {"found": False, "web_alternative": web_alt}
    return path, {"found": True, "web_alternative": web_alt}


class _HTMLToTextParser(HTMLParser):
    """将 HTML 转为可读纯文本（跳过 script/style）。"""

    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in ("script", "style", "noscript", "template"):
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style", "noscript", "template") and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0 and data and data.strip():
            self._chunks.append(data.strip())

    def text(self) -> str:
        raw = " ".join(self._chunks)
        return re.sub(r"\s+", " ", raw).strip()


class TaskCancelledError(Exception):
    """当前任务被用户中止。"""


def _normalize_clarify_choices(raw: Any) -> list[dict[str, str]]:
    out: list[dict[str, str]] = []
    if not isinstance(raw, list):
        return out
    for i, item in enumerate(raw):
        if not isinstance(item, dict):
            continue
        cid = str(item.get("id") or item.get("value") or f"opt{i + 1}").strip()
        label = str(item.get("label") or item.get("text") or cid).strip()
        if cid and label:
            out.append({"id": cid, "label": label})
        if len(out) >= 12:
            break
    return out


# 用户字面中出现下列用语时，才视为「明确要上网检索」；与「本机打开微信/软件」类需求区分。
WEB_INFORMATION_MARKERS_CN = (
    "搜一下",
    "搜点",
    "搜索",
    "查资料",
    "网上找",
    "上网查",
    "帮我搜",
    "去搜",
    "热点新闻",
    "今日新闻",
    "最新新闻",
    "搜材料",
    "找材料",
    "检索",
    "联网查",
    "在线查",
    "外网",
    "百度一下",
    "谷歌一下",
    "搜教程",
    "查教程",
    "找教程",
    "安装教程",
    "官网查",
    "查官网",
    "维基百科",
    "百科一下",
    "网页摘要",
    "总结这个链接",
    "摘要这个网址",
)
WEB_INFORMATION_MARKERS_EN = (
    "search the web",
    "search online",
    "google it",
    "look up online",
    "wikipedia",
    "official website",
)


class ARIAManager:
    ALLOWED_ACTION_TYPES = {
        "kb_delete_all",
        "kb_delete_by_keyword",
        "kb_delete_low_quality",
        "conversation_new",
        "shell_run",
        "file_read",
        "file_write",
        "file_append",
        "file_create_dir",
        "file_move",
        "file_delete",
        "file_list",
        "file_find",
        "clipboard_read",
        "clipboard_write",
        "wechat_check_login",
        "wechat_open_chat",
        "wechat_send_message",
        "screen_watch_start",
        "screen_watch_stop",
        "email_send",
        "email_read",
        "browser_open",
        "browser_click",
        "browser_type",
        "browser_find",
        "browser_hover",
        "browser_select",
        "browser_upload",
        "browser_scroll",
        "browser_wait",
        "browser_js",
        "browser_press",
        "media_summarize",
        "desktop_open_app",
        "desktop_hotkey",
        "desktop_type",
        "desktop_sequence",
        "web_fetch",
        "web_understand",
        "screen_ocr",
        "screen_find_text",
        "screen_click_text",
        "computer_screenshot",
        "computer_click",
        "computer_click_element",
        "computer_double_click",
        "computer_move",
        "computer_drag",
        "computer_scroll",
        "computer_key",
        "computer_type",
        "computer_wait",
        "window_activate",
    }
    HIGH_RISK_ACTION_TYPES = {
        "kb_delete_all",
        "file_delete",
        "shell_run",
        "desktop_hotkey",
        "desktop_type",
        "desktop_sequence",
    }
    USER_GATE_ACTION_TYPES = frozenset(
        {
            "file_write",
            "file_append",
            "file_create_dir",
            "file_move",
            "file_delete",
            "clipboard_write",
            "wechat_open_chat",
            "wechat_send_message",
            "shell_run",
            "desktop_open_app",
            "desktop_hotkey",
            "desktop_type",
            "desktop_sequence",
            "screen_ocr",            "screen_find_text",
            "screen_click_text",
            "computer_screenshot",
            "computer_click",
            "computer_click_element",
            "computer_double_click",
            "computer_move",
            "computer_drag",
            "computer_scroll",
            "computer_key",
            "computer_type",
            "computer_wait",
            "media_summarize",
            "window_activate",
            "email_send",
        }
    )
    TOOL_PROFILES = {
        "local_execute": frozenset(
            {
                "conversation_new",
                "shell_run",
                "file_write",
                "file_move",
                "file_delete",
                "browser_open",
                "browser_click",
                "browser_type",
                "browser_find",
                "browser_hover",
                "browser_select",
                "browser_upload",
                "browser_scroll",
                "browser_wait",
                "browser_js",
                "browser_press",
                "desktop_open_app",
                "desktop_hotkey",
                "desktop_type",
                "desktop_sequence",
                "screen_ocr",
                "screen_find_text",
                "screen_click_text",
                "computer_screenshot",
                "computer_click",
                "computer_click_element",
                "computer_double_click",
                "computer_move",
                "computer_drag",
                "computer_scroll",
                "computer_key",
                "computer_type",
                "computer_wait",
                "media_summarize",
                "window_activate",
            }
        ),
        "web_information": frozenset(
            {
                "browser_open",
                "browser_click",
                "browser_type",
                "browser_find",
                "browser_hover",
                "browser_select",
                "browser_upload",
                "browser_scroll",
                "browser_wait",
                "browser_js",
                "browser_press",
                "web_fetch",
                "web_understand",
                "media_summarize",
            }
        ),
        "qa_only": frozenset(),
        "mixed": frozenset(ALLOWED_ACTION_TYPES),
    }
    WORKSPACE_MODE_PROFILES = {
        "aria": frozenset(ALLOWED_ACTION_TYPES),
        "aria_engineer_autocad": frozenset(
            {
                "conversation_new",
                "file_write",
                "file_move",
                "browser_open",
                "browser_find",
                "browser_wait",
                "desktop_open_app",
                "screen_ocr",
                "screen_find_text",
                "screen_click_text",
                "web_fetch",
                "web_understand",
                "media_summarize",
            }
        ),
    }

    def __init__(self, api_key: Optional[str] = None):
        self.model_pool = MODEL_POOL
        self.execution_log = []  # 全流程日志（给UI展示）
        self.workflow_events = []  # 结构化工作流事件（给实时时间线）
        self.model_thoughts = {}  # 模型思考过程
        self.current_conversation_id = ""
        self.current_task_id = ""
        self.current_request_id = ""
        self.action_screenshots_for_execution = False  # 由异步执行会话线程按前端/会话设置临时打开
        self._screenshot_last_ts = 0.0  # ARIA_SCREENSHOT_MIN_INTERVAL_MS 节流用
        self.cancelled_requests: set[str] = set()
        self.event_sink = None
        self.api_key = api_key or ""
        self.llm = VolcengineLLM(self.api_key) if self.api_key else VolcengineLLM(None)
        env_model = (os.getenv("MODEL_NAME") or "").strip()
        pool_model = str(self.model_pool.get("llm") or "").strip()
        self.unified_model = (
            env_model
            or pool_model
            or str(getattr(self.llm, "model_name", "") or "").strip()
            or "doubao-seed-2-0-lite-260215"
        )
        self.interaction_core = interaction_intelligence.InteractionIntelligenceCore()

        # 初始化三级记忆
        self.stm = ShortTermMemory()  # 短期记忆
        self.mtm = MidTermMemory()  # 中期记忆
        self.ltm = LongTermMemory()  # 长期记忆

        # 加载记忆
        self.mtm.load()
        self.ltm.load()
        self.auto_memory = AutoMemoryManager(self)
        self.exec_agent_name_pool: dict[str, list[str]] = {
            "TextExecAgent": [
                "李楠",
                "张弛",
                "王越",
                "赵晨",
                "陈屿",
                "周航",
                "吴泽",
                "郑川",
                "冯煦",
                "孙启",
                "马岩",
                "朱睿",
                "胡峻",
                "郭湛",
                "何川",
                "高远",
                "林朔",
                "罗尧",
                "梁恺",
                "谢恒",
                "宋川",
                "唐逸",
                "许诺",
                "韩骁",
                "曹峥",
                "彭锐",
                "袁景",
                "邓一鸣",
                "蒋澈",
                "沈奕",
            ],
            "VisionExecAgent": [
                "周岚",
                "顾宁",
                "程澄",
                "苏芮",
                "夏沫",
                "叶青",
                "白露",
                "安禾",
                "姜苒",
                "陆悠",
                "沈禾",
                "温雅",
                "林汐",
                "贺晴",
                "乔然",
                "宋颜",
                "唐婧",
                "许薇",
                "袁念",
                "邱彤",
                "施瑶",
                "徐静",
                "韩悦",
                "罗妍",
                "蔡宁",
                "孔乔",
                "杜曼",
                "陶冉",
                "毛伊",
                "尹澜",
            ],
            "SpeechExecAgent": [
                "赵尧",
                "林嘉",
                "方恺",
                "魏哲",
                "潘越",
                "吕衡",
                "严朗",
                "任川",
                "施博",
                "钟宁",
                "董恺",
                "孟川",
                "祁峰",
                "易然",
                "池恒",
                "裴青",
                "邢岳",
                "鲍骁",
                "洪毅",
                "汪言",
                "贾睿",
                "范哲",
                "樊涛",
                "邹赫",
                "石航",
                "雷靖",
                "龙湛",
                "万川",
                "段驰",
                "侯野",
            ],
        }
        self.action_registry = {
            "kb_delete_all": self._exec_kb_delete_all,
            "kb_delete_by_keyword": self._exec_kb_delete_by_keyword,
            "kb_delete_low_quality": self._exec_kb_delete_low_quality,
            "conversation_new": self._exec_conversation_new,
            "shell_run": self._exec_shell_run,
            "file_read": self._exec_file_read,
            "file_write": self._exec_file_write,
            "file_append": self._exec_file_append,
            "file_create_dir": self._exec_file_create_dir,
            "file_move": self._exec_file_move,
            "file_delete": self._exec_file_delete,
            "file_list": self._exec_file_list,
            "file_find": self._exec_file_find,
            "clipboard_read": self._exec_clipboard_read,
            "clipboard_write": self._exec_clipboard_write,
            "wechat_check_login": self._exec_wechat_check_login,
            "wechat_open_chat": self._exec_wechat_open_chat,
            "wechat_send_message": self._exec_wechat_send_message,
            "screen_watch_start": self._exec_screen_watch_start,
            "screen_watch_stop": self._exec_screen_watch_stop,
            "email_send": self._exec_email_send,
            "email_read": self._exec_email_read,
            "browser_open": self._exec_browser_open,
            "browser_click": self._exec_browser_click,
            "browser_type": self._exec_browser_type,
            "browser_find": self._exec_browser_find,
            "browser_hover": self._exec_browser_hover,
            "browser_select": self._exec_browser_select,
            "browser_upload": self._exec_browser_upload,
            "browser_scroll": self._exec_browser_scroll,
            "browser_wait": self._exec_browser_wait,
            "browser_js": self._exec_browser_js,
            "browser_press": self._exec_browser_press,
            "media_summarize": self._exec_media_summarize,
            "desktop_open_app": self._exec_desktop_open_app,
            "desktop_hotkey": self._exec_desktop_hotkey,
            "desktop_type": self._exec_desktop_type,
            "desktop_sequence": self._exec_desktop_sequence,
            "web_fetch": self._exec_web_fetch,
            "web_understand": self._exec_web_understand,
            "screen_ocr": self._exec_screen_ocr,
            "screen_find_text": self._exec_screen_find_text,
            "screen_click_text": self._exec_screen_click_text,
            "computer_screenshot": self._exec_computer_screenshot,
            "computer_click": self._exec_computer_click,
            "computer_click_element": self._exec_computer_click_element,
            "computer_double_click": self._exec_computer_double_click,
            "computer_move": self._exec_computer_move,
            "computer_drag": self._exec_computer_drag,
            "computer_scroll": self._exec_computer_scroll,
            "computer_key": self._exec_computer_key,
            "computer_type": self._exec_computer_type,
            "computer_wait": self._exec_computer_wait,
            "window_activate": self._exec_window_activate,
        }
        self.execution_sessions: dict[str, dict[str, Any]] = {}
        self.execution_lock = threading.Lock()
        self.allowed_work_root = Path(os.path.abspath("."))
        self.max_action_steps = 30
        self.max_action_retries = 1
        self.max_action_retries = 0
        try:
            self.max_react_iterations = max(1, min(60, int(os.getenv("ARIA_REACT_MAX_STEPS", "20"))))
        except (TypeError, ValueError):
            self.max_react_iterations = 20
        self.default_step_timeout_s = 90
        self._ab_context_by_task: dict[str, dict[str, Any]] = {}
        self.last_model_trace: dict[str, Any] = {}
        self._token_usage_tls = threading.local()
        self._turn_vision_data_urls: list[str] = []
        self._turn_reasoning_effort: str | None = None
        self.workspace_mode = self._normalize_workspace_mode(os.getenv("ARIA_DEFAULT_WORKSPACE_MODE") or "aria")
        from automation.app_framework import ApplicationRegistry
        self.app_registry = ApplicationRegistry()
        from runtime.orchestration import OrchestrationFacade
        self.orchestration = OrchestrationFacade(self)
        self.permission_model = PermissionModel()

        # 人格目录（personality catalog）
        self.personality_map_config: dict[str, Any] = {}
        self.personality_catalog: dict[str, list[dict[str, Any]]] = {}
        self._personality_agents_root: Path | None = None
        self._load_personality_catalog()

        # KAIROS 主动执行引擎（可选）
        self.kairos: Any = None
        if os.getenv("ARIA_KAIROS_ENABLED", "0").strip().lower() in ("1", "true", "yes"):
            try:
                from runtime.kairos import KAIROSEngine
                self.kairos = KAIROSEngine(self)
                self.kairos.start()
            except Exception as _kairos_err:
                logger.warning("KAIROS 启动失败：%s", _kairos_err)

    def record_kairos_activity(self) -> None:
        """通知 AutoDream 重置空闲计时器。每次处理用户请求时调用。"""
        if self.kairos is not None:
            try:
                self.kairos.dream_engine.record_activity()
            except Exception:
                pass

    def _resolve_agency_agents_root(self) -> "Path | None":
        """定位 agency-agents 目录（包含各人格 .md 文件）。"""
        base = Path(__file__).parent
        candidates = [
            base / "third_party" / "agency-agents" / "agency-agents-main",
            base / "third_party" / "agency-agents",
            base / "agency-agents-main" / "agency-agents-main",
            base / "agency-agents-main",
            base / "agency_agents",
            base / "agents",
        ]
        for c in candidates:
            if c.is_dir():
                return c
        return None

    def _load_personality_catalog(self) -> None:
        """从 config/agent_personality_map.json 加载人格目录。"""
        try:
            cfg_path = Path(__file__).parent / "config" / "agent_personality_map.json"
            if not cfg_path.is_file():
                return
            import json as _json
            with cfg_path.open(encoding="utf-8") as f:
                data = _json.load(f)
            self.personality_map_config = data if isinstance(data, dict) else {}
            root = self._resolve_agency_agents_root()
            self._personality_agents_root = root
            profiles = self.personality_map_config.get("global_profiles")
            if not isinstance(profiles, list):
                return
            rows: list[dict[str, Any]] = []
            for p in profiles:
                file_rel = str(p.get("file") or "").strip()
                excerpt = ""
                if root and file_rel:
                    md_path = root / file_rel
                    try:
                        excerpt = md_path.read_text(encoding="utf-8")[:300]
                    except Exception:
                        excerpt = ""
                rows.append({**p, "source_file": file_rel, "excerpt": excerpt})
            for agent_type in ("TextExecAgent", "VisionExecAgent", "SpeechExecAgent"):
                self.personality_catalog[agent_type] = rows
        except Exception as e:
            logger.debug("_load_personality_catalog failed: %s", e)

    def _default_reasoning_effort_from_env(self) -> str:
        raw = (os.getenv("REASONING_EFFORT_DEFAULT") or "medium").strip().lower()
        return _normalize_reasoning_effort(raw) or "medium"

    def set_turn_reasoning_effort(self, effort: str | None) -> None:
        n = _normalize_reasoning_effort(effort)
        self._turn_reasoning_effort = n

    def clear_turn_reasoning_effort(self) -> None:
        self._turn_reasoning_effort = None

    def _reasoning_effort_heuristic(
        self,
        llm_user_input: str,
        dialogue_context: str,
        *,
        has_attachments: bool,
        attachment_exts: list[str],
    ) -> str:
        text = (llm_user_input or "").strip()
        dc = (dialogue_context or "").strip()
        compact = re.sub(r"[\s\W_]+", "", text.lower())
        greeting_keywords = (
            "你好",
            "您好",
            "hello",
            "hi",
            "hey",
            "在吗",
            "谢谢",
            "thank",
            "早上好",
            "晚上好",
            "再见",
            "拜拜",
        )
        high_kw = (
            "深度分析",
            "逐步推导",
            "详细论证",
            "复杂方案",
            "架构设计",
            "完整证明",
            "严谨推导",
            "逐步推理",
        )
        if any(k in text for k in high_kw) or len(text) > 2000:
            return "high"
        exts = {(e or "").lower() for e in (attachment_exts or [])}
        docish = exts & {"png", "jpg", "jpeg", "gif", "webp", "pdf", "docx", "xlsx", "xlsm", "pptx"}
        if has_attachments or docish:
            return "medium"
        if any(
            k in text
            for k in ("步骤", "计划", "实现", "编写", "部署", "调试", "多个", "首先", "然后", "拆分", "子任务")
        ):
            return "medium"
        if len(text) > 800 or len(dc) > 1500:
            return "medium"
        if any(k in text.lower() for k in greeting_keywords) and len(compact) <= 24:
            return "minimal"
        if len(text) < 40 and not has_attachments:
            return "low"
        if len(text) < 200 and not has_attachments:
            return "low"
        return self._default_reasoning_effort_from_env()

    def _reasoning_effort_llm_router(
        self,
        llm_user_input: str,
        dialogue_context: str,
        *,
        has_attachments: bool,
    ) -> str | None:
        if self.is_cancelled():
            return None
        if not getattr(self.llm, "api_key", None) or not callable(getattr(self.llm, "generate", None)):
            return None
        preview = (llm_user_input or "").strip()[:1200]
        messages = [
            {
                "role": "system",
                "content": (
                    "你是ARIA推理强度路由器。根据用户本轮任务复杂度，从 minimal、low、medium、high 中选一。"
                    "minimal：寒暄、极短确认、无实质任务；low：单轮简单问答；medium：含附件、多步、实现类；"
                    "high：明确要求深度推导/复杂方案或输入很长。"
                    '只输出JSON：{"reasoning_effort":"minimal|low|medium|high","reason":"..."}'
                ),
            },
            {
                "role": "user",
                "content": (
                    f"是否有附件: {bool(has_attachments)}\n"
                    f"近期对话字符数: {len((dialogue_context or '').strip())}\n\n"
                    f"本轮输入:\n{preview or '（空）'}"
                ),
            },
        ]
        t_llm = time.perf_counter()
        try:
            raw, usage = self.llm.generate(
                messages,
                model_name=self.unified_model,
                reasoning_effort="minimal",
            )
            self._accumulate_usage_dict(usage)
        except Exception:
            return None
        finally:
            self._accumulate_llm_wall_ms(int((time.perf_counter() - t_llm) * 1000))
        data = self._extract_json_object(raw or "")
        r = _normalize_reasoning_effort(str((data or {}).get("reasoning_effort", "")).strip())
        return r

    def resolve_reasoning_effort_for_turn(
        self,
        llm_user_input: str,
        dialogue_context: str,
        *,
        has_attachments: bool = False,
        attachment_exts: list[str] | None = None,
    ) -> str:
        """指挥官适应性分析：默认规则；ARIA_REASONING_ROUTER=llm 时先尝试一次轻量模型路由。"""
        router = (os.getenv("ARIA_REASONING_ROUTER") or "").strip().lower()
        exts = list(attachment_exts or [])
        if router in ("1", "true", "yes", "on", "llm"):
            routed = self._reasoning_effort_llm_router(
                llm_user_input,
                dialogue_context,
                has_attachments=has_attachments,
            )
            if routed:
                return routed
        return self._reasoning_effort_heuristic(
            llm_user_input,
            dialogue_context,
            has_attachments=has_attachments,
            attachment_exts=exts,
        )

    def set_turn_vision_images(self, data_urls: list[str] | None) -> None:
        """本轮 HTTP 请求内：将上传图片以 data URL 形式随下一条 user 多模态消息发送（由 web_app 设置）。"""
        self._turn_vision_data_urls = [u.strip() for u in (data_urls or []) if isinstance(u, str) and u.strip()]

    def clear_turn_vision_images(self) -> None:
        self._turn_vision_data_urls = []

    def _user_content_with_optional_vision(self, text: str) -> str | list[dict[str, Any]]:
        urls = self._turn_vision_data_urls
        if not urls:
            return text
        parts: list[dict[str, Any]] = [{"type": "text", "text": text}]
        for u in urls:
            parts.append({"type": "image_url", "image_url": {"url": u}})
        return parts

    def reset_token_usage(self) -> None:
        """按线程清空用量计数（每个 HTTP 请求 / 异步执行线程各自一份）。"""
        self._token_usage_tls.accumulator = {
            "prompt_tokens": 0,
            "completion_tokens": 0,
            "total_tokens": 0,
            "llm_calls": 0,
            "llm_wall_ms": 0,
        }

    def get_token_usage_summary(self) -> dict[str, int]:
        acc = getattr(self._token_usage_tls, "accumulator", None)
        if not acc:
            return {
                "prompt_tokens": 0,
                "completion_tokens": 0,
                "total_tokens": 0,
                "llm_calls": 0,
                "llm_wall_ms": 0,
            }
        keys = ("prompt_tokens", "completion_tokens", "total_tokens", "llm_calls", "llm_wall_ms")
        return {k: int(acc.get(k) or 0) for k in keys}

    def _accumulate_llm_wall_ms(self, ms: int) -> None:
        delta = int(ms)
        if delta <= 0:
            return
        if not hasattr(self._token_usage_tls, "accumulator"):
            self.reset_token_usage()
        acc = self._token_usage_tls.accumulator
        acc["llm_wall_ms"] = int(acc.get("llm_wall_ms", 0) or 0) + delta

    def _accumulate_usage_dict(self, usage: dict[str, int]) -> None:
        if not usage:
            return
        if not hasattr(self._token_usage_tls, "accumulator"):
            self.reset_token_usage()
        acc = self._token_usage_tls.accumulator
        for k in ("prompt_tokens", "completion_tokens", "total_tokens", "llm_calls"):
            acc[k] = acc.get(k, 0) + int(usage.get(k) or 0)

    def set_api_key(self, api_key: Optional[str]) -> None:
        api_key = (api_key or "").strip()
        if api_key == self.api_key:
            return
        self.api_key = api_key
        self.llm = VolcengineLLM(api_key if api_key else None)
        env_model = (os.getenv("MODEL_NAME") or "").strip()
        pool_model = str(self.model_pool.get("llm") or "").strip()
        self.unified_model = (
            env_model
            or pool_model
            or str(getattr(self.llm, "model_name", "") or "").strip()
            or "doubao-seed-2-0-lite-260215"
        )

    def set_conversation_context(self, conversation_id: str) -> None:
        self.current_conversation_id = conversation_id or ""

    def _normalize_workspace_mode(self, mode: str | None) -> str:
        raw = str(mode or "").strip().lower().replace("-", "_")
        if raw in self.WORKSPACE_MODE_PROFILES:
            return raw
        if raw in ("aria_engineer", "engineer", "autocad"):
            return "aria_engineer_autocad"
        return "aria"

    def set_workspace_mode(self, mode: str | None) -> str:
        normalized = self._normalize_workspace_mode(mode)
        self.workspace_mode = normalized
        return normalized

    def set_event_sink(self, sink) -> None:
        """注册事件回调，供 SSE 推送。"""
        self.event_sink = sink

    # 记录模型思考过程
    def record_model_thought(self, agent_name, thought):
        if agent_name not in self.model_thoughts:
            self.model_thoughts[agent_name] = []
        self.model_thoughts[agent_name].append({"thought": thought, "timestamp": time.time()})

    # 获取模型思考过程
    def get_model_thoughts(self, agent_name):
        return self.model_thoughts.get(agent_name, [])

    # 清空模型思考过程
    def clear_model_thoughts(self):
        self.model_thoughts = {}

    def _agent_profile(self, agent_code: str) -> dict[str, str]:
        mapping = {
            "TaskParser": {"role": "项目经理PM", "name": "王琳"},
            "MethodSearcher": {"role": "知识专家KS", "name": "陈舟"},
            "SolutionLearner": {"role": "执行专家EXE", "name": "李楠"},
            "TaskSplitter": {"role": "项目经理PM", "name": "王琳"},
            "TextExecAgent": {"role": "执行专家EXE", "name": "李楠"},
            "VisionExecAgent": {"role": "视觉专家VE", "name": "周岚"},
            "SpeechExecAgent": {"role": "语音专家SE", "name": "赵尧"},
            "QualityChecker": {"role": "质检QA", "name": "周启"},
            "MethodSaver": {"role": "知识专家KS", "name": "陈舟"},
        }
        return mapping.get(agent_code, {"role": "执行专家EXE", "name": "李楠"})

    def _pick_exec_agent_name(self, agent_type: str, used_names: dict[str, set[str]]) -> str:
        pool = self.exec_agent_name_pool.get(agent_type, [])
        fallback = self._agent_profile(agent_type)["name"]
        if not pool:
            return fallback
        used = used_names.setdefault(agent_type, set())
        candidates = [n for n in pool if n not in used]
        picked = random.choice(candidates or pool)
        used.add(picked)
        return picked

    def emit_transport_event(self, payload: dict[str, Any]) -> None:
        """经 event_sink 推送非 workflow 事件（如 LLM 流式分片），不写入 workflow_events。"""
        if not self.event_sink or not isinstance(payload, dict):
            return
        evt = {
            "conversation_id": self.current_conversation_id,
            "request_id": self.current_request_id,
            **payload,
        }
        try:
            self.event_sink(evt)
        except Exception:
            pass

    def push_event(
        self,
        stage: str,
        status: str,
        agent_code: str,
        summary: str,
        detail: dict[str, Any] | None = None,
        agent_name_override: str | None = None,
    ) -> None:
        profile = self._agent_profile(agent_code)
        event = {
            "event_id": str(uuid.uuid4()),
            "conversation_id": self.current_conversation_id,
            "task_id": self.current_task_id,
            "request_id": self.current_request_id,
            "stage": stage,
            "status": status,  # pending/running/success/error
            "agent_code": agent_code,
            "agent_name": agent_name_override or profile["name"],
            "agent_role": profile["role"],
            "summary": summary,
            "detail": detail or {},
            "timestamp": time.time(),
        }
        self.workflow_events.append(event)
        if self.event_sink:
            try:
                self.event_sink(event)
            except Exception:
                pass

    def get_workflow_events(self) -> list[dict[str, Any]]:
        return self.workflow_events

    def clear_workflow_events(self) -> None:
        self.workflow_events = []

    def request_cancel(self, request_id: str) -> bool:
        rid = (request_id or "").strip()
        if not rid:
            return False
        self.cancelled_requests.add(rid)
        return True

    def clear_cancel(self, request_id: str) -> None:
        rid = (request_id or "").strip()
        if rid and rid in self.cancelled_requests:
            self.cancelled_requests.remove(rid)

    def is_cancelled(self, request_id: str | None = None) -> bool:
        rid = (request_id or self.current_request_id or "").strip()
        return bool(rid) and rid in self.cancelled_requests

    def check_cancelled(self, stage: str) -> None:
        if self.is_cancelled():
            self.push_event(
                "task_cancelled",
                "success",
                "TaskParser",
                "任务已中止",
                {"stage": stage, "request_id": self.current_request_id},
            )
            self.push_log("系统", f"任务已中止（阶段：{stage}）", "warning")
            raise TaskCancelledError(f"任务已中止（{stage}）")

    def _normalize_query_text(self, text: str) -> str:
        s = (text or "").strip().lower()
        s = re.sub(r"\s+", "", s)
        s = re.sub(r"[，。！？,.!?:;；、\"'“”‘’`~@#$%^&*()（）\[\]{}<>《》\-_=+\\/|]+", "", s)
        return s

    def find_exact_methodology(self, user_input: str) -> dict[str, Any] | None:
        normalized_query = self._normalize_query_text(user_input)
        if not normalized_query:
            return None
        methodologies = getattr(self.ltm, "methodologies", []) or []
        for method in methodologies:
            scene = str(method.get("scene") or method.get("scenario") or "")
            if self._normalize_query_text(scene) == normalized_query:
                return method
        return None

    def _infer_temporal_risk(self, user_input: str) -> str:
        """启发式：强时效事实类任务标 high，其余 low。"""
        text = (user_input or "").strip()
        if not text:
            return "low"
        t = text.lower()
        cn = (
            "天气",
            "气温",
            "降雨",
            "下雪",
            "台风",
            "空气质量",
            "雾霾",
            "股价",
            "涨跌",
            "汇率",
            "比分",
            "赛程",
            "积分榜",
            "金价",
            "油价",
            "期货",
            "今日天气",
            "现在天气",
            "实时",
            "最新股价",
            "赛果",
            "限行",
        )
        if any(k in text for k in cn):
            return "high"
        en = ("weather forecast", "stock price", "live score", "exchange rate")
        if any(k in t for k in en):
            return "high"
        return "low"

    def should_reuse_methodology_without_learn(
        self,
        task_info: dict[str, Any],
        score: float,
        method: Any,
    ) -> bool:
        """
        强时效任务在相似度未达 0.7 时，若已有可用的流程型方法论（分数 ≥ 可配置下限），跳过 learn_from_external 以省 Token。
        """
        if score >= 0.7:
            return False
        if not method or not isinstance(method, dict):
            return False
        temporal = str(task_info.get("temporal_risk") or "low").strip().lower()
        if temporal != "high":
            return False
        try:
            floor = float(os.getenv("ARIA_TEMPORAL_METHOD_MATCH_FLOOR", "0.45"))
        except ValueError:
            floor = 0.45
        if score < floor:
            return False
        steps = method.get("solve_steps") or method.get("steps") or []
        if isinstance(steps, str):
            steps = [s.strip() for s in re.split(r"[\n;；]+", steps) if s.strip()]
        return len(steps) >= 1

    def should_skip_external_methodology_learning(self, task_info: dict[str, Any], user_input: str) -> bool:
        """本机自动化类任务不应走 learn_from_external 生成「上网找步骤」式方法论。"""
        s = str((task_info or {}).get("execution_surface") or "").strip().lower()
        if s in ("local_desktop", "local", "desktop", "desktop_automation"):
            return True
        return False

    def local_automation_methodology_stub(self, task_info: dict[str, Any]) -> dict[str, Any]:
        """跳过外网学习时使用的简短流程占位，引导 Agent 走本机动作而非检索教程。"""
        kws = task_info.get("keywords") if isinstance(task_info.get("keywords"), list) else []
        kw = [str(x).strip() for x in kws if str(x).strip()][:8] or ["本机操作"]
        return {
            "scene": "本机应用/桌面自动化",
            "keywords": kw,
            "solve_steps": [
                "判定：需求在本机已装应用内完成，勿以网页检索教程为主交付。",
                "回复：引导用户查看是否已出现「执行计划」卡片并点击确认执行；若无卡片则上一轮未生成可执行动作。",
                "禁忌：勿编造已发送或已点击；勿展开联网检索步骤。",
            ],
            "applicable_range": "ARIA 本机动作执行",
            "outcome_type": "pure_procedure",
        }

    def _llm_transport_stream_enabled(self) -> bool:
        return os.getenv("ARIA_LLM_STREAM", "0").strip().lower() in ("1", "true", "yes", "on")

    def _call_llm_fast(
        self,
        messages: list[dict[str, Any]],
        fallback_text: str = "",
        *,
        temperature: float = 0.3,
        max_tokens: int = 2048,
    ) -> str:
        """
        快速推理路径：优先使用 Groq（低延迟），不可用时回退到默认 LLM。
        适用于 small_talk、direct_qa 等对延迟敏感但不需要工具调用的场景。
        """
        from llm.groq_llm import GroqLLM, is_groq_enabled
        if is_groq_enabled():
            try:
                groq = GroqLLM()
                result = groq.chat(messages, temperature=temperature, max_tokens=max_tokens)
                return str(result).strip()
            except Exception as e:
                self.push_event("groq_fallback", "warning", "GroqLLM", f"Groq 调用失败，回退主模型: {e}", {})
        return self._call_llm(messages, fallback_text=fallback_text, agent_code="TextExecAgent", reasoning_effort="minimal")

    def _call_llm(
        self,
        messages: list[dict[str, Any]],
        fallback_text: str = "",
        agent_code: str = "",
        reasoning_effort: str | None = None,
        *,
        llm_stream: bool | None = None,
    ) -> str:
        """调用 LLM（全链路统一模型）；失败时返回 fallback_text。reasoning_effort 为 None 时用本轮 set_turn_reasoning_effort 或环境默认。
        llm_stream：显式开关 SSE 文本流；None 时由 ARIA_LLM_STREAM 与 Agent 类型决定。"""
        if self.is_cancelled():
            return fallback_text
        model = self.unified_model
        resolved_eff = reasoning_effort
        if resolved_eff is None:
            resolved_eff = self._turn_reasoning_effort
        eff_arg = _normalize_reasoning_effort(resolved_eff) or self._default_reasoning_effort_from_env()
        ac = agent_code or ""
        if llm_stream is None:
            want_stream = self._llm_transport_stream_enabled() and ac not in _LLM_STREAM_BLOCK_AGENTS
        else:
            want_stream = bool(llm_stream)
        errors: list[str] = []
        max_try = 3
        try:
            if not getattr(self, "llm", None) or not callable(getattr(self.llm, "generate", None)):
                return fallback_text
            stream_fn = getattr(self.llm, "generate_stream_chunks", None)
            for attempt in range(max_try):
                try:
                    t0 = time.perf_counter()
                    text = ""
                    usage: dict[str, int] | None = None
                    try:
                        if (
                            want_stream
                            and callable(stream_fn)
                            and not self.is_cancelled()
                        ):
                            usage_holder: list[dict[str, int]] = []
                            parts: list[str] = []
                            stream_iter = stream_fn(
                                messages,
                                model_name=model,
                                reasoning_effort=eff_arg,
                                usage_holder=usage_holder,
                            )
                            self.emit_transport_event(
                                {"sse_kind": "llm_stream_start", "agent_code": ac, "reasoning_effort": eff_arg}
                            )
                            try:
                                for piece in stream_iter:
                                    if self.is_cancelled():
                                        break
                                    if piece:
                                        parts.append(piece)
                                        self.emit_transport_event(
                                            {"sse_kind": "llm_stream_delta", "delta": piece, "agent_code": ac}
                                        )
                                text = "".join(parts)
                                if usage_holder:
                                    usage = usage_holder[0]
                                if not (text or "").strip():
                                    text, usage = self.llm.generate(
                                        messages, model_name=model, reasoning_effort=eff_arg
                                    )
                            finally:
                                try:
                                    stream_iter.close()
                                except Exception:
                                    pass
                                self.emit_transport_event({"sse_kind": "llm_stream_end", "agent_code": ac})
                        else:
                            text, usage = self.llm.generate(
                                messages, model_name=model, reasoning_effort=eff_arg
                            )
                    finally:
                        self._accumulate_llm_wall_ms(int((time.perf_counter() - t0) * 1000))
                    if usage:
                        self._accumulate_usage_dict(usage)
                    self.last_model_trace = {
                        "agent_code": ac,
                        "model": model,
                        "timestamp": time.time(),
                        "reasoning_effort": eff_arg,
                    }
                    self.push_event(
                        "llm_route",
                        "success",
                        ac or "TaskParser",
                        f"模型: {model}",
                        {"agent_code": ac or "", "model": model, "reasoning_effort": eff_arg},
                    )
                    return text
                except Exception as inner:
                    err = str(inner)
                    errors.append(f"{attempt + 1}:{err}")
                    if attempt + 1 < max_try:
                        time.sleep(min(2 * (attempt + 1), 8))
                        continue
                    raise
            raise RuntimeError("; ".join(errors) if errors else "llm_call_failed")
        except Exception as e:
            self.push_log("LLM", f"LLM调用失败: {str(e)}", "warning")
            self.push_event(
                "llm_route",
                "error",
                ac or "TaskParser",
                f"模型调用失败: {model}",
                {"agent_code": ac or "", "model": model, "error": str(e), "tried": errors},
            )
            return fallback_text

    def _extract_json_object(self, text: str) -> dict[str, Any]:
        if not text:
            return {}
        cleaned = text.strip()
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"```$", "", cleaned).strip()
        m = re.search(r"\{[\s\S]*\}", cleaned)
        if not m:
            return {}
        try:
            return json.loads(m.group(0))
        except Exception:
            return {}

    def _default_persona_brief(self, agent_type: str, step: str, description: str) -> str:
        step_s = (step or "").strip() or "当前步骤"
        desc_s = (description or "").strip()
        if agent_type == "VisionExecAgent":
            return (
                f"你是视觉向执行专家。针对步骤「{step_s}」，从图像/界面/图表角度给出可核验的观察与结论；"
                "说明信息来源假设；勿编造未见内容。"
                f"交付须覆盖：{desc_s or '该步骤要求的视觉分析产出'}。"
            )
        if agent_type == "SpeechExecAgent":
            return (
                f"你是语音/口语向执行专家。针对步骤「{step_s}」，给出适合朗读或对话场景的表述；"
                "标注语气与停顿建议（如需要）；保持简洁可懂。"
                f"交付须覆盖：{desc_s or '该步骤要求的语音相关产出'}。"
            )
        return (
            f"你是文本向执行专家。针对步骤「{step_s}」，进行清晰推理与结构化输出；"
            "区分事实与推断；给出可执行的下一步建议（如适用）。"
            f"交付须覆盖：{desc_s or '该步骤要求的文本产出'}。"
        )

    def _methodology_summary_text(self, method: dict[str, Any] | None) -> str:
        if not method:
            return ""
        scene = str(method.get("scene") or method.get("scenario") or "").strip()
        steps = method.get("solve_steps") or method.get("steps") or []
        if isinstance(steps, str):
            steps = [s.strip() for s in re.split(r"[\n;；]+", steps) if s.strip()]
        lines = [
            "【方法论上下文 — 请与下列子步骤对齐，勿偏离总目标】",
            f"场景(scene)：{scene or '（未命名）'}",
            "步骤纲要(solve_steps)：",
        ]
        for i, s in enumerate(steps, 1):
            lines.append(f"  {i}. {s}")
        return "\n".join(lines)

    def _format_exec_results_as_plain_text(self, results: list[dict[str, Any]]) -> str:
        blocks: list[str] = []
        for idx, r in enumerate(results, start=1):
            agent_type = str(r.get("agent_type") or "")
            display_name = str(r.get("agent_name") or self._agent_profile(agent_type)["name"])
            step = str(r.get("step") or "")
            desc = str(r.get("description") or "")
            body = str(r.get("result") or "")
            blocks.append(
                f"## ARIA_PRIOR_STEP_{idx}\n"
                f"- 执行者显示名: {display_name}\n"
                f"- agent_type: {agent_type}\n"
                f"- step: {step}\n"
                f"- description: {desc}\n"
                f"- 完整产出(result，勿省略理解):\n{body}\n"
            )
        return "\n".join(blocks) if blocks else ""

    def _normalize_keywords(self, keywords: Any) -> list[str]:
        if keywords is None:
            return []
        if isinstance(keywords, list):
            return [str(k).strip() for k in keywords if str(k).strip()]
        if isinstance(keywords, str):
            # 按中文/英文常见分隔符拆分
            parts = re.split(r"[\s,，;；/|]+", keywords.strip())
            return [p.strip() for p in parts if p.strip()]
        return [str(keywords).strip()] if str(keywords).strip() else []

    def _normalize_methodology(self, method: dict[str, Any], task_info: dict[str, Any]) -> dict[str, Any]:
        scene = method.get("scene") or method.get("scenario") or task_info.get("user_input", "")[:50]
        keywords = method.get("keywords")
        if keywords is None:
            keywords = method.get("core_keywords")
        keywords = self._normalize_keywords(keywords) or task_info.get("user_input", "").split()[:5]

        solve_steps = method.get("solve_steps") or method.get("steps") or []
        if isinstance(solve_steps, str):
            solve_steps = [s.strip() for s in re.split(r"[\n;；]+", solve_steps) if s.strip()]

        applicable_range = method.get("applicable_range", "") or method.get("applicability", "")
        text = " ".join([scene, " ".join(keywords), task_info.get("user_input", "")]).lower()
        category = "通用/其他"
        greeting_keywords = ["你好", "您好", "hello", "hi", "在吗", "谢谢", "早上好", "晚上好"]
        if any(g in text for g in greeting_keywords):
            category = "通用/其他"
        elif sum(1 for k in ["分析", "报表", "指标", "趋势", "data", "sql"] if k in text) >= 2:
            category = "数据分析/报表"
        elif sum(1 for k in ["代码", "开发", "接口", "python", "bug", "java", "前端", "后端"] if k in text) >= 2:
            category = "开发工程/代码实现"
        elif sum(1 for k in ["需求", "产品", "交互", "原型", "prd"] if k in text) >= 2:
            category = "产品设计/需求"
        elif sum(1 for k in ["运营", "增长", "市场", "投放", "活动"] if k in text) >= 2:
            category = "运营增长/市场"
        title = method.get("title") or method.get("name") or (scene[:24] if scene else "新方法论")
        event_key = f"{str(scene).strip().lower()}|{'/'.join(sorted({k.lower() for k in keywords})[:4])}"

        out: dict[str, Any] = {
            "method_id": str(uuid.uuid4()),
            "title": title,
            "category": category,
            "scene": scene,
            "scenario": scene,  # 兼容 LongTermMemory 相似度逻辑
            "keywords": keywords,
            "core_keywords": keywords,  # 兼容 LongTermMemory 相似度逻辑
            "solve_steps": solve_steps,
            "applicable_range": applicable_range,
            "event_key": event_key,
            "success_count": int(bool(method.get("is_success", False))),
            "usage_count": 0,
            "create_time": time.strftime("%Y-%m-%d %H:%M:%S"),
        }
        ot = str(method.get("outcome_type") or "").strip().lower()
        if ot in ("stable", "time_bound", "pure_procedure"):
            out["outcome_type"] = ot
        return out

    def _should_save_methodology(
        self,
        task_info: dict[str, Any],
        method: dict[str, Any],
        result_payload: Any,
    ) -> tuple[bool, str, str]:
        """
        判定是否应沉淀方法论。
        返回 (should_save, reason, source)。
        source: llm / heuristic

        三道闸门：
        1. 复杂度太低（complexity_score <= 2）不保存
        2. 强时效性（temporal_risk = high）不保存
        3. outcome_type = time_bound 不保存
        """
        # 闸门 0：检查 plan 中的复杂度评分（如果有）
        complexity = int(task_info.get("complexity_score", 3))
        if complexity <= 2:
            return False, "complexity_too_low", "heuristic"

        # 闸门 1：强时效性
        temporal = str(task_info.get("temporal_risk", "low")).strip().lower()
        if temporal == "high":
            return False, "temporal_risk_high", "heuristic"

        # 闸门 2：time_bound 类型
        outcome = str(task_info.get("outcome_type", "stable")).strip().lower()
        if outcome == "time_bound":
            return False, "outcome_time_bound", "heuristic"

        user_input = str(task_info.get("user_input", "") or "").strip()
        if not user_input:
            return False, "empty_input", "heuristic"

        # 1) 优先走 LLM 判定（更灵活，避免硬编码）
        method_scene = str((method or {}).get("scene") or (method or {}).get("scenario") or "").strip()
        method_keywords = self._normalize_keywords(
            (method or {}).get("keywords") or (method or {}).get("core_keywords")
        )
        step_count = len((method or {}).get("solve_steps") or (method or {}).get("steps") or [])
        is_success = bool(result_payload.get("is_success")) if isinstance(result_payload, dict) else False

        messages = [
            {
                "role": "system",
                "content": (
                    "你是ARIA知识沉淀判定器。目标：判断这次对话是否值得沉淀为可复用方法论。"
                    '请只输出JSON：{"should_save":true/false,"reason":"..."}。'
                    "规则：寒暄/问候/纯礼貌/无明确任务目标 -> false；"
                    "有清晰任务目标、可复用步骤或可迁移经验 -> true。"
                    "若 temporal_risk 为 high（天气/股价等强时效），且用户仅要一次性当下事实、方法论又无「如何查权威源/参数槽位」等可迁移流程，应 should_save:false，避免把过期答案写进知识库。"
                    "若 high 但 solve_steps 明确写了可重复的获取与校验流程（非具体数值），可 should_save:true。"
                ),
            },
            {
                "role": "user",
                "content": (
                    f"user_input: {user_input}\n"
                    f"task_type: {task_info.get('task_type', '')}\n"
                    f"intent: {task_info.get('intent', '')}\n"
                    f"temporal_risk: {task_info.get('temporal_risk', 'low')}\n"
                    f"method_scene: {method_scene}\n"
                    f"method_keywords: {method_keywords}\n"
                    f"method_step_count: {step_count}\n"
                    f"is_success: {is_success}"
                ),
            },
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="MethodSaver")
        llm_data = self._extract_json_object(llm_text)
        if llm_data:
            should_save = bool(llm_data.get("should_save", False))
            reason = str(llm_data.get("reason") or ("valuable_task" if should_save else "small_talk"))
            return should_save, reason[:80], "llm"

        # 2) 回退轻量启发式（仅在 LLM 不可用/失败时）
        compact_input = re.sub(r"[\s\W_]+", "", user_input.lower())
        greeting_keywords = ["你好", "您好", "hello", "hi", "hey", "在吗", "谢谢", "thank", "早上好", "晚上好"]
        has_greeting = any(k in user_input.lower() for k in greeting_keywords)
        if has_greeting and len(compact_input) <= 12:
            return False, "small_talk_fallback", "heuristic"
        return True, "valuable_task_fallback", "heuristic"

    def classify_interaction_mode(self, user_input: str) -> dict[str, Any]:
        text = (user_input or "").strip()
        if not text and not self._turn_vision_data_urls:
            return {"mode": "small_talk", "reason": "empty_input", "source": "heuristic", "confidence": 1.0}

        user_line = (
            f"用户输入：{text}"
            if text
            else "用户输入：（本消息主要为图片附件，无额外文字；请根据画面判断是闲聊问候类还是任务类。）"
        )
        messages = [
            {
                "role": "system",
                "content": (
                    "你是ARIA输入路由器。判断用户输入应该走哪条链路："
                    "- small_talk：纯粹的寒暄/问候/感谢/闲聊，没有任何具体目标或操作意图。"
                    "- task：含有明确目标或操作请求，即使文本中包含问候词（如\"你好，帮我发条消息\"属于 task）。"
                    "判断原则：有具体动作/目标/对象 → task；仅寒暄无目标 → small_talk。"
                    '仅输出JSON：{"mode":"small_talk|task","reason":"...","confidence":0-1}。'
                ),
            },
            {"role": "user", "content": self._user_content_with_optional_vision(user_line)},
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="TaskParser", reasoning_effort="minimal")
        data = self._extract_json_object(llm_text)
        if data:
            mode = str(data.get("mode", "")).strip().lower()
            if mode in ("small_talk", "task"):
                try:
                    confidence = float(data.get("confidence", 0.8))
                except Exception:
                    confidence = 0.8
                return {
                    "mode": mode,
                    "reason": str(data.get("reason") or ""),
                    "source": "llm",
                    "confidence": max(0.0, min(1.0, confidence)),
                }

        lowered = text.lower()
        compact = re.sub(r"[\s\W_]+", "", lowered)
        greeting_keywords = ["你好", "您好", "hello", "hi", "hey", "在吗", "谢谢", "thank", "早上好", "晚上好"]
        # fallback 只拦截极短纯问候（≤6字），避免含任务内容的请求被误判
        if any(k in lowered for k in greeting_keywords) and len(compact) <= 6:
            return {"mode": "small_talk", "reason": "greeting_fallback", "source": "heuristic", "confidence": 0.9}
        return {"mode": "task", "reason": "task_fallback", "source": "heuristic", "confidence": 0.7}

    def classify_interaction_mode_heuristic(self, user_input: str) -> dict[str, Any] | None:
        """
        无 LLM 的快速路由。若为高置信度寒暄（短问候/感谢等），返回 small_talk；
        返回 None 表示应继续走 plan_actions 等完整规划链路。
        含图片附件且无文字时不返回 small_talk（需模型看图）。
        """
        text = (user_input or "").strip()
        has_vision = bool(getattr(self, "_turn_vision_data_urls", None))
        if not text and not has_vision:
            return {"mode": "small_talk", "reason": "empty_input", "source": "heuristic", "confidence": 1.0}
        if not text:
            return None
        lowered = text.lower()
        compact = re.sub(r"[\s\W_]+", "", lowered)
        greeting_keywords = [
            "你好",
            "您好",
            "hello",
            "hi",
            "hey",
            "在吗",
            "谢谢",
            "thank",
            "早上好",
            "晚上好",
            "下午好",
            "晚安",
        ]
        if any(k in lowered for k in greeting_keywords) and len(compact) <= 6:
            return {"mode": "small_talk", "reason": "greeting_heuristic", "source": "heuristic", "confidence": 0.92}
        return None

    def derive_action_risk(self, action_type: str, risk: str) -> str:
        action_type = self._normalize_action_type_alias(action_type)
        normalized = (risk or "").strip().lower()
        if normalized not in ("low", "medium", "high"):
            normalized = "medium"
        if action_type in self.HIGH_RISK_ACTION_TYPES:
            return "high"
        return normalized

    def normalize_action_plan(self, plan: dict[str, Any]) -> dict[str, Any]:
        mode = str(plan.get("mode") or "").strip().lower()
        if mode == "clarify":
            choices = _normalize_clarify_choices(plan.get("choices"))
            return {
                "mode": "clarify",
                "summary": str(plan.get("summary") or "").strip(),
                "requires_confirmation": True,
                "actions": [],
                "requires_double_confirmation": False,
                "choices": choices,
            }
        if mode not in ("action", "qa", "small_talk"):
            mode = "qa"
        actions = plan.get("actions")
        if not isinstance(actions, list):
            actions = []
        normalized_actions = []
        for action in actions:
            if not isinstance(action, dict):
                continue
            action_type = self._normalize_action_type_alias(str(action.get("type") or "").strip())
            if not action_type or action_type not in self.ALLOWED_ACTION_TYPES:
                continue
            params = action.get("params") if isinstance(action.get("params"), dict) else {}
            risk = self.derive_action_risk(action_type, str(action.get("risk") or "medium"))
            if action_type == "shell_run":
                cmd = str(params.get("command") or "")
                if any(k in cmd.lower() for k in ["del ", "rm ", "shutdown", "format ", "reg delete"]):
                    risk = "high"
            normalized_actions.append(
                {
                    "type": action_type,
                    "target": str(action.get("target") or "").strip(),
                    "filters": action.get("filters") if isinstance(action.get("filters"), dict) else {},
                    "params": params,
                    "risk": risk,
                    "reason": str(action.get("reason") or "").strip(),
                }
            )
        normalized_actions = self._normalize_redundant_actions(normalized_actions)
        out: dict[str, Any] = {
            "mode": mode,
            "summary": str(plan.get("summary") or "").strip(),
            "requires_confirmation": True,
            "actions": normalized_actions,
            "requires_double_confirmation": self.requires_double_confirmation(normalized_actions),
        }
        cs = plan.get("complexity_score")
        if cs is not None:
            try:
                out["complexity_score"] = int(float(cs))
            except (TypeError, ValueError):
                pass
        cr = plan.get("complexity_reason")
        if cr is not None and str(cr).strip():
            out["complexity_reason"] = str(cr).strip()
        tr = plan.get("temporal_risk")
        if tr is not None and str(tr).strip():
            out["temporal_risk"] = str(tr).strip().lower()
        ot = plan.get("outcome_type")
        if ot is not None and str(ot).strip():
            out["outcome_type"] = str(ot).strip().lower()
        tf = plan.get("task_form")
        if tf is not None and str(tf).strip():
            tfn = str(tf).strip().lower()
            if tfn in ("local_execute", "web_information", "qa_only", "mixed"):
                out["task_form"] = tfn
        rr = plan.get("react_recommended")
        out["react_recommended"] = rr is True or (
            isinstance(rr, str) and rr.strip().lower() in ("1", "true", "yes", "on")
        )
        rrr = plan.get("react_recommend_reason")
        if rrr is not None and str(rrr).strip():
            out["react_recommend_reason"] = str(rrr).strip()
        rv = plan.get("react_computer_use_vision_recommended")
        out["react_computer_use_vision_recommended"] = rv is True or (
            isinstance(rv, str) and rv.strip().lower() in ("1", "true", "yes", "on")
        )
        return out

    def _user_explicitly_requests_web_information(self, text: str) -> bool:
        """用户是否明确要求上网检索/读网页（与仅操作本机 App 区分）。"""
        t = (text or "").strip()
        if not t:
            return False
        tl = t.lower()
        if re.search(r"https?://[^\s]+", t):
            return True
        for h in WEB_INFORMATION_MARKERS_CN:
            if h in t:
                return True
        for h in WEB_INFORMATION_MARKERS_EN:
            if h in tl:
                return True
        return False

    def _actions_are_web_research_only(self, actions: list[Any]) -> bool:
        """计划动作是否仅为联网检索（无桌面/微信等），用于微信直达执行时覆盖误规划。"""
        if not actions:
            return True
        for a in actions:
            if not isinstance(a, dict):
                return False
            t = self._normalize_action_type_alias(str(a.get("type") or ""))
            if t in ("web_understand", "web_fetch"):
                continue
            if t == "browser_open":
                raw = self._browser_open_raw_url(a)
                if self._is_search_engine_results_url(raw):
                    continue
                return False
            return False
        return True

    def _strip_contradictory_web_actions(
        self, user_text: str, plan: dict[str, Any]
    ) -> None:
        """当模型判定为本地执行或启发式为本地任务，且用户未明确要求上网时，移除误加的联网检索类动作。"""
        if not isinstance(plan, dict) or plan.get("mode") == "clarify":
            return
        if self._user_explicitly_requests_web_information(user_text):
            return
        tf = str(plan.get("task_form") or "").strip().lower()
        if tf == "web_information":
            return
        local_tf = tf in ("local_execute", "qa_only")
        local_h = self._user_intent_local_doc_edit(user_text)
        should_strip = local_tf or local_h or (tf == "mixed" and local_h)
        if not should_strip:
            return
        actions = plan.get("actions")
        if not isinstance(actions, list) or not actions:
            return
        kept: list[dict[str, Any]] = []
        for a in actions:
            if not isinstance(a, dict):
                continue
            t = self._normalize_action_type_alias(str(a.get("type") or ""))
            if t in ("web_understand", "web_fetch"):
                continue
            if t == "browser_open":
                raw = self._browser_open_raw_url(a)
                if self._is_search_engine_results_url(raw):
                    continue
            kept.append(a)
        if len(kept) == len(actions):
            return
        plan["actions"] = kept
        plan["requires_double_confirmation"] = self.requires_double_confirmation(kept)
        if not kept and str(plan.get("mode") or "") == "action":
            plan["mode"] = "qa"
            plan["requires_double_confirmation"] = False
            if not str(plan.get("summary") or "").strip():
                plan["summary"] = "已移除非必需的联网检索动作，请用对话回答或生成本机可执行步骤"

    def _browser_open_raw_url(self, action: dict[str, Any]) -> str:
        params = action.get("params") if isinstance(action.get("params"), dict) else {}
        u = str(params.get("url") or "").strip()
        if u:
            return u
        return str(action.get("target") or "").strip()

    def _is_concrete_browser_open_url(self, url: str) -> bool:
        if not url or len(url) > 2048:
            return False
        lowered = url.lower().strip()
        bogus_substrings = (
            "default browser",
            "the browser",
            "默认浏览器",
            "打开浏览器",
            "open browser",
        )
        if any(s in lowered for s in bogus_substrings):
            return False
        if " " in lowered and not lowered.startswith(("http://", "https://")):
            return False
        if lowered.startswith(("http://", "https://")):
            return True
        if "localhost" in lowered:
            return True
        if re.match(r"^\d{1,3}(\.\d{1,3}){3}(:\d+)?(/|$)", lowered):
            return True
        if re.match(r"^[a-z0-9][a-z0-9.-]*\.[a-z]{2,}(/|\?|$)", lowered):
            return True
        return False

    def _default_search_url(self, query: str) -> str:
        q = (query or "").strip() or "今日新闻"
        enc = quote_plus(q)
        if re.search(r"[\u4e00-\u9fff]", q):
            return f"https://www.baidu.com/s?wd={enc}"
        return f"https://www.bing.com/search?q={enc}"

    def _default_search_url_for_fetch(self, query: str) -> str:
        """供 web_understand/web_fetch 使用：优先 DuckDuckGo HTML 版，服务端可解析；Bing/百度 SERP 常为空壳或强 JS。"""
        q = (query or "").strip() or "news"
        enc = quote_plus(q)
        return f"https://html.duckduckgo.com/html/?q={enc}"

    def _bing_search_url_for_fetch(self, query: str) -> str:
        q = (query or "").strip() or "news"
        return f"https://www.bing.com/search?q={quote_plus(q)}"

    def _is_search_engine_results_url(self, url: str) -> bool:
        if not url:
            return False
        u = url.lower()
        if "baidu.com" in u and ("wd=" in u or "/s?" in u or "word=" in u):
            return True
        if "bing.com" in u and "search" in u:
            return True
        if "google." in u and "/search" in u:
            return True
        if "duckduckgo.com" in u:
            return True
        if "sogou.com" in u and ("query=" in u or "keyword=" in u):
            return True
        if "html.duckduckgo.com" in u and "/html/" in u:
            return True
        return False

    def _extract_search_query_from_url(self, url: str) -> str:
        try:
            p = urlparse(url)
            qs = parse_qs(p.query)
            for key in ("q", "wd", "word", "query", "keyword"):
                vals = qs.get(key)
                if vals and str(vals[0]).strip():
                    return unquote_plus(str(vals[0]).strip())
        except Exception:
            pass
        return ""

    def _understand_fetch_url_candidates(self, primary_url: str, question: str) -> list[str]:
        """搜索引擎页可能单源抓取失败，按序尝试主 URL、DuckDuckGo HTML、Bing。"""
        seen: set[str] = set()
        ordered: list[str] = []

        def add(u: str) -> None:
            u = (u or "").strip()
            if not u or u in seen:
                return
            seen.add(u)
            ordered.append(u)

        add(primary_url)
        q = self._extract_search_query_from_url(primary_url).strip()
        if not q and (question or "").strip():
            q = (question or "").strip()[:500]
        if q:
            add(f"https://html.duckduckgo.com/html/?q={quote_plus(q)}")
            add(self._bing_search_url_for_fetch(q))
        return ordered

    def _user_intent_browser_only(self, text: str) -> bool:
        """用户明确只要打开网页、不要服务端摘要时，不自动插入 web_understand。"""
        t = (text or "").strip()
        if not t:
            return False
        if any(p in t for p in ("只打开", "只要打开", "仅打开", "打开就行", "just open", "only open")):
            return True
        return False

    def _user_intent_local_doc_edit(self, text: str) -> bool:
        """用户要改工作区/本地 docx 样式时，勿自动插入 web_understand 搜索教程。"""
        t = (text or "").strip()
        if ".docx" not in t:
            return False
        if not any(k in t for k in ("字体", "字号", "加粗", "斜体", "格式", "样式", "修改", "改成", "调整")):
            return False
        if any(k in t for k in ("教程", "怎么改", "如何改", "教我", "教我改", "查资料", "上网搜", "搜索怎么")):
            return False
        return True

    def _heuristic_docx_style_plan(self, text: str) -> dict[str, Any] | None:
        """检测到「路径 + .docx + 样式词」时直接给出 file_write（docx_style），避免先走联网。"""
        t = (text or "").strip()
        if ".docx" not in t or not self._user_intent_local_doc_edit(t):
            return None
        path: str | None = None
        for m in re.finditer(r"([\w.\\/\-\u4e00-\u9fff]*(?:data|artifacts)[\w.\\/\-\u4e00-\u9fff]*\.docx)", t, re.I):
            candidate = m.group(1).strip().strip("\"'").replace("\\", "/")
            if ".." in candidate or candidate.startswith("http"):
                continue
            path = candidate
            break
        if not path:
            for m in re.finditer(r"(\S+\.docx)", t):
                candidate = m.group(1).strip().strip("\"'").replace("\\", "/")
                if candidate.startswith("http") or ".." in candidate:
                    continue
                path = candidate
                break
        if not path:
            return None
        docx_style: dict[str, Any] = {}
        tl = t.lower()
        if "加粗" in t or "bold" in tl:
            docx_style["bold"] = True
        if "斜体" in t or "italic" in tl:
            docx_style["italic"] = True
        msize = re.search(r"(\d{1,2})\s*(?:号|pt|磅)", t, re.I)
        if msize:
            docx_style["font_size_pt"] = int(msize.group(1))
        for font in ("宋体", "黑体", "楷体", "仿宋", "微软雅黑", "Times New Roman", "Arial", "Calibri"):
            if font in t:
                docx_style["font_name"] = font
                break
        actions = [
            {
                "type": "file_write",
                "target": path,
                "filters": {},
                "params": {"path": path, "docx_style": docx_style, "mode": "overwrite"},
                "risk": "low",
                "reason": "工作区 .docx 样式调整：直接写回文件，不联网查教程",
            }
        ]
        return {
            "mode": "action",
            "summary": f"将调整文档样式：{path}",
            "requires_confirmation": True,
            "actions": actions,
            "requires_double_confirmation": False,
        }

    def _prefer_qa_for_weather_without_explicit_web(self, text: str) -> bool:
        """无明确「打开/搜索/网上」等意图时，天气类问题走 qa，避免固定触发 web_understand + browser_open。"""
        t = (text or "").strip()
        if not t:
            return False
        if not any(
            k in t for k in ("天气", "气温", "下雨", "降温", "升温", "台风", "雾霾", "空气质量", "冷不冷", "热不热")
        ):
            return False
        if any(
            k in t
            for k in (
                "搜索",
                "打开",
                "链接",
                "网址",
                "浏览器",
                "网上",
                "联网",
                "搜一下",
                "查一下",
                "百度",
            )
        ):
            return False
        tl = t.lower()
        if "google" in tl or "bing" in tl:
            return False
        return True

    def _enrich_research_actions(self, user_input: str, plan: dict[str, Any]) -> None:
        """若计划仅有「打开搜索引擎」的 browser_open，则前置 web_understand（服务端可解析搜索页 + 模型作答）。"""
        if plan.get("mode") != "action":
            return
        actions = plan.get("actions")
        if not isinstance(actions, list) or not actions:
            return
        text = (user_input or "").strip()
        if self._user_intent_local_doc_edit(text):
            return
        if self._user_intent_browser_only(text):
            return
        for action in actions:
            if not isinstance(action, dict):
                return
            t = self._normalize_action_type_alias(str(action.get("type") or ""))
            if t in ("web_fetch", "web_understand"):
                return
            if t != "browser_open":
                return
            raw = self._browser_open_raw_url(action)
            if not self._is_search_engine_results_url(raw):
                return
        fetch_url = self._default_search_url_for_fetch(text)
        question = text[:2000] if len(text) > 2000 else text
        wu: dict[str, Any] = {
            "type": "web_understand",
            "target": fetch_url,
            "filters": {},
            "params": {"url": fetch_url, "question": question},
            "risk": "low",
            "reason": "服务端抓取搜索摘要并结合模型作答（优先可解析的搜索页；本机打开的地址可能为百度等，与抓取源可能不同）",
        }
        plan["actions"] = [wu] + list(actions)
        plan["requires_double_confirmation"] = self.requires_double_confirmation(plan["actions"])

    def _mend_browser_open_actions(self, user_input: str, plan: dict[str, Any]) -> None:
        """模型常把网址写在 target 却留空 params，或填「默认浏览器」；补成可打开的搜索引擎 URL（服务端不代抓网页）。"""
        if plan.get("mode") != "action":
            return
        actions = plan.get("actions")
        if not isinstance(actions, list):
            return
        text = (user_input or "").strip()
        for action in actions:
            if not isinstance(action, dict):
                continue
            t = self._normalize_action_type_alias(str(action.get("type") or ""))
            if t != "browser_open":
                continue
            raw = self._browser_open_raw_url(action)
            if self._is_concrete_browser_open_url(raw):
                continue
            url = self._default_search_url(text)
            params = action.get("params") if isinstance(action.get("params"), dict) else {}
            action["params"] = {**params, "url": url}
            action["target"] = url
            if not str(action.get("reason") or "").strip():
                action["reason"] = "在本机浏览器打开搜索结果页（ARIA 不在服务器上抓取网页，需在本机查看）"

    def _resolve_task_id_for_turn(self, reuse_task_id: str | None) -> str:
        t = (reuse_task_id or "").strip()
        if t:
            try:
                uuid.UUID(t)
                return t
            except ValueError:
                pass
        return str(uuid.uuid4())

    def _normalize_redundant_actions(self, actions: list[dict[str, Any]]) -> list[dict[str, Any]]:
        """按 automation/app_profiles/merge_rules.yaml 合并相邻重复链（如 open + send）。"""
        return normalize_actions_with_merge_rules(actions, self._normalize_action_type_alias)

    def _allowed_action_types_for_task_form(self, task_form: str) -> frozenset[str]:
        tf = (task_form or "").strip().lower()
        return self.TOOL_PROFILES.get(tf, self.TOOL_PROFILES["mixed"])

    def _allowed_action_types_for_workspace_mode(self, workspace_mode: str | None) -> frozenset[str]:
        mode = self._normalize_workspace_mode(workspace_mode or self.workspace_mode)
        return self.WORKSPACE_MODE_PROFILES.get(mode, self.WORKSPACE_MODE_PROFILES["aria"])

    def _apply_task_form_tool_allowlist(self, plan: dict[str, Any]) -> dict[str, Any]:
        if not isinstance(plan, dict):
            return {"mode": "qa", "summary": "invalid_plan", "actions": []}
        actions = plan.get("actions")
        if not isinstance(actions, list) or not actions:
            return plan
        task_form = str(plan.get("task_form") or "mixed").strip().lower()
        workspace_mode = self._normalize_workspace_mode(plan.get("workspace_mode") or self.workspace_mode)
        task_allowed = self._allowed_action_types_for_task_form(task_form)
        mode_allowed = self._allowed_action_types_for_workspace_mode(workspace_mode)
        allowed = task_allowed & mode_allowed
        filtered: list[dict[str, Any]] = []
        dropped_by_task_form: list[str] = []
        dropped_by_workspace_mode: list[str] = []
        for action in actions:
            if not isinstance(action, dict):
                continue
            normalized_type = self._normalize_action_type_alias(str(action.get("type") or ""))
            if normalized_type in allowed:
                if normalized_type != str(action.get("type") or ""):
                    action = dict(action)
                    action["type"] = normalized_type
                filtered.append(action)
            else:
                if normalized_type not in task_allowed:
                    dropped_by_task_form.append(normalized_type or "unknown")
                elif normalized_type not in mode_allowed:
                    dropped_by_workspace_mode.append(normalized_type or "unknown")
                else:
                    dropped_by_task_form.append(normalized_type or "unknown")
        if dropped_by_task_form or dropped_by_workspace_mode:
            plan["actions"] = filtered
            plan["workspace_mode"] = workspace_mode
            plan["requires_double_confirmation"] = self.requires_double_confirmation(filtered)
            if dropped_by_task_form:
                self.push_event(
                    "tool_allowlist",
                    "warning",
                    "TaskParser",
                    "已按 task_form 收敛可用动作集合",
                    {"task_form": task_form, "dropped_actions": dropped_by_task_form, "kept_count": len(filtered)},
                )
            if dropped_by_workspace_mode:
                self.push_event(
                    "workspace_mode_allowlist",
                    "warning",
                    "TaskParser",
                    "已按 workspace_mode 收敛可用动作集合",
                    {
                        "workspace_mode": workspace_mode,
                        "dropped_actions": dropped_by_workspace_mode,
                        "kept_count": len(filtered),
                    },
                )
            dropped_count = len(dropped_by_task_form) + len(dropped_by_workspace_mode)
            self.push_log("TaskParser", f"工具白名单过滤：移除 {dropped_count} 个动作", "warning")
            if not filtered and plan.get("mode") == "action":
                plan["mode"] = "clarify"
                plan["summary"] = (
                    "当前任务形态下可执行动作已被安全策略过滤。请确认是否切换为联网检索、或补充本机执行约束。"
                )
                plan["choices"] = [
                    {"id": "switch_local", "label": "改为本机执行（local_execute）"},
                    {"id": "switch_web", "label": "改为联网检索（web_information）"},
                ]
        else:
            plan["workspace_mode"] = workspace_mode
        return plan

    def plan_actions(self, user_input: str, dialogue_context: str = "") -> dict[str, Any]:
        text = (user_input or "").strip()
        if not text:
            return {"mode": "small_talk", "summary": "empty_input", "requires_confirmation": True, "actions": []}

        if self._prefer_qa_for_weather_without_explicit_web(text):
            return {
                "mode": "qa",
                "summary": "weather_or_ambient_query_no_explicit_web",
                "requires_confirmation": False,
                "actions": [],
                "requires_double_confirmation": False,
            }

        doc_style_plan = self._heuristic_docx_style_plan(text)
        if doc_style_plan:
            self._mend_browser_open_actions(text, doc_style_plan)
            return doc_style_plan

        # 注入当前真实日期，解决时间认知问题
        current_time_str = time.strftime("%Y年%m月%d日 %H:%M，%A")

        messages = [
            {
                "role": "system",
                "content": (
                    f"【当前时间】{current_time_str}（以此为准判别「今天」「明天」等）。\n\n"
                    "你是ARIA动作规划器。请把用户输入解析为：small_talk / qa / action / clarify。"
                    f"若用户提到具体日期，请先与此对比判断是过去/现在/未来。\n\n"
                    "【任务形式 task_form（必填）】在输出 JSON 前必须先归类，并填入 task_form 字段，取值只能是："
                    "local_execute（本机程序/自动化即可完成：微信发消息、打开桌面软件、改本地文档、终端命令、desktop/wechat/screen_* 等）；"
                    "web_information（用户明确要上网查资料/新闻/网页摘要/搜教程等，或给出了 http(s) 链接要求读取）；"
                    "qa_only（仅需对话回答，不要任何可执行动作）；"
                    "mixed（同时需要本机动作与联网信息，少见）。"
                    "规则：用户点名具体 App 并要求发送/打开/点击/输入等，一律 local_execute，禁止用 web_understand 代替本机执行；"
                    "仅当用户字面出现检索/网页意图或 task_form=web_information 时才输出 web_understand/web_fetch；"
                    "模棱两可时优先 local_execute，可配合 clarify 追问，不要默认联网搜索。\n\n"
                    "clarify 时可选附带 choices 数组（最多6条）："
                    '"choices":[{"id":"a","label":"选项文案"},...]，供用户点击，不必打字。'
                    "禁止编造执行结果。"
                    "【ReAct 建议 react_recommended（必填布尔）】判断是否适合在用户确认后启用逐步「Thought→Action→Observation」循环，而非一次性批量执行。"
                    "置 true：GUI/桌面自动化多步且状态依赖屏幕反馈、Computer Use（computer_*）、路径不明需试探、调试/排查、"
                    "或计划动作链较长且中间结果影响下一步。置 false：单一确定动作（如一次 file_write、参数明确的 wechat_send_message）、"
                    "纯 web_fetch/web_understand、低风险可预测流水线。须填 react_recommend_reason（一句中文理由）。"
                    "react_computer_use_vision_recommended：仅当 react_recommended 为 true 且任务确实需要每步看桌面截图辅助推理时为 true；"
                    "否则 false（避免无谓多模态成本）。含 computer_* 动作时通常应为 true。"
                    "仅输出JSON："
                    '{"mode":"small_talk|qa|action|clarify",'
                    '"summary":"...",'
                    '"task_form":"local_execute|web_information|qa_only|mixed",'
                    '"react_recommended":false,'
                    '"react_recommend_reason":"...",'
                    '"react_computer_use_vision_recommended":false,'
                    '"complexity_score":1,'
                    '"complexity_reason":"一句话说明评分理由",'
                    '"temporal_risk":"high",'
                    '"outcome_type":"time_bound",'
                    '"choices":[],'
                    '"actions":[{"type":"...","target":"...","filters":{},"params":{},"risk":"low|medium|high","reason":"..."}]}'
                    "可用动作类型示例：kb_delete_all,kb_delete_by_keyword,kb_delete_low_quality,conversation_new,shell_run,file_write,file_move,file_delete,browser_open,browser_click,browser_type,browser_find,browser_hover,browser_select,browser_upload,browser_scroll,browser_wait,browser_js,browser_press,media_summarize,desktop_open_app,desktop_hotkey,desktop_type,desktop_sequence,wechat_send_message,wechat_open_chat,wechat_check_login,screen_ocr,screen_find_text,screen_click_text,web_fetch,web_understand。"
                    "能力边界（须严格遵守，勿向用户承诺未启用的能力）："
                    + browser_driver.capability_summary_for_planner()
                    + desktop_uia.capability_summary_for_planner()
                    + screen_ocr.get_capability_summary()
                    + "desktop_open_app：params.app 或顶层 target 填应用名/路径；Windows 下会扫描本机桌面（含 OneDrive 桌面）、开始菜单与微信/WPS 等常见安装路径，再启动；仍失败时让用户提供 .exe 完整路径。"
                    "本地文档/表格：能确定相对工作区的路径与内容时优先 file_write（params.path、params.content、mode overwrite|append）；格式/路径/是否覆盖不明时用 mode=clarify 追问，勿直接拒绝。"
                    "若用户已给出工作区相对路径的 .docx，且仅要求改字体/字号/加粗/斜体等样式：必须 mode=action，使用 file_write，params.path 为原文档路径，params.docx_style 为对象（font_name、font_size_pt、bold、italic）；禁止同时输出 web_understand 查教程；仅当用户明确要教程或原理时才 web_understand。"
                    'Office 二进制（服务端生成，用户确认执行后可下载）：path 以 .docx 结尾时 params.content 为正文、可选 params.title；.xlsx 时优先 params.rows 为二维数组[["列1","列2"],["a","b"]]，否则用 content 多行、制表符或逗号分列；.pptx 时 params.title 与 params.bullets 字符串数组或 content 多行（首行可作标题）。路径建议 data/artifacts/文件名。'
                    "禁止因「没有 Word/Office」就声明无法完成：须列举替代（记事本、Markdown、WPS、LibreOffice、VS Code、工作区 file_write 写 .md/.txt/.csv 等）。"
                    "用户可能已通过 Web 上传文件：【本轮输入】中「抽取摘要」含 txt/md/pdf/docx/xlsx/pptx 等正文摘录及 data/artifacts/uploads/… 相对路径；图片除摘要外，在同轮多模态请求中会以 image 形式一并传入（请直接根据画面回答/规划）。可据此做总结、改写、或 file_write 写入新文件；勿声称「已保存到用户本机文档」除非实际执行了 file_write。禁止虚构未执行动作的结果。"
                    "缺软件或不会配置时：可输出 web_understand（检索安装/配置步骤，params.question 写清系统与软件名）或 clarify 让用户选已装软件；不要只给失败理由。"
                    "检测到用户要打开的应用可能未安装时：应输出 clarify 模式询问用户「未找到 XX 应用，是否打开网页版？」并给出网页版 URL 建议（如微信→https://web.wechat.com，WPS→https://www.kdocs.cn，Office→https://www.office.com，Adobe→https://www.adobe.com 等）。"
                    "browser_open：若 ARIA_PLAYWRIGHT=1 且已安装 Playwright 与 Chromium，则用受控浏览器导航；否则用系统默认浏览器打开 URL。"
                    "browser_click/browser_type：Playwright 启用且包已安装时，params.selector 为 CSS 选择器，可选 params.url 先打开页面再操作；未启用时为模拟，不得声称已点击或已输入。"
                    "browser_find：params.selector 为 CSS 选择器，可选 params.text_contains 过滤文本；返回元素列表及位置信息。"
                    "browser_hover：鼠标悬停到元素，params.selector 为 CSS 选择器，用于触发下拉菜单等。"
                    "browser_select：选择下拉框，params.selector 为 CSS 选择器，params.value 为选项值或文本。"
                    "browser_upload：上传文件，params.selector 为文件输入框选择器，params.file_path 为文件路径（需先用 file_write 生成）。"
                    "browser_scroll：滚动页面，params.selector 可选（不填则滚动到底部）。"
                    "browser_wait：等待元素出现，params.selector 必填，params.timeout_ms 可选（默认 30000ms），适合动态加载的页面。"
                    "browser_js：执行自定义 JavaScript，params.script 为代码，返回结果（需谨慎使用，避免安全风险）。"
                    "browser_press：Playwright 下模拟键盘，params.key 如 Playwright 键名 Enter/Tab/Escape 等，可选 params.selector 先聚焦元素。"
                    "media_summarize：视频摘要。params.url（YouTube 等公开页，服务端 oEmbed 元数据）或 params.path（工作区内 .mp4/.webm，需本机 ffmpeg 抽帧 + 多模态）；勿承诺完整听懂对白若无字幕。"
                    "desktop_hotkey/desktop_type/desktop_sequence：仅当 ARIA_DESKTOP_UIA=1 且 Windows 上已安装 pywinauto 时为真实注入，否则为模拟。"
                    + load_planner_fragment("planner_desktop_apps_fragment.txt")
                    + self._computer_use_capability_summary()
                    + "screen_ocr/screen_find_text/screen_click_text：屏幕 OCR 识别与操作；须安装 Tesseract OCR，未安装时执行会失败；本机 Tesseract 可用性见上方预检结果。"
                    "web_fetch：服务端抓取并抽取正文（params.url）；web_understand：抓取后由 params.question 指定要让模型回答的问题，并生成理解与摘要（需 API Key）。"
                    "仅当 task_form=web_information 或 mixed（且联网部分确有必要）时，才输出 web_understand/web_fetch；"
                    "task_form=local_execute 时 actions 中不得出现 web_understand/web_fetch，也不得用「打开搜索引擎」代替本机动作。"
                    "用户要查资料、搜材料、总结要点、对比依据且依赖外部信息且 task_form 已判为 web_information 时：必须包含 web_understand（params.url 为可访问页面 URL，params.question 填用户任务），不要只输出 browser_open。"
                    "若用户给了具体文章/文档链接并要求摘要，用 web_understand 或 web_fetch；可附加 browser_open 方便本机查看。"
                    "常识、概念解释、无需「今日/最新/实时」联网信息时：用 mode=qa、task_form=qa_only、actions 为空数组，不要生成打开浏览器的动作。"
                    "多步示例（仅 web_information）：web_understand（搜索或文章 URL + question）后可跟 browser_open（同一或不同 URL）。"
                    "browser_open 的 target 或 params.url 必须是可打开的具体网址，禁止「默认浏览器」等描述；搜索类应给出完整搜索 URL（如 https://www.baidu.com/s?wd=关键词）。"
                ),
            },
            {
                "role": "user",
                "content": self._user_content_with_optional_vision(
                    
                        f"【本会话近期对话】\n{(dialogue_context or '').strip()}\n\n【本轮输入】\n{text}"
                        if (dialogue_context or "").strip()
                        else text
                    
                ),
            },
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="TaskParser")
        llm_plan = self._extract_json_object(llm_text)
        plan = self.normalize_action_plan(llm_plan) if llm_plan else {}
        plan = self._apply_task_form_tool_allowlist(plan)
        if plan.get("mode") != "clarify":
            self._strip_contradictory_web_actions(text, plan)
        if plan.get("mode") == "clarify":
            return plan
        if plan and plan.get("mode") not in ("qa", "clarify"):
            acts_early = plan.get("actions") if isinstance(plan.get("actions"), list) else []
            if acts_early:
                self._mend_browser_open_actions(text, plan)
                self._enrich_research_actions(text, plan)
                return plan
            # small_talk / action 但动作列表为空：勿提前 return，继续走微信与其它规则兜底
        if plan and plan.get("mode") == "qa" and plan.get("actions"):
            self._mend_browser_open_actions(text, plan)
            self._enrich_research_actions(text, plan)
            return plan

        lowered = text.lower()
        if any(k in lowered for k in ["清理", "清空", "删除", "移除"]) and any(
            k in text for k in ["知识库", "方法论", "经验"]
        ):
            action_type = "kb_delete_low_quality"
            if "清空" in text:
                action_type = "kb_delete_all"
            actions = [
                {
                    "type": action_type,
                    "target": "knowledge_base",
                    "filters": {},
                    "params": {},
                    "risk": self.derive_action_risk(action_type, "medium"),
                    "reason": "用户请求清理知识库",
                }
            ]
            plan_kb = {
                "mode": "action",
                "summary": "识别为知识库清理操作",
                "requires_confirmation": True,
                "actions": actions,
                "requires_double_confirmation": self.requires_double_confirmation(actions),
            }
            self._mend_browser_open_actions(text, plan_kb)
            return plan_kb
        if any(k in lowered for k in ["执行命令", "运行命令", "run command", "terminal"]) and any(
            s in text for s in ["`", "cmd", "powershell", "python", "pip "]
        ):
            cmd = text
            m = re.search(r"`([^`]+)`", text)
            if m:
                cmd = m.group(1).strip()
            actions = [
                {
                    "type": "shell_run",
                    "target": "terminal",
                    "filters": {},
                    "params": {"command": cmd, "cwd": "."},
                    "risk": self.derive_action_risk("shell_run", "medium"),
                    "reason": "用户请求执行终端命令",
                }
            ]
            plan_sh = {
                "mode": "action",
                "summary": "识别为终端命令执行",
                "requires_confirmation": True,
                "actions": actions,
                "requires_double_confirmation": self.requires_double_confirmation(actions),
            }
            self._mend_browser_open_actions(text, plan_sh)
            return plan_sh
        if any(k in lowered for k in ["创建文件", "写入文件", "保存到文件", "write file"]):
            actions = [
                {
                    "type": "file_write",
                    "target": "filesystem",
                    "filters": {},
                    "params": {"path": "output.txt", "content": text, "mode": "overwrite"},
                    "risk": "low",
                    "reason": "用户请求写入文件",
                }
            ]
            plan_fw = {
                "mode": "action",
                "summary": "识别为文件写入操作",
                "requires_confirmation": True,
                "actions": actions,
                "requires_double_confirmation": False,
            }
            self._mend_browser_open_actions(text, plan_fw)
            return plan_fw
        if self._user_explicitly_requests_web_information(text):
            fetch_url = self._default_search_url_for_fetch(text)
            open_url = self._default_search_url(text)
            q = text[:2000] if len(text) > 2000 else text
            actions = [
                {
                    "type": "web_understand",
                    "target": fetch_url,
                    "filters": {},
                    "params": {"url": fetch_url, "question": q},
                    "risk": "low",
                    "reason": "服务端抓取可解析搜索页（DuckDuckGo HTML 等）并结合模型摘要（对话内查看；无法自动点开 SERP 具体条目）",
                },
                {
                    "type": "browser_open",
                    "target": open_url,
                    "filters": {},
                    "params": {"url": open_url},
                    "risk": "low",
                    "reason": "在本机浏览器打开搜索引擎（中文默认百度）",
                },
            ]
            plan_ws = {
                "mode": "action",
                "summary": "识别为联网检索：先摘要后本机打开搜索页",
                "requires_confirmation": True,
                "actions": actions,
                "requires_double_confirmation": self.requires_double_confirmation(actions),
            }
            self._mend_browser_open_actions(text, plan_ws)
            self._enrich_research_actions(text, plan_ws)
            return plan_ws
        return {"mode": "qa", "summary": "未识别为可执行动作", "requires_confirmation": True, "actions": []}

    def format_clarify_plan_for_user(self, plan: dict[str, Any]) -> str:
        s = str(plan.get("summary") or "").strip()
        choices = plan.get("choices") if isinstance(plan.get("choices"), list) else []
        choice_lines = ""
        if choices:
            choice_lines = "\n\n快捷选项（也可点击下方按钮）：\n" + "\n".join(
                f"- {c.get('label', c.get('id', ''))}" for c in choices if isinstance(c, dict)
            )
        tail = (
            "\n\n说明：可在聊天输入框旁用回形针上传文件（保存到工作区 data/artifacts/uploads/…）；"
            "若要在本机得到可下载生成文件，请通过「确认执行」后的 file_write 写入 ARIA 工作区，"
            "执行完成后回复里会给出 /api/workspace_file?path=… 下载链接。"
        )
        if not s:
            return "为继续推进，请补充更具体的需求（例如保存路径、文件格式、使用的软件名称等）。" + choice_lines + tail
        return "为准确帮你完成，请先确认或补充下列信息（可直接逐条回复）：\n\n" + s + choice_lines + tail

    def format_action_plan_for_user(self, plan: dict[str, Any]) -> str:
        actions = plan.get("actions") or []
        if not actions:
            return "我理解了你的请求，但当前未识别出可执行动作。请再具体一点，比如“删除关键词为xx的方法论”。"
        lines = ["我已理解为可执行操作，执行前请确认："]
        for idx, action in enumerate(actions, start=1):
            lines.append(
                f"{idx}. type={action.get('type')} target={action.get('target')} risk={action.get('risk')} reason={action.get('reason')}"
            )
        if self.requires_double_confirmation(actions):
            lines.append("检测到高风险动作：需要二次确认。先回复“确认执行”，再回复“二次确认”后才会实际执行。")
        else:
            lines.append("请回复“确认执行”后我再实际执行。")
        return "\n".join(lines)

    def requires_double_confirmation(self, actions: list[dict[str, Any]]) -> bool:
        return self.evaluate_action_risk_level(actions) == "high"

    def evaluate_action_risk_level(self, actions: list[dict[str, Any]]) -> str:
        """
        三层风险策略：
        - safe: 可自动执行
        - medium: 需一次确认
        - high: 需二次确认
        """
        level = "safe"
        for a in actions or []:
            if not isinstance(a, dict):
                continue
            action_type = self._normalize_action_type_alias(str(a.get("type") or ""))
            risk = str(a.get("risk") or "medium").strip().lower()
            if risk == "high" or action_type in self.HIGH_RISK_ACTION_TYPES:
                return "high"
            if risk == "medium" or action_type in self.USER_GATE_ACTION_TYPES:
                level = "medium"
        return level

    def actions_require_user_gate(self, actions: list[dict[str, Any]]) -> bool:
        """本地写文件、删文件、终端、打开桌面应用等须用户确认后才执行，不参与自动执行。"""
        for a in actions or []:
            if not isinstance(a, dict):
                continue
            t = self._normalize_action_type_alias(str(a.get("type") or ""))
            if t in self.USER_GATE_ACTION_TYPES:
                return True
        return False

    def taor_action_blocked_for_dispatch(self, action: dict[str, Any]) -> tuple[bool, str, str]:
        """
        TAOR 自主循环无 UI 确认链：默认拦截与主流程「需确认」同类的动作。
        受控环境可设 ARIA_TAOR_ALLOW_GATED_ACTIONS=1 放行（仍尊重 PLAN 只读模式）。
        返回 (blocked, error_code, message)。
        """
        if not isinstance(action, dict):
            return True, "invalid_action", "动作格式无效"
        raw_type = str(action.get("type") or "")
        action_type = self._normalize_action_type_alias(raw_type)
        if not action_type or action_type not in self.ALLOWED_ACTION_TYPES:
            return False, "", ""
        pm = self.permission_model
        if pm.is_readonly_only() and not pm.allows_under_plan_mode(action_type):
            return (
                True,
                "permission_denied",
                f"当前权限级别（plan）不允许执行 {action_type}",
            )
        allow_gated = (os.getenv("ARIA_TAOR_ALLOW_GATED_ACTIONS", "0").strip().lower() in ("1", "true", "yes"))
        if allow_gated:
            return False, "", ""
        risk = self.derive_action_risk(action_type, str(action.get("risk") or "medium"))
        bad_eval = self.evaluate_action_risk_level([action]) != "safe"
        bad_conf = pm.requires_confirmation(action_type, risk)
        if bad_eval or bad_conf:
            msg = (
                "TAOR 模式下不允许未经确认的动作；请关闭 TAOR 使用主流程，或在受控环境设置 ARIA_TAOR_ALLOW_GATED_ACTIONS=1。"
            )
            return (True, "user_gate_blocked" if bad_eval else "confirmation_required", msg)
        return False, "", ""

    def _ensure_safe_path(self, raw_path: str) -> Path:
        p = Path(raw_path or "").expanduser()
        if not p.is_absolute():
            p = (self.allowed_work_root / p).resolve()
        else:
            p = p.resolve()
        root = self.allowed_work_root.resolve()
        if root != p and root not in p.parents:
            raise ValueError(f"path_outside_allowlist: {p}")
        return p

    def _sanitize_shell_command(self, command: str) -> str:
        cmd = (command or "").strip()
        if not cmd:
            raise ValueError("empty_command")
        reason = shell_command_blocked_reason(cmd)
        if reason:
            raise ValueError(f"blocked_command:{reason}")
        return cmd

    def _capture_screenshot(self, prefix: str) -> list[str]:
        """满足以下任一即截屏：环境变量 ARIA_ACTION_SCREENSHOT，或本次执行请求中 action_screenshots 为真。"""
        env_on = (os.getenv("ARIA_ACTION_SCREENSHOT") or "").strip().lower() in ("1", "true", "yes", "on")
        client_on = bool(getattr(self, "action_screenshots_for_execution", False))
        if not env_on and not client_on:
            return []
        min_iv = int(os.getenv("ARIA_SCREENSHOT_MIN_INTERVAL_MS", "0") or "0")
        if min_iv > 0:
            now = time.time()
            if now - self._screenshot_last_ts < min_iv / 1000.0:
                return []
            self._screenshot_last_ts = now
        ts = int(time.time() * 1000)
        out_dir = self.allowed_work_root / "data" / "artifacts" / "screenshots"
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / f"{prefix}_{ts}.png"
        try:
            from PIL import ImageGrab  # type: ignore

            img = ImageGrab.grab()
            img.save(str(out_path))
            return [str(out_path)]
        except Exception:
            return []

    def _normalize_action_type_alias(self, action_type: str) -> str:
        mapping = {
            "terminal_run": "shell_run",
            "cmd_run": "shell_run",
            "file_create": "file_write",
            "file_rename": "file_move",
            "browser_navigate": "browser_open",
            "app_open": "desktop_open_app",
            "http_get": "web_fetch",
            "fetch_url": "web_fetch",
            "web_read": "web_fetch",
            "read_webpage": "web_fetch",
            "summarize_url": "web_understand",
            "web_summarize": "web_understand",
            "browser_press_key": "browser_press",
        }
        t = (action_type or "").strip()
        return mapping.get(t, t)

    def _capability_unavailable_result(self, action_type: str, capability: str, hint: str = "") -> dict[str, Any]:
        return {
            "success": False,
            "message": "unavailable_capability",
            "error_code": "unavailable_capability",
            "action_type": action_type,
            "capability": capability,
            "retryable": True,
            "needs_manual_takeover": False,
            "stderr": hint or capability,
            "artifacts": [],
            "screenshots": [],
        }

    def _verify_action_result(self, action_type: str, action: dict[str, Any], result: dict[str, Any]) -> dict[str, Any]:
        try:
            if action_type == "file_write" and bool(result.get("success")):
                path = str((action.get("params") or {}).get("path") or action.get("target") or "").strip()
                if path:
                    ok = self._ensure_safe_path(path).is_file()
                    return {"checked": True, "ok": bool(ok), "method": "file_exists"}
            if action_type == "file_move" and bool(result.get("success")):
                dst = str((action.get("params") or {}).get("dst") or "").strip()
                if dst:
                    ok = self._ensure_safe_path(dst).exists()
                    return {"checked": True, "ok": bool(ok), "method": "dst_exists"}
            if action_type == "file_delete" and bool(result.get("success")):
                path = str((action.get("params") or {}).get("path") or action.get("target") or "").strip()
                if path:
                    ok = not self._ensure_safe_path(path).exists()
                    return {"checked": True, "ok": bool(ok), "method": "path_not_exists"}
        except Exception:
            return {"checked": True, "ok": False, "method": "verification_exception"}
        return {"checked": False, "ok": None, "method": "none"}

    def _normalize_action_result(
        self,
        action_type: str,
        action: dict[str, Any],
        raw_result: dict[str, Any],
    ) -> dict[str, Any]:
        result = dict(raw_result or {})
        msg = str(result.get("message") or "")
        if msg.endswith("_simulated") or "_simulated:" in msg:
            hint = "当前能力未启用或依赖未安装，已拒绝用模拟成功掩盖真实执行状态"
            if action_type.startswith("browser_"):
                hint = "browser 实际控制不可用，请启用 ARIA_PLAYWRIGHT=1 并安装 playwright"
            elif action_type.startswith("desktop_"):
                hint = "desktop 实际注入不可用，请启用 ARIA_DESKTOP_UIA=1 并安装 pywinauto"
            return self._capability_unavailable_result(action_type, f"{action_type}_runtime", hint=hint)
        result.setdefault("error_code", "" if result.get("success") is not False else "execution_failed")
        _RETRYABLE_ERROR_CODES = {"timeout", "network_error", "transient_error", "execution_exception"}
        _error_code = result.get("error_code") or ""
        result.setdefault(
            "retryable",
            _error_code in _RETRYABLE_ERROR_CODES or bool(result.get("retryable") is True),
        )
        result.setdefault("needs_manual_takeover", False)
        result["verification"] = self._verify_action_result(action_type, action, result)
        return self.interaction_core.normalize_result(action_type, result)

    def execute_actions(
        self,
        actions: list[dict[str, Any]],
        conversation_id: str,
        request_id: str,
        methodology_manager: Any,
        conversation_manager: Any,
    ) -> dict[str, Any]:
        report = []
        merged = self._normalize_redundant_actions(list(actions or []))
        bounded_actions = merged[: self.max_action_steps]
        for idx, action in enumerate(bounded_actions, start=1):
            if self.is_cancelled(request_id):
                break
            raw_type = str(action.get("type") or "")
            action_type = self._normalize_action_type_alias(raw_type)
            handler = self.action_registry.get(action_type)
            started = time.time()
            row = {
                "step_id": idx,
                "action": action_type,
                "input": action,
                "status": "error",
                "duration_ms": 0,
                "error_code": "",
                "retryable": False,
                "needs_manual_takeover": False,
                "verification": {"checked": False, "ok": None, "method": "none"},
                "stdout": "",
                "stderr": "",
                "artifacts": [],
                "screenshots": [],
                "strategy_path": "rule_path",
                "confidence": 0.0,
                "safe_block_reason": "",
                "fallback_used": False,
                "decision_trace": [],
            }
            if not handler:
                row["stderr"] = "unsupported_action"
                row["error_code"] = "unsupported_action"
                row["result"] = {"success": False, "message": "unsupported_action", "error_code": "unsupported_action"}
                row["duration_ms"] = int((time.time() - started) * 1000)
                report.append(row)
                continue
            self.push_event(
                "computer_action",
                "running",
                "TaskParser",
                f"执行动作 #{idx}: {action_type}",
                {"step_id": idx, "action_type": action_type, "input": action},
            )
            try:
                action_ctx = dict(action)
                action_ctx["_request_id"] = request_id
                result = handler(action_ctx, conversation_id, methodology_manager, conversation_manager)
                base = result if isinstance(result, dict) else {"success": True, "output": result}
                # 统一内核：规则失败后，浏览器动作自动尝试文本语义兜底。
                if self.interaction_core.should_try_browser_fallback(action_type, action, base):
                    fb = self.interaction_core.try_browser_fallback(action_type, action, base)
                    if isinstance(fb, dict):
                        base = fb
                row["result"] = self._normalize_action_result(action_type, action, base)
                row["status"] = "success" if row["result"].get("success") is not False else "error"
                row["error_code"] = str(row["result"].get("error_code") or "")
                row["retryable"] = bool(row["result"].get("retryable"))
                row["needs_manual_takeover"] = bool(row["result"].get("needs_manual_takeover"))
                row["verification"] = row["result"].get("verification") or row["verification"]
                row["stdout"] = str(row["result"].get("stdout") or row["result"].get("message") or "")
                row["stderr"] = str(row["result"].get("stderr") or "")
                row["artifacts"] = row["result"].get("artifacts") or []
                row["screenshots"] = row["result"].get("screenshots") or []
                row["strategy_path"] = str(row["result"].get("strategy_path") or "rule_path")
                try:
                    row["confidence"] = float(row["result"].get("confidence", 0.0) or 0.0)
                except Exception:
                    row["confidence"] = 0.0
                row["safe_block_reason"] = str(row["result"].get("safe_block_reason") or "")
                row["fallback_used"] = bool(row["result"].get("fallback_used"))
                row["decision_trace"] = row["result"].get("decision_trace") or []
            except Exception as e:
                row["stderr"] = str(e)
                row["error_code"] = "execution_exception"
                row["retryable"] = True
                row["result"] = {
                    "success": False,
                    "message": str(e),
                    "error_code": "execution_exception",
                    "retryable": True,
                }
            row["duration_ms"] = int((time.time() - started) * 1000)

            # ── 重试 & outcome_state 逻辑 ──────────────────────────────── #
            attempt = 1
            max_retries = int(getattr(self, "max_action_retries", 0) or 0)
            while True:
                row["attempt"] = attempt
                success = row["status"] == "success"
                retryable = bool(row.get("retryable"))
                verify_ok = (row.get("verification") or {}).get("ok")

                # 取消检测（重试前）
                if not success and retryable and attempt <= max_retries:
                    if self.is_cancelled(request_id):
                        row["outcome_state"] = "cancelled"
                        row["recovery_decision"] = "cancelled_by_user"
                        break

                if success:
                    # 验证失败 → verify_failed
                    if verify_ok is False:
                        row["outcome_state"] = "verify_failed"
                        row["needs_manual_takeover"] = True
                        row["recovery_decision"] = "manual_takeover"
                        if attempt <= max_retries:
                            # 追加一条 verify_failed 记录后继续
                            report.append(dict(row))
                            # 重新执行
                            attempt += 1
                            started2 = time.time()
                            try:
                                action_ctx2 = dict(action)
                                action_ctx2["_request_id"] = request_id
                                result2 = handler(action_ctx2, conversation_id, methodology_manager, conversation_manager)
                                base2 = result2 if isinstance(result2, dict) else {"success": True, "output": result2}
                                row = dict(report[-1])  # 基于上一条
                                row["result"] = self._normalize_action_result(action_type, action, base2)
                                row["status"] = "success" if row["result"].get("success") is not False else "error"
                                row["error_code"] = str(row["result"].get("error_code") or "")
                                row["retryable"] = bool(row["result"].get("retryable"))
                                row["needs_manual_takeover"] = bool(row["result"].get("needs_manual_takeover"))
                                row["verification"] = row["result"].get("verification") or row["verification"]
                                row["stdout"] = str(row["result"].get("stdout") or row["result"].get("message") or "")
                                row["stderr"] = str(row["result"].get("stderr") or "")
                                row["artifacts"] = row["result"].get("artifacts") or []
                                row["screenshots"] = row["result"].get("screenshots") or []
                                row["duration_ms"] = int((time.time() - started2) * 1000)
                                row["attempt"] = attempt
                            except Exception as e2:
                                row["stderr"] = str(e2)
                                row["status"] = "error"
                                row["error_code"] = "execution_exception"
                                row["retryable"] = True
                                row["duration_ms"] = int((time.time() - started2) * 1000)
                            continue
                    else:
                        row["outcome_state"] = "success"
                        row["recovery_decision"] = ""
                    break
                elif retryable and attempt <= max_retries:
                    row["outcome_state"] = "recoverable_error"
                    row["recovery_decision"] = f"retry_scheduled:{attempt}"
                    report.append(dict(row))
                    # 重新执行
                    attempt += 1
                    started2 = time.time()
                    try:
                        action_ctx2 = dict(action)
                        action_ctx2["_request_id"] = request_id
                        result2 = handler(action_ctx2, conversation_id, methodology_manager, conversation_manager)
                        base2 = result2 if isinstance(result2, dict) else {"success": True, "output": result2}
                        row = {
                            "step_id": idx, "action": action_type, "input": action,
                            "status": "success" if base2.get("success") is not False else "error",
                            "duration_ms": int((time.time() - started2) * 1000),
                            "error_code": str(base2.get("error_code") or ""),
                            "retryable": bool(base2.get("retryable")),
                            "needs_manual_takeover": bool(base2.get("needs_manual_takeover")),
                            "verification": base2.get("verification") or {"checked": False, "ok": None, "method": "none"},
                            "stdout": str(base2.get("stdout") or base2.get("message") or ""),
                            "stderr": str(base2.get("stderr") or ""),
                            "artifacts": base2.get("artifacts") or [],
                            "screenshots": base2.get("screenshots") or [],
                            "strategy_path": str(base2.get("strategy_path") or "rule_path"),
                            "confidence": float(base2.get("confidence") or 0.0),
                            "safe_block_reason": str(base2.get("safe_block_reason") or ""),
                            "fallback_used": bool(base2.get("fallback_used")),
                            "decision_trace": base2.get("decision_trace") or [],
                            "result": self._normalize_action_result(action_type, action, base2),
                            "attempt": attempt,
                        }
                    except Exception as e2:
                        row = {
                            "step_id": idx, "action": action_type, "input": action,
                            "status": "error", "duration_ms": int((time.time() - started2) * 1000),
                            "error_code": "execution_exception", "retryable": True,
                            "needs_manual_takeover": False,
                            "verification": {"checked": False, "ok": None, "method": "none"},
                            "stdout": "", "stderr": str(e2), "artifacts": [], "screenshots": [],
                            "strategy_path": "rule_path", "confidence": 0.0,
                            "safe_block_reason": "", "fallback_used": False, "decision_trace": [],
                            "result": {"success": False, "message": str(e2), "error_code": "execution_exception"},
                            "attempt": attempt,
                        }
                    continue
                else:
                    row["outcome_state"] = "failed"
                    row["recovery_decision"] = "give_up"
                    break

            # ── 取消检测（执行后）──────────────────────────────────────── #
            if self.is_cancelled(request_id) and row.get("outcome_state") not in ("cancelled",):
                row["outcome_state"] = "cancelled"
                row["recovery_decision"] = "cancelled_by_user"

            self.push_event(
                "computer_action",
                "success" if row["status"] == "success" else "error",
                "TaskParser",
                f"动作 #{idx} 结束: {action_type}",
                {
                    "step_id": idx,
                    "action_type": action_type,
                    "status": row["status"],
                    "error_code": row["error_code"],
                    "duration_ms": row["duration_ms"],
                    "stdout": row.get("stdout", "")[:300],
                    "stderr": row.get("stderr", "")[:300],
                    "artifacts": row.get("artifacts", []),
                    "screenshots": row.get("screenshots", []),
                    "strategy_path": row.get("strategy_path", "rule_path"),
                    "confidence": row.get("confidence", 0.0),
                    "safe_block_reason": row.get("safe_block_reason", ""),
                    "fallback_used": row.get("fallback_used", False),
                },
            )
            report.append(row)
            if self.is_cancelled(request_id):
                break

        success_count = sum(1 for r in report if r.get("status") == "success")
        unavailable_count = sum(
            1
            for r in report
            if str((r.get("result") or {}).get("error_code") or r.get("error_code") or "") == "unavailable_capability"
        )
        manual_takeover = any(r.get("needs_manual_takeover") for r in report)
        return {
            "success_count": success_count,
            "total": len(report),
            "failed_count": max(0, len(report) - success_count),
            "unavailable_count": unavailable_count,
            "manual_takeover_required": manual_takeover,
            "report": report,
        }

    def execute_based_on_complexity(
        self,
        plan: dict[str, Any],
        conversation_id: str,
        request_id: str,
        methodology_manager: Any,
        conversation_manager: Any,
    ) -> dict[str, Any]:
        """基于复杂度评分选择执行路径"""
        score = int(plan.get("complexity_score", 3))
        temporal = str(plan.get("temporal_risk", "low")).lower()
        outcome = str(plan.get("outcome_type", "stable")).lower()

        # 确保 plan 中有 actions
        actions = plan.get("actions") if isinstance(plan.get("actions"), list) else []
        if not actions:
            return {"success": False, "message": "no_actions_to_execute", "fast_path": False}

        if score <= 2:
            # 快路径：直接执行，不保存方法论
            result = self.execute_actions(
                actions, conversation_id, request_id, methodology_manager, conversation_manager
            )
            return {**result, "fast_path": True, "methodology_saved": False, "complexity_score": score}

        elif score <= 4:
            # 标准路径
            result = self.execute_actions(
                actions, conversation_id, request_id, methodology_manager, conversation_manager
            )
            # 强时效性不保存方法论
            save_method = (temporal != "high") and (outcome != "time_bound")
            if save_method and result.get("success_count", 0) > 0:
                # 调用 save_methodology 保存
                method = {
                    "scene": plan.get("summary", "")[:500],
                    "solve_steps": [
                        f"{i + 1}. {a.get('type')}: {a.get('reason', '')[:200]}" for i, a in enumerate(actions)
                    ],
                    "keywords": [],
                    "risk_level": self.evaluate_action_risk_level(actions),
                    "quality_metrics": {
                        "total_steps": int(result.get("total", 0) or 0),
                        "success_steps": int(result.get("success_count", 0) or 0),
                        "failed_steps": int(result.get("failed_count", 0) or 0),
                        "unavailable_count": int(result.get("unavailable_count", 0) or 0),
                    },
                    "score": (
                        float(result.get("success_count", 0) or 0.0) / max(1, float(result.get("total", 0) or 1.0))
                    ),
                }
                try:
                    methodology_manager.add_methodology(method)
                    self.push_log("MethodSaver", "方法论已保存", "completed")
                except Exception as e:
                    self.push_log("MethodSaver", f"方法论保存失败：{e}", "warning")
            return {**result, "fast_path": False, "methodology_saved": save_method, "complexity_score": score}

        else:
            # 重路径：完整流程（已有逻辑保持不变）
            # 这里调用原有的完整执行流程
            return self.execute_full_workflow(
                plan, conversation_id, request_id, methodology_manager, conversation_manager
            )

    def execute_full_workflow(
        self,
        plan: dict[str, Any],
        conversation_id: str,
        request_id: str,
        methodology_manager: Any,
        conversation_manager: Any,
    ) -> dict[str, Any]:
        """完整工作流程：适用于复杂度 5+ 的任务"""
        # 保持原有逻辑不变，这里简单调用 execute_actions
        actions = plan.get("actions") if isinstance(plan.get("actions"), list) else []
        result = self.execute_actions(actions, conversation_id, request_id, methodology_manager, conversation_manager)
        return {
            **result,
            "fast_path": False,
            "full_workflow": True,
            "complexity_score": plan.get("complexity_score", 5),
        }

    def create_execution_session(
        self,
        conversation_id: str,
        request_id: str,
        actions: list[dict[str, Any]],
        methodology_manager: Any,
        conversation_manager: Any,
        *,
        action_screenshots: bool = False,
        plan_summary: str = "",
        plan_risk_level: str = "medium",
    ) -> str:
        session_id = str(uuid.uuid4())
        with self.execution_lock:
            self.execution_sessions[session_id] = {
                "session_id": session_id,
                "session_kind": "batch",
                "conversation_id": conversation_id,
                "request_id": request_id,
                "created_at": time.time(),
                "updated_at": time.time(),
                "status": "pending",
                "paused": False,
                "aborted": False,
                "manual_takeover": False,
                "actions": actions or [],
                "report": [],
                "error": "",
                "action_screenshots": bool(action_screenshots),
                "plan_summary": str(plan_summary or ""),
                "plan_risk_level": str(plan_risk_level or "medium"),
                "_methodology_manager": methodology_manager,
                "_conversation_manager": conversation_manager,
            }
        return session_id

    def create_react_execution_session(
        self,
        conversation_id: str,
        request_id: str,
        user_goal: str,
        dialogue_context: str,
        methodology_manager: Any,
        conversation_manager: Any,
        *,
        action_screenshots: bool = False,
        plan_summary: str = "",
        plan_risk_level: str = "medium",
        react_computer_use_vision: bool = False,
    ) -> str:
        """ReAct：异步线程内逐步 Thought→Action→Observation，与批量 execute_actions 会话并存。"""
        session_id = str(uuid.uuid4())
        with self.execution_lock:
            self.execution_sessions[session_id] = {
                "session_id": session_id,
                "session_kind": "react",
                "conversation_id": conversation_id,
                "request_id": request_id,
                "created_at": time.time(),
                "updated_at": time.time(),
                "status": "pending",
                "paused": False,
                "aborted": False,
                "manual_takeover": False,
                "actions": [],
                "report": [],
                "error": "",
                "action_screenshots": bool(action_screenshots),
                "plan_summary": str(plan_summary or ""),
                "plan_risk_level": str(plan_risk_level or "medium"),
                "react_computer_use_vision": bool(react_computer_use_vision),
                "_methodology_manager": methodology_manager,
                "_conversation_manager": conversation_manager,
                "react_goal": (user_goal or "").strip(),
                "react_dialogue": (dialogue_context or "").strip(),
                "react_trace": [],
                "react_final_message": "",
                "react_iteration_cap": int(self.max_react_iterations),
            }
        return session_id

    def _normalize_react_action_dict(self, raw: dict[str, Any]) -> dict[str, Any] | None:
        if not isinstance(raw, dict):
            return None
        action_type = self._normalize_action_type_alias(str(raw.get("type") or "").strip())
        if not action_type or action_type not in self.ALLOWED_ACTION_TYPES:
            return None
        params = raw.get("params") if isinstance(raw.get("params"), dict) else {}
        risk = self.derive_action_risk(action_type, str(raw.get("risk") or "medium"))
        if action_type == "shell_run":
            cmd = str(params.get("command") or "")
            if any(k in cmd.lower() for k in ["del ", "rm ", "shutdown", "format ", "reg delete"]):
                risk = "high"
        return {
            "type": action_type,
            "target": str(raw.get("target") or "").strip(),
            "filters": raw.get("filters") if isinstance(raw.get("filters"), dict) else {},
            "params": params,
            "risk": risk,
            "reason": str(raw.get("reason") or "").strip(),
        }

    def _react_capability_prompt_fragment(self) -> str:
        return (
            browser_driver.capability_summary_for_planner()
            + desktop_uia.capability_summary_for_planner()
            + screen_ocr.get_capability_summary()
            + self._computer_use_capability_summary()
        )

    def _react_coordinate_contract_prompt_fragment(self) -> str:
        """ReAct/TAOR：与 computer_use.virtual_screen_metrics 同源的坐标说明，避免模型与执行端理解不一致。"""
        from automation import computer_use

        if not computer_use.is_computer_use_enabled():
            return ""
        m = computer_use.virtual_screen_metrics()
        return (
            f"\n【Computer 屏幕与坐标契约】虚拟屏原点=({m['left']},{m['top']})，"
            f"宽={m['width']}，高={m['height']}（多显示器时为合并后的虚拟矩形）。"
            "computer_* 的 x,y 默认 coord_space=absolute（相对该虚拟屏的像素）；"
            "或使用 coord_space=normalized_1000 且 x,y∈[0,1000] 表示在虚拟矩形内的比例位置。\n"
        )

    def _computer_use_capability_summary(self) -> str:
        """预检 computer_use 和 OCR 可用性，返回给 planner 的状态说明。"""
        from automation import computer_use
        lines: list[str] = []
        # computer_use 可用性
        if computer_use.is_computer_use_enabled():
            lines.append("【computer_use 已启用】computer_screenshot/click/type/key 等动作可用。")
        else:
            lines.append("【computer_use 已禁用】所有 computer_* 动作不可用（ARIA_COMPUTER_USE=0）。")
        # OCR (Tesseract) 可用性
        try:
            import pytesseract
            pytesseract.get_tesseract_version()
            lines.append("【Tesseract OCR 可用】screen_ocr/screen_find_text/screen_click_text 可用。")
        except Exception:
            lines.append(
                "【Tesseract OCR 不可用】screen_ocr/screen_find_text/screen_click_text 均会失败，"
                "禁止在计划中使用这三个动作；改用 computer_screenshot + computer_click 视觉方案。"
            )
        return "\n".join(lines) + "\n"


    def _react_format_trace_for_prompt(self, trace: list[dict[str, Any]], max_obs_chars: int = 3500) -> str:
        if not trace:
            return "（尚无先前步骤）"
        lines: list[str] = []
        for row in trace:
            it = int(row.get("iteration") or 0)
            th = str(row.get("thought") or "").strip()
            act = row.get("action")
            obs = str(row.get("observation") or "").strip()
            if len(obs) > max_obs_chars:
                obs = obs[: max_obs_chars - 80] + "\n…[观测已截断]"
            lines.append(f"--- 第 {it} 步 ---\nThought: {th}")
            if isinstance(act, dict) and act.get("type"):
                lines.append(
                    f"Action: type={act.get('type')} target={act.get('target')} reason={act.get('reason', '')[:200]}"
                )
            else:
                lines.append("Action: （无有效动作）")
            lines.append(f"Observation: {obs}\n")
        return "\n".join(lines)

    def _react_observation_from_row(self, row: dict[str, Any]) -> str:
        if not isinstance(row, dict):
            return "（无执行结果）"
        st = str(row.get("status") or "")
        act = str(row.get("action") or "")
        out = str(row.get("stdout") or "").strip()
        err = str(row.get("stderr") or "").strip()
        ec = str(row.get("error_code") or "").strip()
        ver = row.get("verification") if isinstance(row.get("verification"), dict) else {}
        vok = ver.get("ok")
        parts = [f"status={st}", f"action={act}"]
        if ec:
            parts.append(f"error_code={ec}")
        if out:
            parts.append(f"stdout={out[:6000]}")
        if err:
            parts.append(f"stderr={err[:2000]}")
        if vok is not None:
            parts.append(f"verification_ok={vok}")
        res = row.get("result") if isinstance(row.get("result"), dict) else {}
        diag = res.get("computer_diagnostic")
        if isinstance(diag, dict) and diag:
            try:
                parts.append("computer_diagnostic=" + json.dumps(diag, ensure_ascii=False, default=str))
            except Exception:
                parts.append(f"computer_diagnostic={diag!r}")
        return "\n".join(parts)

    def react_infer_next_step(
        self,
        user_goal: str,
        dialogue_context: str,
        trace: list[dict[str, Any]],
        iteration: int,
    ) -> dict[str, Any]:
        """单轮 ReAct：返回 thought、finish、final_message、action（可选）。"""
        current_time_str = time.strftime("%Y年%m月%d日 %H:%M，%A")
        trace_block = self._react_format_trace_for_prompt(trace)
        sys_content = (
            f"【当前时间】{current_time_str}\n\n"
            "你是 ARIA 的 ReAct 执行器：在每一步先推理（Thought），再决定是否调用**一个**自动化动作，或结束任务。\n"
            "你必须只输出一个 JSON 对象，禁止 markdown 代码围栏、禁止前后缀说明文字。\n\n"
            "JSON 字段：\n"
            '- "thought": 字符串，本步推理（下一步要做什么、为什么）。\n'
            '- "finish": 布尔。若为 true，表示任务已在对话内完成或无需再执行动作。\n'
            '- "final_message": 字符串。当 finish=true 时填写给用户的最终说明（可总结已完成步骤）。\n'
            '- "action": 对象或 null。当 finish=false 且需要执行环境动作时，输出**恰好一个**动作；若仅需继续思考下一步则仍应给出 action 或置 finish=true。\n'
            "action 对象字段：type, target, filters(对象), params(对象), risk(low|medium|high), reason。\n\n"
            "规则：\n"
            "- 若上一步 Observation 显示失败，应在 thought 中分析原因并调整策略（换动作、换参数或结束并说明）。\n"
            "- 不要编造未执行的结果。\n"
            "- 单次只输出一个 action；不要输出多个动作。\n"
            "- 能力边界须遵守：\n"
            + self._react_capability_prompt_fragment()
            + self._react_coordinate_contract_prompt_fragment()
            + "\n可用 type 列表："
            + ", ".join(sorted(self.ALLOWED_ACTION_TYPES))
            + "。\n"
        )
        user_block = (
            f"【用户总目标】\n{user_goal.strip()}\n\n"
            f"【会话上下文摘录】\n{(dialogue_context or '').strip() or '（无）'}\n\n"
            f"【已完成 ReAct 轨迹】\n{trace_block}\n\n"
            f"当前迭代序号：{iteration}。请输出下一步 JSON。"
        )
        messages = [
            {"role": "system", "content": sys_content},
            {"role": "user", "content": self._user_content_with_optional_vision(user_block)},
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="ReActAgent")
        data = self._extract_json_object(llm_text)
        if not isinstance(data, dict):
            return {
                "thought": "（解析 JSON 失败）",
                "finish": True,
                "final_message": "ReAct 规划器未返回有效 JSON，请重试或改用普通执行模式。",
                "action": None,
            }
        thought = str(data.get("thought") or "").strip()
        finish = data.get("finish") is True or str(data.get("finish")).lower() in ("1", "true", "yes")
        final_message = str(data.get("final_message") or "").strip()
        raw_action = data.get("action")
        action: dict[str, Any] | None = None
        if isinstance(raw_action, dict) and raw_action.get("type"):
            action = raw_action
        if finish:
            return {"thought": thought, "finish": True, "final_message": final_message or thought, "action": None}
        if not action:
            return {
                "thought": thought or "（未给出动作）",
                "finish": True,
                "final_message": final_message or "模型未提供有效 action 且未标记 finish，已停止 ReAct 循环。",
                "action": None,
            }
        return {"thought": thought, "finish": False, "final_message": "", "action": action}

    def _react_session_desktop_vision_on(self, sess: dict[str, Any]) -> bool:
        """会话级或环境变量 ARIA_REACT_COMPUTER_USE_VISION 任一为真则每步 ReAct 推理前注入桌面截图。"""
        if bool(sess.get("react_computer_use_vision")):
            return True
        v = (os.getenv("ARIA_REACT_COMPUTER_USE_VISION") or "").strip().lower()
        return v in ("1", "true", "yes", "on")

    def _react_cooperative_wait(self, session_id: str, request_id: str) -> bool:
        """若返回 False 表示应结束线程（abort / manual_takeover / cancel）。"""
        while True:
            with self.execution_lock:
                cur = self.execution_sessions.get(session_id) or {}
                if cur.get("aborted"):
                    cur["status"] = "aborted"
                    cur["updated_at"] = time.time()
                    cur["token_usage"] = self.get_token_usage_summary()
                    return False
                if cur.get("manual_takeover"):
                    cur["status"] = "manual_takeover"
                    cur["updated_at"] = time.time()
                    cur["token_usage"] = self.get_token_usage_summary()
                    return False
                is_paused = bool(cur.get("paused"))
            if self.is_cancelled(request_id):
                with self.execution_lock:
                    cur2 = self.execution_sessions.get(session_id) or {}
                    cur2["aborted"] = True
                    cur2["status"] = "aborted"
                    cur2["error"] = "request_cancelled_by_user"
                    cur2["updated_at"] = time.time()
                    cur2["token_usage"] = self.get_token_usage_summary()
                return False
            if not is_paused:
                return True
            time.sleep(0.2)

    def _run_react_execution_session(self, session_id: str) -> None:
        self.reset_token_usage()
        with self.execution_lock:
            sess = self.execution_sessions.get(session_id)
            if not sess:
                return
            methodology_manager = sess.get("_methodology_manager")
            conversation_manager = sess.get("_conversation_manager")
            conversation_id = sess.get("conversation_id") or ""
            request_id = sess.get("request_id") or ""
            shot_flag = bool(sess.get("action_screenshots"))
            plan_summary = str(sess.get("plan_summary") or "")
            plan_risk_level = str(sess.get("plan_risk_level") or "medium")
            goal = str(sess.get("react_goal") or "")
            dialogue = str(sess.get("react_dialogue") or "")
            cap = int(sess.get("react_iteration_cap") or self.max_react_iterations)
            react_vision_on = self._react_session_desktop_vision_on(sess)
        prev_shots = bool(getattr(self, "action_screenshots_for_execution", False))
        self.action_screenshots_for_execution = shot_flag
        report: list[dict[str, Any]] = []
        trace: list[dict[str, Any]] = []
        accumulated_actions: list[dict[str, Any]] = []
        final_msg = ""
        try:
            for it in range(1, cap + 1):
                if not self._react_cooperative_wait(session_id, request_id):
                    return

                self.push_event(
                    "react_iteration",
                    "running",
                    "ReActAgent",
                    f"ReAct 第 {it}/{cap} 轮：推理与决策",
                    {"iteration": it, "cap": cap},
                )
                self.clear_turn_vision_images()
                if react_vision_on:
                    try:
                        from automation import computer_use

                        if computer_use.is_computer_use_enabled():
                            try:
                                jpeg_max = int(os.getenv("ARIA_REACT_COMPUTER_USE_JPEG_MAX", "1280") or "1280")
                            except (TypeError, ValueError):
                                jpeg_max = 1280
                            try:
                                jpeg_q = int(os.getenv("ARIA_REACT_COMPUTER_USE_JPEG_QUALITY", "75") or "75")
                            except (TypeError, ValueError):
                                jpeg_q = 75
                            url = computer_use.capture_jpeg_data_url(
                                max_side=max(320, min(4096, jpeg_max)),
                                quality=max(30, min(95, jpeg_q)),
                            )
                            if url:
                                self.set_turn_vision_images([url])
                    except Exception:
                        pass
                step = self.react_infer_next_step(goal, dialogue, trace, it)
                self.clear_turn_vision_images()
                thought = str(step.get("thought") or "").strip()
                self.record_model_thought("ReActAgent", f"[{it}] {thought[:2000]}")
                self.push_event(
                    "react_thought",
                    "success",
                    "ReActAgent",
                    thought[:500] + ("…" if len(thought) > 500 else ""),
                    {"iteration": it, "thought": thought},
                )

                if step.get("finish"):
                    final_msg = str(step.get("final_message") or "").strip()
                    with self.execution_lock:
                        cur = self.execution_sessions.get(session_id)
                        if cur:
                            cur["react_final_message"] = final_msg
                            cur["react_trace"] = list(trace)
                    break

                raw_action = step.get("action")
                if not isinstance(raw_action, dict):
                    obs = "模型未返回有效 action 对象"
                    trace.append({"iteration": it, "thought": thought, "action": None, "observation": obs})
                    continue

                norm = self._normalize_react_action_dict(raw_action)
                if not norm:
                    obs = f"动作类型不在允许列表或格式无效：{raw_action.get('type')!r}"
                    trace.append({"iteration": it, "thought": thought, "action": raw_action, "observation": obs})
                    self.push_event("react_observation", "warning", "ReActAgent", obs, {"iteration": it})
                    continue

                self._mend_browser_open_actions(goal, {"actions": [norm]})
                accumulated_actions.append(norm)
                self.push_event(
                    "react_action",
                    "running",
                    "ReActAgent",
                    f"执行动作: {norm.get('type')}",
                    {"iteration": it, "action": norm},
                )

                row_list = self.execute_actions(
                    [norm], conversation_id, request_id, methodology_manager, conversation_manager
                ).get("report", [])
                row = row_list[0] if row_list else {}
                if isinstance(row, dict):
                    row["step_id"] = len(report) + 1
                    report.append(row)
                obs = self._react_observation_from_row(row) if isinstance(row, dict) else "（执行无返回）"
                trace.append({"iteration": it, "thought": thought, "action": norm, "observation": obs})

                self.push_event(
                    "react_observation",
                    "success" if str(row.get("status") or "") == "success" else "warning",
                    "ReActAgent",
                    obs[:600] + ("…" if len(obs) > 600 else ""),
                    {"iteration": it, "status": row.get("status"), "full_observation": obs},
                )

                with self.execution_lock:
                    cur = self.execution_sessions.get(session_id)
                    if cur:
                        cur["report"] = list(report)
                        cur["actions"] = list(accumulated_actions)
                        cur["react_trace"] = list(trace)
                        cur["updated_at"] = time.time()

                if not self._react_cooperative_wait(session_id, request_id):
                    return

            with self.execution_lock:
                cur = self.execution_sessions.get(session_id)
                if not cur:
                    return
                cur["react_trace"] = list(trace)
                cur["react_final_message"] = final_msg
                cur["status"] = "completed"
                cur["updated_at"] = time.time()
                cur["token_usage"] = self.get_token_usage_summary()
                cur["quality_metrics"] = self._compute_execution_quality_metrics(report, risk_level=plan_risk_level)
                cur["quality_score"] = self._score_from_quality_metrics(cur["quality_metrics"])
            self.push_event(
                "react_complete",
                "success",
                "ReActAgent",
                f"ReAct 执行完成，共 {len(trace)} 步，动作 {len(report)} 个",
                {"iterations": len(trace), "executed_actions": len(report)},
            )

            if report:
                try:
                    learn = self._auto_learn_from_execution_session(
                        actions=accumulated_actions,
                        report=report,
                        summary=plan_summary or goal[:500],
                        risk_level=plan_risk_level,
                    )
                    self.push_event(
                        "method_auto_iterate",
                        "success",
                        "MethodSaver",
                        "ReAct 执行后已自动更新方法论版本",
                        {"score": learn.get("score"), "metrics": learn.get("metrics")},
                    )
                except Exception as e:
                    self.push_log("MethodSaver", f"自动迭代保存失败：{e}", "warning")

            body = self._format_react_chat_message(trace, report, final_msg)
            tu = self.get_token_usage_summary()
            logs = self.get_execution_log()
            wfe = self.get_workflow_events()
            if conversation_manager and conversation_id:
                conversation_manager.append_message(
                    conversation_id,
                    "assistant",
                    body,
                    {
                        "logs": logs,
                        "workflow_events": wfe,
                        "token_usage": tu,
                        "execution_summary": True,
                        "react_trace": trace,
                        "react_mode": True,
                    },
                )
                conversation_manager.replace_workflow_events(conversation_id, wfe)
        finally:
            self.action_screenshots_for_execution = prev_shots

    def _format_react_chat_message(
        self,
        trace: list[dict[str, Any]],
        report: list[dict[str, Any]],
        final_message: str,
    ) -> str:
        # 技术类型（对用户无意义的中间操作，不在摘要里逐步展示）
        _SKIP_OBS_TYPES = {"computer_screenshot", "computer_move", "screen_ocr", "window_list"}
        parts: list[str] = ["【ReAct 执行摘要】"]
        for row in trace:
            it = int(row.get("iteration") or 0)
            act = row.get("action")
            obs_raw = str(row.get("observation") or "").strip()
            act_type = (act.get("type") or "") if isinstance(act, dict) else ""

            # 跳过纯技术操作的展示（截图/OCR等对用户无意义）
            if act_type in _SKIP_OBS_TYPES:
                continue

            parts.append(f"\n--- 第 {it} 步 ---")
            if isinstance(act, dict) and act.get("type"):
                reason = str(act.get("reason") or "").strip()[:200]
                parts.append(f"操作：{act.get('type')}" + (f"  — {reason}" if reason else ""))
            # 过滤observation中的纯技术调试行
            obs_lines = []
            for line in obs_raw.splitlines():
                # 跳过只含技术键值对的行（如 image_size=... virtual_screen=...）
                if "image_size=" in line or "virtual_screen=" in line or "origin=(" in line:
                    continue
                obs_lines.append(line)
            obs_clean = "\n".join(obs_lines).strip()
            if obs_clean:
                parts.append(f"结果：{obs_clean[:2000]}")
        if (final_message or "").strip():
            parts.append("\n【结束说明】\n" + final_message.strip())
        tail = self._format_execution_report_chat_text(report)
        if tail.strip():
            parts.append("\n\n【环境与工具输出汇总】\n" + tail)
        return "\n".join(parts).strip()

    def format_react_plan_for_user(self, plan: dict[str, Any], user_goal: str) -> str:
        """待确认时展示 ReAct 说明（不逐条列出一次性计划中的全部动作）。"""
        summary = str(plan.get("summary") or "").strip()
        reason = str(plan.get("react_recommend_reason") or "").strip()
        forced = bool(plan.get("react_user_forced"))
        source = str(plan.get("react_mode_source") or "").strip().lower()
        lines: list[str] = []
        if forced:
            lines.append(
                "您已开启「强制 ReAct」：确认后将逐步执行「Thought → Action → Observation」，"
                "每一步根据环境反馈再决定下一步。"
            )
        elif source == "computer":
            lines.append(
                "本计划包含桌面坐标类自动化（Computer Use），已自动采用 ReAct：确认后将按「Thought → Action → Observation」"
                "逐步执行，以便根据屏幕反馈调整下一步。"
            )
        else:
            lines.append(
                "ARIA 判断本任务适合使用 ReAct（逐步推理执行）：确认后将按「Thought → Action → Observation」循环推进，"
                "而非一次性跑完固定动作列表。"
            )
        if reason:
            lines.append(f"\n判断说明：{reason[:1500]}")
        if summary:
            lines.append(f"\n规划摘要：{summary}")
        lines.append(f"\n目标原文：{user_goal.strip()[:2000]}")
        lines.append("\n请回复「确认执行」以开始 ReAct 循环（仍适用风险确认策略：高风险需二次确认）。")
        return "\n".join(lines)

    def start_execution_session(self, session_id: str) -> dict[str, Any]:
        with self.execution_lock:
            sess = self.execution_sessions.get(session_id)
            if not sess:
                return {"success": False, "message": "session_not_found"}
            if sess["status"] in ("running", "completed"):
                return {"success": True, "status": sess["status"], "session_id": session_id}
            sess["status"] = "running"
            sess["updated_at"] = time.time()
            kind = str(sess.get("session_kind") or "batch")
        target = self._run_react_execution_session if kind == "react" else self._run_execution_session
        t = threading.Thread(target=target, args=(session_id,), daemon=True)
        t.start()
        return {"success": True, "status": "running", "session_id": session_id}

    def _format_execution_report_chat_text(self, report: list[dict[str, Any]], max_total: int = 32000) -> str:
        # 纯技术操作，对用户无展示价值
        _SKIP_REPORT_TYPES = {"computer_screenshot", "computer_move", "screen_ocr", "window_list", "window_activate"}
        parts: list[str] = []
        for r in report:
            if not isinstance(r, dict):
                continue
            act = str(r.get("action") or "")
            st = str(r.get("status") or "")
            inp = r.get("input") if isinstance(r.get("input"), dict) else {}
            out = str(r.get("stdout") or "").strip()
            err = str(r.get("stderr") or "").strip()
            if act in _SKIP_REPORT_TYPES:
                # 失败时仍提示用户，但不展示纯技术截图参数
                if st != "success":
                    parts.append(f"【{act}】失败：{err or '操作未成功'}")
                continue
            if act in ("web_understand", "web_fetch"):
                if st == "success" and out:
                    parts.append(f"【{act}】\n{out}")
                else:
                    detail = err or out or "未知错误"
                    parts.append(f"【{act}】失败：{detail}")
            elif act == "browser_open":
                url = self._browser_open_raw_url(inp) if isinstance(inp, dict) else ""
                if st == "success":
                    parts.append(f"【browser_open】已尝试在本机打开：{url or '（未解析到 URL）'}")
                else:
                    parts.append(f"【browser_open】失败：{err or '未知错误'}")
            elif act == "file_write":
                rel = ""
                if isinstance(inp, dict):
                    pp = inp.get("params") if isinstance(inp.get("params"), dict) else {}
                    rel = str(pp.get("path") or "").strip().replace("\\", "/")
                if st == "success":
                    lines_fw = [
                        "【file_write】已在 ARIA 工作区内写入文件（非「我的文档」等系统路径，除非您显式写了绝对路径）。"
                    ]
                    if out:
                        lines_fw.append(out)
                    if rel and ".." not in rel and not rel.startswith("/"):
                        lines_fw.append(
                            "【下载】在同一浏览器打开本 ARIA 站点，新标签访问（或复制到地址栏）：\n"
                            f"/api/workspace_file?path={quote(rel, safe='')}"
                        )
                    parts.append("\n".join(lines_fw))
                else:
                    parts.append(f"【file_write】失败：{err or out or '未知错误'}")
            else:
                if st != "success":
                    parts.append(f"【{act}】失败：{err or out or '未知错误'}")
                elif out:
                    parts.append(f"【{act}】\n{out}")
        joined = "\n\n".join(parts).strip()
        if not joined:
            return "动作已执行完毕。未生成可展示的步骤输出，请在协作日志或执行报告中查看详情。"
        if len(joined) > max_total:
            return joined[: max_total - 120] + f"\n\n... [已截断，原长约 {len(joined)} 字符]"
        return joined

    def _append_execution_report_to_conversation(
        self,
        conversation_id: str,
        conversation_manager: Any,
        report: list[dict[str, Any]],
    ) -> None:
        if not conversation_id or not conversation_manager:
            return
        body = self._format_execution_report_chat_text(report)
        tu = self.get_token_usage_summary()
        logs = self.get_execution_log()
        wfe = self.get_workflow_events()
        conversation_manager.append_message(
            conversation_id,
            "assistant",
            body,
            {
                "logs": logs,
                "workflow_events": wfe,
                "token_usage": tu,
                "execution_summary": True,
            },
        )
        conversation_manager.replace_workflow_events(conversation_id, wfe)

    def _compute_execution_quality_metrics(
        self, report: list[dict[str, Any]], risk_level: str = "medium"
    ) -> dict[str, Any]:
        total = len(report or [])
        success = sum(1 for r in (report or []) if str(r.get("status")) == "success")
        unavailable = sum(
            1
            for r in (report or [])
            if str((r.get("result") or {}).get("error_code") or r.get("error_code") or "") == "unavailable_capability"
        )
        retryable = sum(1 for r in (report or []) if bool(r.get("retryable")))
        confirmations = (
            2 if str(risk_level or "medium") == "high" else (1 if str(risk_level or "medium") == "medium" else 0)
        )
        duration_ms = sum(int(r.get("duration_ms", 0) or 0) for r in (report or []))
        success_rate = (success / total) if total > 0 else 0.0
        return {
            "total_steps": total,
            "success_steps": success,
            "failed_steps": max(0, total - success),
            "success_rate": round(success_rate, 4),
            "unavailable_count": unavailable,
            "retryable_failures": retryable,
            "confirmations": confirmations,
            "duration_ms": duration_ms,
        }

    def _score_from_quality_metrics(self, metrics: dict[str, Any]) -> float:
        success_rate = float(metrics.get("success_rate", 0.0) or 0.0)
        unavailable_penalty = float(metrics.get("unavailable_count", 0) or 0) * 0.1
        retry_penalty = float(metrics.get("retryable_failures", 0) or 0) * 0.05
        score = max(0.0, min(1.0, success_rate - unavailable_penalty - retry_penalty))
        return round(score, 4)

    def _refine_method_steps_from_report(
        self, actions: list[dict[str, Any]], report: list[dict[str, Any]]
    ) -> list[str]:
        steps: list[str] = []
        for i, a in enumerate(actions or [], start=1):
            t = str(a.get("type") or "")
            reason = str(a.get("reason") or "").strip()
            steps.append(f"{i}. {t} - {reason or '执行该步骤'}")
        failed = [r for r in (report or []) if str(r.get("status") or "") != "success"]
        if failed:
            steps.append(f"{len(steps) + 1}. 若失败，先读取 error_code / stderr，优先重试 retryable 步骤。")
            steps.append(f"{len(steps) + 1}. 若能力不可用（unavailable_capability），切换可用驱动或请求用户手动接管。")
        return steps

    def _auto_learn_from_execution_session(
        self,
        *,
        actions: list[dict[str, Any]],
        report: list[dict[str, Any]],
        summary: str,
        risk_level: str,
    ) -> dict[str, Any]:
        metrics = self._compute_execution_quality_metrics(report)
        score = self._score_from_quality_metrics(metrics)
        steps = self._refine_method_steps_from_report(actions, report)
        scene = (summary or "执行计划").strip()[:500]
        keywords = sorted(
            {str(a.get("type") or "").strip() for a in (actions or []) if str(a.get("type") or "").strip()}
        )[:10]
        method = {
            "scene": scene,
            "scenario": scene,
            "keywords": keywords,
            "core_keywords": keywords,
            "solve_steps": steps,
            "outcome_type": "stable",
            "risk_level": (risk_level or "medium"),
            "score": score,
            "quality_metrics": metrics,
            "evidence_refs": [f"action_report_steps:{len(report or [])}"],
            "is_success": metrics.get("failed_steps", 0) == 0,
        }
        saved = self.ltm.add_methodology(method)
        return {"saved": saved, "metrics": metrics, "score": score}

    def _run_execution_session(self, session_id: str) -> None:
        self.reset_token_usage()
        with self.execution_lock:
            sess = self.execution_sessions.get(session_id)
            if not sess:
                return
            actions = list(sess.get("actions") or [])[: self.max_action_steps]
            methodology_manager = sess.get("_methodology_manager")
            conversation_manager = sess.get("_conversation_manager")
            conversation_id = sess.get("conversation_id") or ""
            request_id = sess.get("request_id") or ""
            shot_flag = bool(sess.get("action_screenshots"))
            plan_summary = str(sess.get("plan_summary") or "")
            plan_risk_level = str(sess.get("plan_risk_level") or "medium")
        prev_shots = bool(getattr(self, "action_screenshots_for_execution", False))
        self.action_screenshots_for_execution = shot_flag
        report: list[dict[str, Any]] = []
        try:
            for idx, action in enumerate(actions, start=1):
                # cooperative pause/abort checks
                while True:
                    with self.execution_lock:
                        cur = self.execution_sessions.get(session_id) or {}
                        if cur.get("aborted"):
                            cur["status"] = "aborted"
                            cur["updated_at"] = time.time()
                            cur["token_usage"] = self.get_token_usage_summary()
                            return
                        if cur.get("manual_takeover"):
                            cur["status"] = "manual_takeover"
                            cur["updated_at"] = time.time()
                            cur["token_usage"] = self.get_token_usage_summary()
                            return
                        is_paused = bool(cur.get("paused"))
                    if self.is_cancelled(request_id):
                        with self.execution_lock:
                            cur2 = self.execution_sessions.get(session_id) or {}
                            cur2["aborted"] = True
                            cur2["status"] = "aborted"
                            cur2["error"] = "request_cancelled_by_user"
                            cur2["updated_at"] = time.time()
                            cur2["token_usage"] = self.get_token_usage_summary()
                        return
                    if not is_paused:
                        break
                    time.sleep(0.2)

                row = self.execute_actions(
                    [action], conversation_id, request_id, methodology_manager, conversation_manager
                ).get("report", [])
                if row:
                    row[0]["step_id"] = idx
                    report.append(row[0])
                with self.execution_lock:
                    cur = self.execution_sessions.get(session_id)
                    if not cur:
                        return
                    cur["report"] = list(report)
                    cur["updated_at"] = time.time()

            with self.execution_lock:
                cur = self.execution_sessions.get(session_id)
                if not cur:
                    return
                cur["status"] = "completed"
                cur["updated_at"] = time.time()
                cur["token_usage"] = self.get_token_usage_summary()
                cur["quality_metrics"] = self._compute_execution_quality_metrics(report, risk_level=plan_risk_level)
                cur["quality_score"] = self._score_from_quality_metrics(cur["quality_metrics"])

            if report:
                try:
                    learn = self._auto_learn_from_execution_session(
                        actions=actions,
                        report=report,
                        summary=plan_summary,
                        risk_level=plan_risk_level,
                    )
                    self.push_event(
                        "method_auto_iterate",
                        "success",
                        "MethodSaver",
                        "执行后已自动更新方法论版本",
                        {"score": learn.get("score"), "metrics": learn.get("metrics")},
                    )
                except Exception as e:
                    self.push_log("MethodSaver", f"自动迭代保存失败：{e}", "warning")
            self._append_execution_report_to_conversation(conversation_id, conversation_manager, report)
        finally:
            self.action_screenshots_for_execution = prev_shots

    def pause_execution_session(self, session_id: str) -> dict[str, Any]:
        with self.execution_lock:
            sess = self.execution_sessions.get(session_id)
            if not sess:
                return {"success": False, "message": "session_not_found"}
            sess["paused"] = True
            sess["status"] = "paused"
            sess["updated_at"] = time.time()
            return {"success": True, "status": "paused"}

    def resume_execution_session(self, session_id: str) -> dict[str, Any]:
        with self.execution_lock:
            sess = self.execution_sessions.get(session_id)
            if not sess:
                return {"success": False, "message": "session_not_found"}
            sess["paused"] = False
            if sess.get("status") == "paused":
                sess["status"] = "running"
            sess["updated_at"] = time.time()
            return {"success": True, "status": sess.get("status")}

    def abort_execution_session(self, session_id: str) -> dict[str, Any]:
        with self.execution_lock:
            sess = self.execution_sessions.get(session_id)
            if not sess:
                return {"success": False, "message": "session_not_found"}
            sess["aborted"] = True
            sess["status"] = "aborted"
            sess["updated_at"] = time.time()
            return {"success": True, "status": "aborted"}

    def get_execution_session(self, session_id: str) -> dict[str, Any]:
        with self.execution_lock:
            sess = self.execution_sessions.get(session_id)
            if not sess:
                return {"success": False, "message": "session_not_found"}
            report = sess.get("report") or []
            fallback_used_steps = sum(1 for r in report if bool(r.get("fallback_used")))
            safe_block_steps = sum(1 for r in report if str(r.get("safe_block_reason") or "").strip())
            avg_conf = 0.0
            if report:
                vals: list[float] = []
                for r in report:
                    try:
                        vals.append(float(r.get("confidence", 0.0) or 0.0))
                    except Exception:
                        continue
                if vals:
                    avg_conf = sum(vals) / len(vals)
            out: dict[str, Any] = {
                "success": True,
                "session_id": session_id,
                "session_kind": str(sess.get("session_kind") or "batch"),
                "conversation_id": sess.get("conversation_id"),
                "request_id": sess.get("request_id"),
                "status": sess.get("status"),
                "paused": sess.get("paused"),
                "aborted": sess.get("aborted"),
                "manual_takeover": sess.get("manual_takeover"),
                "report": report,
                "created_at": sess.get("created_at"),
                "updated_at": sess.get("updated_at"),
                "token_usage": sess.get("token_usage"),
                "quality_metrics": sess.get("quality_metrics") or {},
                "quality_score": sess.get("quality_score"),
                "intelligence_metrics": {
                    "fallback_used_steps": fallback_used_steps,
                    "safe_block_steps": safe_block_steps,
                    "avg_confidence": round(avg_conf, 4),
                },
            }
            if str(sess.get("session_kind") or "") == "react":
                out["react_trace"] = sess.get("react_trace") or []
                out["react_final_message"] = str(sess.get("react_final_message") or "")
            ca = float(sess.get("created_at") or 0)
            ua = float(sess.get("updated_at") or 0)
            st = str(sess.get("status") or "")
            terminal = st in ("completed", "aborted", "manual_takeover")
            if terminal and ca > 0 and ua >= ca:
                elapsed_wall = int((ua - ca) * 1000)
            elif ca > 0:
                elapsed_wall = int(max(0, (time.time() - ca) * 1000))
            else:
                elapsed_wall = 0
            tu_sess = sess.get("token_usage") if isinstance(sess.get("token_usage"), dict) else {}
            llm_ms = int((tu_sess or {}).get("llm_wall_ms", 0) or 0)
            qm = sess.get("quality_metrics") if isinstance(sess.get("quality_metrics"), dict) else {}
            action_ms = int((qm or {}).get("duration_ms", 0) or 0)
            if action_ms <= 0 and report:
                action_ms = sum(int((r or {}).get("duration_ms", 0) or 0) for r in report)
            out["elapsed_ms"] = elapsed_wall
            out["timing_breakdown"] = compute_timing_breakdown(
                elapsed_ms=elapsed_wall,
                llm_ms=llm_ms,
                local_action_ms=action_ms,
            )
            return out

    def _exec_kb_delete_all(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        all_methods = methodology_manager.get_all_methodologies() or []
        ids = [m.get("method_id") for m in all_methods if m.get("method_id")]
        return methodology_manager.delete_methodologies_batch(ids)

    def _exec_kb_delete_by_keyword(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        kw = str((action.get("filters") or {}).get("keyword") or "").strip()
        candidates = methodology_manager.search_methodologies(kw) if kw else []
        ids = [m.get("method_id") for m in candidates if m.get("method_id")]
        result = methodology_manager.delete_methodologies_batch(ids)
        result["keyword"] = kw
        return result

    def _exec_kb_delete_low_quality(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        all_methods = methodology_manager.get_all_methodologies() or []
        ids = [
            m.get("method_id")
            for m in all_methods
            if m.get("method_id") and int(m.get("success_count", 0)) <= 0 and int(m.get("usage_count", 0)) <= 0
        ]
        return methodology_manager.delete_methodologies_batch(ids)

    def _exec_conversation_new(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        title = str((action.get("params") or {}).get("title") or "新会话")
        conv = conversation_manager.create_conversation(title)
        return {"success": True, "conversation_id": conv.get("conversation_id")}

    def _exec_shell_run(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        command = self._sanitize_shell_command(str(params.get("command") or ""))
        cwd = str(params.get("cwd") or ".")
        timeout_s = int(params.get("timeout_s") or self.default_step_timeout_s)
        safe_cwd = self._ensure_safe_path(cwd)
        try:
            proc = subprocess.run(
                command,
                cwd=str(safe_cwd),
                shell=True,
                capture_output=True,
                text=True,
                timeout=max(1, timeout_s),
            )
            raw_stdout = proc.stdout or ""
            truncated = len(raw_stdout) > 4000
            return {
                "success": proc.returncode == 0,
                "returncode": proc.returncode,
                "stdout": raw_stdout[:4000] + (" ...[truncated]" if truncated else ""),
                "stderr": (proc.stderr or "")[:4000],
                "artifacts": [],
                "screenshots": [],
            }
        except subprocess.TimeoutExpired:
            return {
                "success": False,
                "returncode": -1,
                "stdout": "",
                "stderr": f"Command timed out after {timeout_s}s",
                "error_code": "timeout",
                "retryable": True,
                "artifacts": [],
                "screenshots": [],
            }

    @staticmethod
    def _office_scalar_cell(v: Any) -> Any:
        if v is None:
            return ""
        if isinstance(v, (str, int, float, bool)):
            return v
        return str(v)

    def _patch_office_docx_styles(self, path: Path, style: dict[str, Any]) -> None:
        """修改已有 .docx 的正文与表格单元格内 run 的字体样式（尽力覆盖常见结构）。"""
        from docx import Document
        from docx.shared import Pt

        doc = Document(str(path))
        font_name = str(style.get("font_name") or "").strip()
        size_pt = style.get("font_size_pt")
        bold = style.get("bold")
        italic = style.get("italic")

        def apply_run(run: Any) -> None:
            if font_name:
                run.font.name = font_name
            if size_pt is not None:
                try:
                    run.font.size = Pt(int(size_pt))
                except (TypeError, ValueError):
                    pass
            if bold is not None:
                run.bold = bool(bold)
            if italic is not None:
                run.italic = bool(italic)

        for para in doc.paragraphs:
            for run in para.runs:
                apply_run(run)
        for tbl in doc.tables:
            for row in tbl.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            apply_run(run)
        doc.save(str(path))

    def _write_office_docx(self, path: Path, params: dict[str, Any]) -> None:
        from docx import Document

        doc = Document()
        title = str(params.get("title") or "").strip()
        if title:
            doc.add_heading(title, 0)
        content = str(params.get("content") or "")
        if len(content) > _OFFICE_MAX_CONTENT_CHARS:
            content = content[:_OFFICE_MAX_CONTENT_CHARS]
        for line in content.splitlines():
            doc.add_paragraph(line)
        doc.save(str(path))

    def _write_office_xlsx(self, path: Path, params: dict[str, Any]) -> None:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        rows_raw = params.get("rows")
        if isinstance(rows_raw, str):
            try:
                rows_raw = json.loads(rows_raw)
            except Exception:
                rows_raw = None
        if isinstance(rows_raw, list) and rows_raw:
            for r_idx, row in enumerate(rows_raw[:_OFFICE_MAX_ROWS], 1):
                if isinstance(row, (list, tuple)):
                    cells = list(row)[:_OFFICE_MAX_COLS]
                else:
                    cells = [row]
                for c_idx, cell in enumerate(cells, 1):
                    ws.cell(row=r_idx, column=c_idx, value=self._office_scalar_cell(cell))
        else:
            content = str(params.get("content") or "")
            if len(content) > _OFFICE_MAX_CONTENT_CHARS:
                content = content[:_OFFICE_MAX_CONTENT_CHARS]
            for r_idx, line in enumerate(content.splitlines()[:_OFFICE_MAX_ROWS], 1):
                if "\t" in line:
                    parts = line.split("\t")
                else:
                    parts = re.split(r",", line)
                parts = [p.strip() for p in parts][:_OFFICE_MAX_COLS]
                for c_idx, val in enumerate(parts, 1):
                    ws.cell(row=r_idx, column=c_idx, value=val)
        wb.save(str(path))

    def _write_office_pptx(self, path: Path, params: dict[str, Any]) -> None:
        from pptx import Presentation

        prs = Presentation()
        content = str(params.get("content") or "")
        if len(content) > _OFFICE_MAX_CONTENT_CHARS:
            content = content[:_OFFICE_MAX_CONTENT_CHARS]
        lines = [ln.strip() for ln in content.splitlines() if ln.strip()]
        title = str(params.get("title") or "").strip()
        bullets_raw = params.get("bullets")
        if isinstance(bullets_raw, str):
            try:
                bullets_raw = json.loads(bullets_raw)
            except Exception:
                bullets_raw = None
        if isinstance(bullets_raw, list):
            b_lines = [str(b).strip() for b in bullets_raw if str(b).strip()][:_OFFICE_MAX_PPT_BULLETS]
        else:
            b_lines = []
        if not title and lines:
            title = lines[0]
            lines = lines[1:]
        if not title:
            title = "演示"
        if not b_lines:
            b_lines = lines[:_OFFICE_MAX_PPT_BULLETS]

        slide0 = prs.slides.add_slide(prs.slide_layouts[0])
        slide0.shapes.title.text = title

        if b_lines:
            slide1 = prs.slides.add_slide(prs.slide_layouts[1])
            slide1.shapes.title.text = "要点"
            body = slide1.shapes.placeholders[1].text_frame
            body.text = b_lines[0]
            for bl in b_lines[1:]:
                p = body.add_paragraph()
                p.text = bl
                p.level = 0
        prs.save(str(path))

    def _exec_file_write(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        path = self._ensure_safe_path(str(params.get("path") or ""))
        mode = str(params.get("mode") or "overwrite").lower()
        ext = path.suffix.lower()
        path.parent.mkdir(parents=True, exist_ok=True)

        if ext == ".docx":
            p_dict = params if isinstance(params, dict) else {}
            docx_style = p_dict.get("docx_style") if isinstance(p_dict.get("docx_style"), dict) else {}
            has_style = bool(docx_style) and any(v not in (None, "", []) for v in docx_style.values())
            try:
                if path.is_file() and has_style:
                    self._patch_office_docx_styles(path, docx_style)
                    return {
                        "success": True,
                        "message": f"docx_style_patched:{path}",
                        "artifacts": [str(path)],
                        "screenshots": [],
                    }
                if path.is_file() and isinstance(p_dict.get("docx_style"), dict) and not has_style:
                    return {
                        "success": False,
                        "message": "docx_style_empty",
                        "stderr": "文件已存在但未给出有效的 docx_style（如 font_name、font_size_pt、bold），避免误覆盖全文。",
                        "artifacts": [],
                        "screenshots": [],
                    }
                self._write_office_docx(path, p_dict)
            except ImportError:
                return {
                    "success": False,
                    "message": "missing_dependency",
                    "stderr": "未安装 python-docx，请执行：pip install python-docx",
                    "stdout": "",
                    "artifacts": [],
                    "screenshots": [],
                }
            except Exception as e:
                return {
                    "success": False,
                    "message": "docx_write_failed",
                    "stderr": str(e),
                    "stdout": "",
                    "artifacts": [],
                    "screenshots": [],
                }
            return {"success": True, "message": f"file_written:{path}", "artifacts": [str(path)], "screenshots": []}

        if ext in (".xlsx", ".xlsm"):
            try:
                self._write_office_xlsx(path, params if isinstance(params, dict) else {})
            except ImportError:
                return {
                    "success": False,
                    "message": "missing_dependency",
                    "stderr": "未安装 openpyxl，请执行：pip install openpyxl",
                    "stdout": "",
                    "artifacts": [],
                    "screenshots": [],
                }
            except Exception as e:
                return {
                    "success": False,
                    "message": "xlsx_write_failed",
                    "stderr": str(e),
                    "stdout": "",
                    "artifacts": [],
                    "screenshots": [],
                }
            return {"success": True, "message": f"file_written:{path}", "artifacts": [str(path)], "screenshots": []}

        if ext in (".xls", ".doc"):
            return {
                "success": False,
                "message": "legacy_office_format",
                "stderr": "服务端仅支持生成 .xlsx / .docx / .pptx，请将路径改为对应扩展名。",
                "stdout": "",
                "artifacts": [],
                "screenshots": [],
            }

        if ext == ".pptx":
            try:
                self._write_office_pptx(path, params if isinstance(params, dict) else {})
            except ImportError:
                return {
                    "success": False,
                    "message": "missing_dependency",
                    "stderr": "未安装 python-pptx，请执行：pip install python-pptx",
                    "stdout": "",
                    "artifacts": [],
                    "screenshots": [],
                }
            except Exception as e:
                return {
                    "success": False,
                    "message": "pptx_write_failed",
                    "stderr": str(e),
                    "stdout": "",
                    "artifacts": [],
                    "screenshots": [],
                }
            return {"success": True, "message": f"file_written:{path}", "artifacts": [str(path)], "screenshots": []}

        try:
            content = str(params.get("content") or "")
            with open(path, "a" if mode == "append" else "w", encoding="utf-8") as f:
                f.write(content)
            return {"success": True, "message": f"file_written:{path}", "artifacts": [str(path)], "screenshots": []}
        except PermissionError as e:
            return {"success": False, "message": f"Permission denied: {e}", "error_code": "permission_denied", "retryable": False, "artifacts": [], "screenshots": []}
        except OSError as e:
            return {"success": False, "message": f"File write failed: {e}", "error_code": "io_error", "retryable": False, "artifacts": [], "screenshots": []}

    def _exec_file_move(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        src = self._ensure_safe_path(str(params.get("src") or ""))
        dst = self._ensure_safe_path(str(params.get("dst") or ""))
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(src), str(dst))
        return {"success": True, "message": f"moved:{src}->{dst}", "artifacts": [str(dst)], "screenshots": []}

    def _exec_file_delete(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        target = self._ensure_safe_path(str(params.get("path") or ""))
        if target.is_dir():
            shutil.rmtree(str(target))
        elif target.exists():
            target.unlink()
        return {"success": True, "message": f"deleted:{target}", "artifacts": [], "screenshots": []}

    def _exec_file_read(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        path = self._ensure_safe_path(str(params.get("path") or ""))
        encoding = str(params.get("encoding") or "utf-8")
        max_bytes = int(params.get("max_bytes") or 51200)  # 默认 50 KB
        try:
            size = path.stat().st_size
            truncated = size > max_bytes
            with open(path, "r", encoding=encoding, errors="replace") as f:
                content = f.read(max_bytes)
            return {
                "success": True,
                "message": f"file_read:{path}",
                "content": content,
                "size_bytes": size,
                "truncated": truncated,
                "artifacts": [],
                "screenshots": [],
            }
        except FileNotFoundError:
            return {"success": False, "message": f"file_not_found:{path}", "error_code": "file_not_found", "artifacts": [], "screenshots": []}
        except (PermissionError, OSError) as e:
            return {"success": False, "message": f"file_read_failed:{e}", "error_code": "io_error", "artifacts": [], "screenshots": []}

    def _exec_file_append(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        path = self._ensure_safe_path(str(params.get("path") or ""))
        content = str(params.get("content") or "")
        path.parent.mkdir(parents=True, exist_ok=True)
        try:
            with open(path, "a", encoding="utf-8") as f:
                f.write(content)
            return {"success": True, "message": f"file_appended:{path}", "artifacts": [str(path)], "screenshots": []}
        except (PermissionError, OSError) as e:
            return {"success": False, "message": f"file_append_failed:{e}", "error_code": "io_error", "artifacts": [], "screenshots": []}

    def _exec_file_create_dir(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        path = self._ensure_safe_path(str(params.get("path") or ""))
        try:
            path.mkdir(parents=True, exist_ok=True)
            return {"success": True, "message": f"dir_created:{path}", "artifacts": [str(path)], "screenshots": []}
        except (PermissionError, OSError) as e:
            return {"success": False, "message": f"file_create_dir_failed:{e}", "error_code": "io_error", "artifacts": [], "screenshots": []}

    def _exec_file_list(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        path = self._ensure_safe_path(str(params.get("path") or "."))
        recursive = bool(params.get("recursive", False))
        pattern = str(params.get("pattern") or "*")
        try:
            if recursive:
                matches = list(path.rglob(pattern))
            else:
                matches = list(path.glob(pattern))
            entries = []
            for p in sorted(matches):
                try:
                    stat = p.stat()
                    entries.append({
                        "name": p.name,
                        "path": str(p),
                        "type": "dir" if p.is_dir() else "file",
                        "size_bytes": stat.st_size if p.is_file() else None,
                    })
                except OSError:
                    pass
            return {"success": True, "message": f"file_list:{path}", "entries": entries, "count": len(entries), "artifacts": [], "screenshots": []}
        except (PermissionError, OSError) as e:
            return {"success": False, "message": f"file_list_failed:{e}", "error_code": "io_error", "artifacts": [], "screenshots": []}

    def _exec_file_find(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        directory = self._ensure_safe_path(str(params.get("directory") or "."))
        name_pattern = str(params.get("name_pattern") or "*")
        content_contains = str(params.get("content_contains") or "").strip()
        max_results = int(params.get("max_results") or 100)
        try:
            matches = []
            for p in sorted(directory.rglob(name_pattern)):
                if len(matches) >= max_results:
                    break
                if p.is_dir():
                    continue
                if content_contains:
                    try:
                        text = p.read_text(encoding="utf-8", errors="replace")
                        if content_contains not in text:
                            continue
                    except OSError:
                        continue
                try:
                    matches.append({"name": p.name, "path": str(p), "size_bytes": p.stat().st_size})
                except OSError:
                    pass
            return {"success": True, "message": f"file_find:{directory}", "matches": matches, "count": len(matches), "artifacts": [], "screenshots": []}
        except (PermissionError, OSError) as e:
            return {"success": False, "message": f"file_find_failed:{e}", "error_code": "io_error", "artifacts": [], "screenshots": []}

    def _exec_wechat_check_login(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation.apps.wechat_automation import wechat_check_login
        params = action.get("params") or {}
        result = wechat_check_login(is_enterprise=bool(params.get("is_enterprise", False)))
        result.setdefault("artifacts", [])
        result.setdefault("screenshots", [])
        return result

    def _exec_wechat_open_chat(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation.apps.wechat_automation import wechat_open_chat
        params = action.get("params") or {}
        contact = str(params.get("contact_name") or action.get("target") or "").strip()
        result = wechat_open_chat(contact_name=contact, is_enterprise=bool(params.get("is_enterprise", False)))
        result.setdefault("artifacts", [])
        result.setdefault("screenshots", [])
        return result

    def _exec_wechat_send_message(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation.apps.wechat_automation import wechat_send_message
        params = action.get("params") or {}
        contact = str(params.get("contact_name") or action.get("target") or "").strip()
        message = str(params.get("message") or "").strip()
        result = wechat_send_message(contact_name=contact, message=message, is_enterprise=bool(params.get("is_enterprise", False)))
        result.setdefault("artifacts", [])
        result.setdefault("screenshots", [])
        return result

    def _exec_screen_watch_start(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation.screen_watcher import get_or_create_watcher, ScreenChangeEvent
        params = action.get("params") or {}
        interval = float(params.get("interval", 2.0))
        threshold = float(params.get("diff_threshold", 0.02))
        save_dir = str(params.get("save_dir") or "").strip() or None

        def _on_change(event: ScreenChangeEvent) -> None:
            self.push_event(
                "screen_watch", "info", "ScreenWatcher",
                f"屏幕变化检测：{event.diff_ratio:.1%}",
                {"diff_ratio": event.diff_ratio, "screenshot_path": event.screenshot_path,
                 "timestamp": event.timestamp},
            )

        watcher = get_or_create_watcher(
            _on_change, interval=interval, diff_threshold=threshold, save_dir=save_dir
        )
        return {
            "success": True,
            "message": f"ScreenWatcher 已启动（interval={interval}s, threshold={threshold:.1%}）",
            "artifacts": [], "screenshots": [],
        }

    def _exec_screen_watch_stop(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation.screen_watcher import stop_global_watcher
        stop_global_watcher()
        return {"success": True, "message": "ScreenWatcher 已停止", "artifacts": [], "screenshots": []}

    def _exec_email_send(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation.email_client import email_send
        params = action.get("params") or {}
        result = email_send(params)
        result.setdefault("artifacts", [])
        result.setdefault("screenshots", [])
        return result

    def _exec_email_read(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation.email_client import email_read
        params = action.get("params") or {}
        result = email_read(params)
        result.setdefault("artifacts", [])
        result.setdefault("screenshots", [])
        return result

    def _exec_clipboard_read(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        try:
            import pyperclip  # type: ignore
            content = pyperclip.paste()
            return {"success": True, "message": "clipboard_read", "content": content, "content_type": "text", "artifacts": [], "screenshots": []}
        except ImportError:
            return {"success": False, "message": "missing_dependency", "stderr": "未安装 pyperclip，请执行：pip install pyperclip", "artifacts": [], "screenshots": []}
        except Exception as e:
            return {"success": False, "message": f"clipboard_read_failed:{e}", "error_code": "io_error", "artifacts": [], "screenshots": []}

    def _exec_clipboard_write(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        text = str(params.get("text") or "")
        try:
            import pyperclip  # type: ignore
            pyperclip.copy(text)
            return {"success": True, "message": "clipboard_written", "artifacts": [], "screenshots": []}
        except ImportError:
            return {"success": False, "message": "missing_dependency", "stderr": "未安装 pyperclip，请执行：pip install pyperclip", "artifacts": [], "screenshots": []}
        except Exception as e:
            return {"success": False, "message": f"clipboard_write_failed:{e}", "error_code": "io_error", "artifacts": [], "screenshots": []}

    def _resolve_fetch_url(self, action: dict[str, Any]) -> str:
        params = action.get("params") if isinstance(action.get("params"), dict) else {}
        u = str(params.get("url") or "").strip()
        if not u:
            u = str(action.get("target") or "").strip()
        if not u:
            raise ValueError("missing_url")
        if not u.startswith(("http://", "https://")):
            u = "https://" + u.lstrip("/")
        return u

    def _assert_url_host_is_public(self, url: str) -> None:
        p = urlparse(url)
        if p.scheme not in ("http", "https"):
            raise ValueError("fetch_only_http_https")
        host = p.hostname
        if not host:
            raise ValueError("missing_host")
        try:
            for res in socket.getaddrinfo(host, None):
                ip = ipaddress.ip_address(res[4][0])
                if not ip.is_global:
                    raise ValueError("blocked_network")
        except ValueError as e:
            if str(e) in ("blocked_network", "fetch_only_http_https", "missing_host"):
                raise
            raise ValueError("dns_failed") from e

    def _http_get_bytes_capped(self, url: str, max_bytes: int = 2_000_000) -> tuple[str, bytes, str, Optional[str]]:
        self._assert_url_host_is_public(url)
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
            ),
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
        }
        with requests.get(
            url,
            timeout=25,
            allow_redirects=True,
            headers=headers,
            stream=True,
        ) as r:
            r.raise_for_status()
            final_url = r.url
            self._assert_url_host_is_public(final_url)
            ctype = (r.headers.get("content-type") or "").split(";")[0].strip().lower()
            enc = requests.utils.get_encoding_from_headers(r.headers) or r.encoding
            chunks: list[bytes] = []
            total = 0
            for chunk in r.iter_content(chunk_size=65536):
                if self.is_cancelled():
                    raise TaskCancelledError()
                if not chunk:
                    continue
                remain = max_bytes - total
                if remain <= 0:
                    break
                if len(chunk) <= remain:
                    chunks.append(chunk)
                    total += len(chunk)
                else:
                    chunks.append(chunk[:remain])
                    break
        return final_url, b"".join(chunks), ctype, enc

    def _html_to_plain_text(self, html: str, max_chars: int) -> str:
        parser = _HTMLToTextParser()
        try:
            parser.feed(html[: min(len(html), 3_000_000)])
            parser.close()
        except Exception:
            pass
        text = parser.text()
        if len(text) > max_chars:
            return text[:max_chars] + f"\n... [正文已截断，共约 {len(text)} 字]"
        return text

    def _exec_web_fetch(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") if isinstance(action.get("params"), dict) else {}
        url = self._resolve_fetch_url(action)
        max_chars = int(params.get("max_chars") or 32000)
        max_chars = max(1000, min(max_chars, 120_000))
        final_url, raw, ctype, enc = self._http_get_bytes_capped(url)
        enc = enc or "utf-8"
        try:
            body = raw.decode(enc, errors="replace")
        except LookupError:
            body = raw.decode("utf-8", errors="replace")
        head = body[:4000].lower()
        looks_html = "html" in ctype or body.lstrip().lower().startswith("<!doctype") or "<html" in head
        if not looks_html:
            snippet = body[:max_chars]
            out = self._cap_fetch_stdout(snippet)
            return {
                "success": True,
                "message": f"fetched_non_html:{final_url}",
                "stdout": out,
                "artifacts": [],
                "screenshots": [],
            }
        text = self._html_to_plain_text(body, max_chars)
        out = self._cap_fetch_stdout(text)
        return {
            "success": True,
            "message": f"fetched:{final_url} chars={len(text)}",
            "stdout": out,
            "artifacts": [],
            "screenshots": [],
        }

    def _cap_fetch_stdout(self, s: str, lim: int = 96_000) -> str:
        if len(s) <= lim:
            return s
        return s[:lim] + f"\n... [stdout 已截断，总长 {len(s)}]"

    def _exec_web_understand(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") if isinstance(action.get("params"), dict) else {}
        url = self._resolve_fetch_url(action)
        question = str(params.get("question") or params.get("prompt") or "").strip()
        if not question:
            question = "请用简洁的要点总结该网页的主要内容；若正文不足以判断，请说明。"
        max_chars = int(params.get("max_chars") or 28000)
        max_chars = max(4000, min(max_chars, 80_000))

        final_url = url
        text = ""
        last_err = ""
        last_snip = ""
        fetch_attempts = 0
        for try_url in self._understand_fetch_url_candidates(url, question):
            fetch_attempts += 1
            try:
                fu, raw, _ctype, enc = self._http_get_bytes_capped(try_url)
            except Exception as e:
                last_err = str(e)
                continue
            enc = enc or "utf-8"
            try:
                html = raw.decode(enc, errors="replace")
            except (LookupError, UnicodeDecodeError):
                html = raw.decode("utf-8", errors="replace")
            candidate = self._html_to_plain_text(html, max_chars)
            final_url = fu
            text = candidate
            if len(candidate.strip()) >= 80:
                break
            last_snip = candidate[:2000]
            last_err = "抓取到的正文过短（常见于需登录或强 JS 渲染的页面）"

        if len(text.strip()) < 80:
            detail = last_err or "抓取失败"
            if last_snip:
                detail += f"；末次摘录预览：{last_snip[:400]}"
            if fetch_attempts > 1:
                detail += "（已按序尝试多个搜索源，含 DuckDuckGo HTML 与 Bing）"
            return {
                "success": False,
                "message": "page_too_empty_or_not_html",
                "stderr": detail,
                "stdout": (text or last_snip)[:2000],
                "artifacts": [],
                "screenshots": [],
            }
        messages = [
            {
                "role": "system",
                "content": (
                    "你是网页阅读助手。只根据「网页摘录」作答，不要编造页面上没有的信息。"
                    "若摘录不完整或无法回答，如实说明。用中文、条理清晰、尽量简洁。"
                    "禁止使用任何 Markdown 格式（禁止 #、##、**、*、`、--- 等符号）；直接用纯文字和数字编号输出。" + _MATH_NOTATION_FOR_CHAT
                ),
            },
            {
                "role": "user",
                "content": f"页面最终 URL：{final_url}\n\n网页摘录：\n{text}\n\n用户问题：{question}",
            },
        ]
        answer = self._call_llm(
            messages,
            fallback_text="（未配置 API Key 或模型调用失败，无法生成理解摘要。可先执行 web_fetch 查看原始正文。）",
            agent_code="TextExecAgent",
        )
        answer = (answer or "").strip()
        out = self._cap_fetch_stdout(answer, lim=48_000)
        return {
            "success": True,
            "message": "web_understood",
            "stdout": out,
            "stderr": "",
            "artifacts": [],
            "screenshots": [],
        }

    def _exec_browser_open(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        url = str(params.get("url") or "").strip()
        if not url:
            url = str(action.get("target") or "").strip()
        if not url:
            raise ValueError("missing_url")
        if not url.startswith(("http://", "https://")):
            url = "https://" + url.lstrip("/")
        shots = self._capture_screenshot("browser_open")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.navigate(url)
            if ok:
                return {
                    "success": True,
                    "message": f"browser_playwright_opened:{url}",
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_open_failed",
                "stderr": err or "playwright_navigate_failed",
                "artifacts": [],
                "screenshots": shots,
            }
        if browser_driver.is_playwright_enabled() and not browser_driver.playwright_package_installed():
            webbrowser.open(url)
            return {
                "success": True,
                "message": f"browser_opened:{url} (ARIA_PLAYWRIGHT=1 but playwright package missing; used default browser)",
                "stderr": "install: pip install playwright && playwright install chromium",
                "artifacts": [],
                "screenshots": shots,
            }
        webbrowser.open(url)
        return {"success": True, "message": f"browser_opened:{url}", "artifacts": [], "screenshots": shots}

    def _exec_browser_click(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip()
        nav = str(params.get("url") or "").strip() or None
        shots = self._capture_screenshot("browser_click")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.click(selector, navigate_url=nav)
            if ok:
                return {
                    "success": True,
                    "message": f"browser_click:{selector}",
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_click_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {
            "success": True,
            "message": f"browser_click_simulated:{selector}",
            "artifacts": [],
            "screenshots": shots,
        }

    def _exec_browser_type(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip()
        text = str(params.get("text") or "")
        nav = str(params.get("url") or "").strip() or None
        shots = self._capture_screenshot("browser_type")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.fill(selector, text, navigate_url=nav)
            if ok:
                return {
                    "success": True,
                    "message": f"browser_type:{selector}",
                    "stdout": text[:200],
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_type_failed",
                "stdout": text[:200],
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {
            "success": True,
            "message": f"browser_type_simulated:{selector}",
            "stdout": text[:200],
            "artifacts": [],
            "screenshots": shots,
        }

    def _exec_browser_find(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip()
        text_contains = str(params.get("text_contains") or "").strip() or None
        shots = self._capture_screenshot("browser_find")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, results, err = browser_driver.find_elements(selector, text_contains)
            if ok:
                return {
                    "success": True,
                    "message": f"browser_find:{selector}",
                    "stdout": f"Found {len(results)} element(s): " + json.dumps(results[:5], ensure_ascii=False),
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_find_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {"success": True, "message": f"browser_find_simulated:{selector}", "artifacts": [], "screenshots": shots}

    def _exec_browser_hover(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip()
        shots = self._capture_screenshot("browser_hover")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.hover(selector)
            if ok:
                return {"success": True, "message": f"browser_hover:{selector}", "artifacts": [], "screenshots": shots}
            return {
                "success": False,
                "message": "browser_hover_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {
            "success": True,
            "message": f"browser_hover_simulated:{selector}",
            "artifacts": [],
            "screenshots": shots,
        }

    def _exec_browser_select(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip()
        value = str(params.get("value") or "").strip()
        shots = self._capture_screenshot("browser_select")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.select_option(selector, value)
            if ok:
                return {
                    "success": True,
                    "message": f"browser_select:{selector}={value}",
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_select_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {
            "success": True,
            "message": f"browser_select_simulated:{selector}",
            "artifacts": [],
            "screenshots": shots,
        }

    def _exec_browser_upload(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip()
        file_path = str(params.get("file_path") or "").strip()
        shots = self._capture_screenshot("browser_upload")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.upload_file(selector, file_path)
            if ok:
                return {
                    "success": True,
                    "message": f"browser_upload:{selector} -> {file_path}",
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_upload_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {
            "success": True,
            "message": f"browser_upload_simulated:{selector}",
            "artifacts": [],
            "screenshots": shots,
        }

    def _exec_browser_scroll(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip() or None
        shots = self._capture_screenshot("browser_scroll")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.scroll_to(selector)
            if ok:
                return {
                    "success": True,
                    "message": f"browser_scroll:{selector or 'bottom'}",
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_scroll_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {
            "success": True,
            "message": f"browser_scroll_simulated:{selector or 'bottom'}",
            "artifacts": [],
            "screenshots": shots,
        }

    def _exec_browser_wait(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        selector = str(params.get("selector") or "").strip()
        timeout_ms = int(params.get("timeout_ms") or 30000)
        shots = self._capture_screenshot("browser_wait")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.wait_for_element(selector, timeout_ms)
            if ok:
                return {"success": True, "message": f"browser_wait:{selector}", "artifacts": [], "screenshots": shots}
            return {
                "success": False,
                "message": "browser_wait_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {"success": True, "message": f"browser_wait_simulated:{selector}", "artifacts": [], "screenshots": shots}

    def _exec_browser_js(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        script = str(params.get("script") or "").strip()
        shots = self._capture_screenshot("browser_js")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, result, err = browser_driver.execute_javascript(script)
            if ok:
                return {
                    "success": True,
                    "message": "browser_js executed",
                    "stdout": json.dumps(result, ensure_ascii=False)[:500] if result else "",
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "browser_js_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {"success": True, "message": "browser_js_simulated", "artifacts": [], "screenshots": shots}

    def _path_to_image_data_url(self, path: Path) -> str | None:
        ext = path.suffix.lower().lstrip(".")
        mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg", "webp": "image/webp"}.get(ext)
        if not mime:
            return None
        try:
            data = path.read_bytes()
        except OSError:
            return None
        if len(data) > 4 * 1024 * 1024:
            return None
        b64 = base64.standard_b64encode(data).decode("ascii")
        return f"data:{mime};base64,{b64}"

    def _youtube_oembed_blurb(self, url: str) -> str | None:
        try:
            r = requests.get(
                "https://www.youtube.com/oembed",
                params={"url": url, "format": "json"},
                timeout=12,
            )
            if not r.ok:
                return None
            j = r.json()
            title = str(j.get("title") or "").strip()
            author = str(j.get("author_name") or "").strip()
            if not title:
                return None
            return f"视频标题：{title}\n上传者：{author or '（未知）'}\n（来自 YouTube oEmbed，非逐字稿；完整内容需观看视频或提供字幕。）"
        except Exception:
            return None

    def _exec_browser_press(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        key = str(params.get("key") or "").strip()
        selector = str(params.get("selector") or "").strip() or None
        shots = self._capture_screenshot("browser_press")
        if browser_driver.is_playwright_enabled() and browser_driver.playwright_package_installed():
            ok, err = browser_driver.press_key(key, selector=selector)
            if ok:
                return {"success": True, "message": f"browser_press:{key}", "artifacts": [], "screenshots": shots}
            return {
                "success": False,
                "message": "browser_press_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {"success": True, "message": f"browser_press_simulated:{key}", "artifacts": [], "screenshots": shots}

    def _exec_media_summarize(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") if isinstance(action.get("params"), dict) else {}
        url = str(params.get("url") or action.get("target") or "").strip()
        relpath = str(params.get("path") or "").strip()
        question = str(params.get("question") or "请用中文简要归纳主题与可视要点。").strip()
        if relpath:
            try:
                video_path = self._ensure_safe_path(relpath)
            except ValueError as e:
                return {
                    "success": False,
                    "message": "invalid_path",
                    "stderr": str(e),
                    "artifacts": [],
                    "screenshots": [],
                }
            if not video_path.is_file():
                return {
                    "success": False,
                    "message": "video_not_found",
                    "stderr": str(video_path),
                    "artifacts": [],
                    "screenshots": [],
                }
            suf = video_path.suffix.lower()
            if suf not in (".mp4", ".webm", ".mov", ".mkv"):
                return {
                    "success": False,
                    "message": "unsupported_video_ext",
                    "stderr": "仅支持工作区内 mp4/webm/mov/mkv（需本机 ffmpeg 抽帧）",
                    "artifacts": [],
                    "screenshots": [],
                }
            frames_dir = self.allowed_work_root / "data" / "artifacts" / "video_frames"
            frames_dir.mkdir(parents=True, exist_ok=True)
            prefix = frames_dir / f"vf_{int(time.time() * 1000)}"
            out_pattern = str(prefix) + "_%03d.png"
            try:
                proc = subprocess.run(
                    [
                        "ffmpeg",
                        "-y",
                        "-i",
                        str(video_path),
                        "-vf",
                        "fps=0.2",
                        "-frames:v",
                        "6",
                        out_pattern,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                )
            except FileNotFoundError:
                return {
                    "success": False,
                    "message": "ffmpeg_not_found",
                    "stderr": "未找到 ffmpeg，请安装后重试（并将 ffmpeg 加入 PATH）",
                    "artifacts": [],
                    "screenshots": [],
                }
            except subprocess.TimeoutExpired:
                return {
                    "success": False,
                    "message": "ffmpeg_timeout",
                    "stderr": "ffmpeg 超时",
                    "artifacts": [],
                    "screenshots": [],
                }
            if proc.returncode != 0:
                return {
                    "success": False,
                    "message": "ffmpeg_failed",
                    "stderr": (proc.stderr or proc.stdout or "")[:2000],
                    "artifacts": [],
                    "screenshots": [],
                }
            frame_paths = sorted(frames_dir.glob(f"{prefix.name}_*.png"))[:6]
            parts: list[dict[str, Any]] = [{"type": "text", "text": question}]
            for fp in frame_paths:
                du = self._path_to_image_data_url(fp)
                if du:
                    parts.append({"type": "image_url", "image_url": {"url": du}})
            if len(parts) <= 1:
                return {
                    "success": False,
                    "message": "no_frames",
                    "stderr": "未能生成抽帧图片",
                    "artifacts": [],
                    "screenshots": [],
                }
            messages = [{"role": "user", "content": parts}]
            summary = self._call_llm(messages, fallback_text="（模型不可用或未返回摘要）", agent_code="MediaSummarizer")
            return {
                "success": True,
                "message": "media_summarize_video",
                "stdout": summary[:8000],
                "artifacts": [str(p) for p in frame_paths],
                "screenshots": [],
            }
        ulow = url.lower()
        if "youtube.com" in ulow or "youtu.be" in ulow:
            blurb = self._youtube_oembed_blurb(url)
            if blurb:
                follow = self._call_llm(
                    [
                        {
                            "role": "user",
                            "content": f"根据下列视频元数据回答用户问题（可补充说明无法从元数据得知的细节）：\n{blurb}\n\n用户问题：{question}",
                        }
                    ],
                    fallback_text=blurb,
                    agent_code="MediaSummarizer",
                )
                return {
                    "success": True,
                    "message": "media_summarize_oembed",
                    "stdout": follow[:8000],
                    "artifacts": [],
                    "screenshots": [],
                }
        return {
            "success": False,
            "message": "media_summarize_need_path",
            "stderr": "请提供 params.path（工作区视频文件）或可 oEmbed 的 YouTube 链接。",
            "artifacts": [],
            "screenshots": [],
        }

    def _exec_desktop_open_app(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        app = str(params.get("app") or "").strip()
        if not app:
            app = str(action.get("target") or "").strip()
        if not app:
            raise ValueError("missing_app")
        launched = app
        web_alternative = None
        # 检查应用是否已在运行（Windows），避免重复打开
        if os.name == "nt":
            try:
                import psutil
                # 用关键词扩展表匹配（"微信" → ["wechat","WeChat","weixin",...] 等多语言别名）
                kws = [k.lower().replace(".exe", "") for k in _windows_open_app_keywords(app)]
                if not kws:
                    kws = [app.lower().replace(".exe", "")]
                for proc in psutil.process_iter(["name"]):
                    try:
                        pname = (proc.info.get("name") or "").lower().replace(".exe", "")
                    except (psutil.NoSuchProcess, psutil.AccessDenied):
                        continue
                    if not pname:
                        continue
                    if any(kw and (kw in pname or pname in kw) for kw in kws):
                        # 已在运行：直接返回，不截图（避免额外延迟）
                        return {
                            "success": True,
                            "message": f"app_already_running:{app}",
                            "stdout": f"app_already_running:{app}",
                            "artifacts": [],
                            "screenshots": [],
                            "strategy_path": "rule_path",
                            "confidence": 0.95,
                            "fallback_used": False,
                        }
            except Exception:
                pass  # psutil 不可用，继续正常启动
        try:
            if os.name == "nt":
                resolved, info = _windows_resolve_app_executable(app)
                web_alternative = info.get("web_alternative")
                if resolved is not None:
                    os.startfile(str(resolved))
                    launched = f"{app} -> {resolved}"
                else:
                    # 未在桌面/开始菜单/常见路径解析到文件时仍尝试 shell 解析（注册的应用别名等）
                    subprocess.Popen(["cmd", "/c", "start", "", app], shell=False)
            else:
                subprocess.Popen(app, shell=True)
        except OSError as e:
            # 如果打开失败且有网页版替代方案，返回特殊响应
            if web_alternative:
                return {
                    "success": False,
                    "message": f"open_failed:{e}",
                    "web_alternative": web_alternative,
                    "suggestion": f"未找到 {app}，已为您准备网页版：{web_alternative}",
                    "strategy_path": "rule_path",
                    "confidence": 0.35,
                    "safe_block_reason": "unresolved_target",
                }
            raise ValueError(
                f"open_failed:{e}（已尝试匹配桌面/开始菜单快捷方式与常见安装路径；仍失败请改用 .exe 完整路径）"
            ) from e
        # 启动成功：稍等应用渲染后截图（避免截到空白启动画面）
        wait_s = float((action.get("params") or {}).get("wait_s") or 1.2)
        time.sleep(max(0.5, min(10.0, wait_s)))
        shots = self._capture_screenshot("desktop_open_app")
        return {
            "success": True,
            "message": f"desktop_app_opened:{launched}",
            "stdout": f"desktop_app_opened:{launched}",
            "artifacts": [],
            "screenshots": shots,
            "strategy_path": "rule_path",
            "confidence": 0.9,
            "fallback_used": False,
        }

    def _exec_desktop_hotkey(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        hotkey = str(params.get("hotkey") or "").strip()
        shots = self._capture_screenshot("desktop_hotkey")
        if desktop_uia.is_uia_enabled() and os.name == "nt" and desktop_uia.pywinauto_package_installed():
            ok, err = desktop_uia.send_hotkey(hotkey)
            if ok:
                return {
                    "success": True,
                    "message": f"desktop_hotkey:{hotkey}",
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "desktop_hotkey_failed",
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {"success": True, "message": f"desktop_hotkey_simulated:{hotkey}", "artifacts": [], "screenshots": shots}

    def _exec_desktop_type(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        params = action.get("params") or {}
        text = str(params.get("text") or params.get("content") or "")
        shots = self._capture_screenshot("desktop_type")
        if desktop_uia.is_uia_enabled() and os.name == "nt" and desktop_uia.pywinauto_package_installed():
            ok, err = desktop_uia.type_text(text)
            if ok:
                return {
                    "success": True,
                    "message": "desktop_type_sent",
                    "stdout": text[:400],
                    "artifacts": [],
                    "screenshots": shots,
                }
            return {
                "success": False,
                "message": "desktop_type_failed",
                "stdout": text[:400],
                "stderr": err or "unknown",
                "artifacts": [],
                "screenshots": shots,
            }
        return {
            "success": True,
            "message": "desktop_type_simulated",
            "stdout": text[:400],
            "artifacts": [],
            "screenshots": shots,
        }

    def _exec_desktop_sequence(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        """顺序执行多步：sleep / hotkey / type，用于任意桌面应用脚本（params.steps）。"""
        params = action.get("params") or {}
        steps = params.get("steps")
        if not isinstance(steps, list) or not steps:
            return {
                "success": False,
                "message": "desktop_sequence_empty",
                "stderr": "missing_or_empty_steps",
                "artifacts": [],
                "screenshots": [],
                "strategy_path": "rule_path",
                "confidence": 0.2,
                "safe_block_reason": "unsafe_to_continue",
            }
        shots = self._capture_screenshot("desktop_sequence")
        real = desktop_uia.is_uia_enabled() and os.name == "nt" and desktop_uia.pywinauto_package_installed()
        log_parts: list[str] = []
        for step in steps:
            if not isinstance(step, dict):
                continue
            stype = str(step.get("type") or "").strip().lower()
            if stype == "sleep":
                time.sleep(float(step.get("seconds", step.get("sec", 0)) or 0))
                log_parts.append(f"sleep:{step.get('seconds', step.get('sec', 0))}")
            elif stype in ("hotkey", "desktop_hotkey"):
                hk = str(step.get("hotkey") or step.get("keys") or "").strip()
                if not hk:
                    return {
                        "success": False,
                        "message": "desktop_sequence_failed",
                        "stderr": "empty_hotkey_step",
                        "artifacts": [],
                        "screenshots": shots,
                        "strategy_path": "rule_path",
                        "confidence": 0.25,
                        "safe_block_reason": "unsafe_to_continue",
                    }
                if real:
                    ok, err = desktop_uia.send_hotkey(hk)
                    if not ok:
                        return {
                            "success": False,
                            "message": "desktop_sequence_failed",
                            "stderr": err or "hotkey_failed",
                            "artifacts": [],
                            "screenshots": shots,
                            "strategy_path": "rule_path",
                            "confidence": 0.25,
                            "safe_block_reason": "unsafe_to_continue",
                        }
                log_parts.append(f"hotkey:{hk}")
            elif stype in ("type", "desktop_type", "text"):
                txt = str(step.get("text") or step.get("content") or "")
                if real:
                    ok, err = desktop_uia.type_text(txt)
                    if not ok:
                        return {
                            "success": False,
                            "message": "desktop_sequence_failed",
                            "stderr": err or "type_failed",
                            "artifacts": [],
                            "screenshots": shots,
                            "strategy_path": "rule_path",
                            "confidence": 0.25,
                            "safe_block_reason": "unsafe_to_continue",
                        }
                log_parts.append(f"type:{len(txt)}chars")
            else:
                return {
                    "success": False,
                    "message": "desktop_sequence_unknown_step",
                    "stderr": f"unknown_step_type:{stype}",
                    "artifacts": [],
                    "screenshots": shots,
                    "strategy_path": "rule_path",
                    "confidence": 0.2,
                    "safe_block_reason": "unsafe_to_continue",
                }
        out = ";".join(log_parts)
        if real:
            return {
                "success": True,
                "message": "desktop_sequence_ok",
                "stdout": out[:800],
                "artifacts": [],
                "screenshots": shots,
                "strategy_path": "rule_path",
                "confidence": 0.88,
                "fallback_used": False,
            }
        return {
            "success": True,
            "message": f"desktop_sequence_simulated:{out[:400]}",
            "stdout": out[:800],
            "artifacts": [],
            "screenshots": shots,
            "strategy_path": "rule_path",
            "confidence": 0.6,
            "fallback_used": False,
        }

    def _exec_screen_ocr(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        """
        执行屏幕 OCR 识别

        params:
            region: 可选，(left, top, width, height) 或 None（全屏）
            lang: OCR 语言，默认 'chi_sim+eng'
        """
        params = action.get("params") or {}
        region = params.get("region")
        lang = params.get("lang", "chi_sim+eng")

        shots = self._capture_screenshot("screen_ocr")

        result = screen_ocr.ocr_screen(region, lang)

        if result.get("success"):
            text_preview = result.get("text", "")[:500]
            blocks_count = len(result.get("blocks", []))
            return {
                "success": True,
                "message": f"screen_ocr_completed:识别到{blocks_count}个文字块",
                "stdout": f"识别结果（前 500 字）：\n{text_preview}",
                "artifacts": [],
                "screenshots": shots,
            }
        else:
            error = str(result.get("error") or "unknown")
            hint = str(result.get("hint") or "")
            msg = error + (f"。{hint}" if hint else "")
            return {
                "success": False,
                "message": msg,
                "stderr": msg,
                "artifacts": [],
                "screenshots": shots,
            }

    def _exec_screen_find_text(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        """
        执行屏幕文字查找

        params:
            text:           要查找的文字（必填）
            region:         可选，(left, top, width, height) 或 None（全屏）
            lang:           OCR 语言，默认 'chi_sim+eng'
            scale:          截图放大倍数，默认 2.0（有利于 UI 小字号识别）
            min_confidence: Tesseract 最低置信度，默认 30
        """
        params = action.get("params") or {}
        search_text = params.get("text", "")
        region = params.get("region")
        lang = params.get("lang", "chi_sim+eng")
        try:
            scale = float(params.get("scale") or 2.0)
        except (TypeError, ValueError):
            scale = 2.0
        try:
            min_confidence = int(params.get("min_confidence") or 30)
        except (TypeError, ValueError):
            min_confidence = 30

        shots = self._capture_screenshot("screen_find_text")

        if not search_text.strip():
            return {
                "success": False,
                "message": "missing_search_text:请指定要查找的文字",
                "stderr": "missing_search_text",
                "artifacts": [],
                "screenshots": shots,
            }

        result = screen_ocr.find_text_on_screen(search_text, region, lang, scale=scale, min_confidence=min_confidence)

        if result.get("success"):
            matches = result.get("matches", [])
            if matches:
                positions = [f"{m['text']} @ ({m['center'][0]}, {m['center'][1]})" for m in matches[:5]]
                return {
                    "success": True,
                    "message": f"screen_find_text_found:找到{len(matches)}处匹配",
                    "stdout": "找到位置：\n" + "\n".join(positions),
                    "artifacts": [],
                    "screenshots": shots,
                }
            else:
                return {
                    "success": False,
                    "message": f"text_not_found:未在屏幕上找到「{search_text}」",
                    "stderr": f"text_not_found:{search_text}",
                    "artifacts": [],
                    "screenshots": shots,
                }
        else:
            error = result.get("error", "unknown")
            return {
                "success": False,
                "message": error,
                "stderr": error,
                "artifacts": [],
                "screenshots": shots,
            }

    def _exec_screen_click_text(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        """
        执行点击屏幕文字

        params:
            text: 要点击的文字（必填）
            region: 可选，(left, top, width, height) 或 None（全屏）
            lang: OCR 语言，默认 'chi_sim+eng'
            button: 鼠标按钮，'left' / 'right' / 'middle'，默认 'left'
        """
        params = action.get("params") or {}
        search_text = params.get("text", "")
        region = params.get("region")
        lang = params.get("lang", "chi_sim+eng")
        button = params.get("button", "left")

        shots = self._capture_screenshot("screen_click_text")

        if not search_text.strip():
            return {
                "success": False,
                "message": "missing_search_text:请指定要点击的文字",
                "stderr": "missing_search_text",
                "artifacts": [],
                "screenshots": shots,
            }

        result = screen_ocr.click_text(search_text, region, lang, button)

        if result.get("success"):
            position = result.get("position")
            pos_str = f"位置 ({position[0]}, {position[1]})" if position else ""
            return {
                "success": True,
                "message": result.get("message", "screen_click_text_completed"),
                "stdout": f"已点击「{search_text}」{pos_str}",
                "artifacts": [],
                "screenshots": shots,
            }
        else:
            error = result.get("error", "unknown")
            return {
                "success": False,
                "message": error,
                "stderr": error,
                "artifacts": [],
                "screenshots": shots,
            }

    # ------------------------------------------------------------------ #
    # computer_use 动作执行器                                               #
    # ------------------------------------------------------------------ #

    def _exec_computer_screenshot(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        result = computer_use.run_screenshot_info(params)
        # 本步仅保留一张图，避免 ReAct / 多轮累积撑爆上下文
        try:
            try:
                jpeg_max = int(os.getenv("ARIA_REACT_COMPUTER_USE_JPEG_MAX", "1280") or "1280")
            except (TypeError, ValueError):
                jpeg_max = 1280
            try:
                jpeg_q = int(os.getenv("ARIA_REACT_COMPUTER_USE_JPEG_QUALITY", "75") or "75")
            except (TypeError, ValueError):
                jpeg_q = 75
            data_url = computer_use.capture_jpeg_data_url(
                max_side=max(320, min(4096, jpeg_max)),
                quality=max(30, min(95, jpeg_q)),
            )
            if data_url:
                self._turn_vision_data_urls = [data_url]
                result["screenshot_data_url"] = data_url
        except Exception:
            pass
        return result

    def _exec_computer_click(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        button = str(params.get("button") or "left")
        return computer_use.run_click(params, button=button, clicks=1)

    def _exec_computer_click_element(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        """
        用 OS-Atlas grounding 模型按自然语言描述定位并点击 UI 元素。

        params:
            query       : 元素描述，如"搜索框"、"发送按钮"（必填）
            button      : left / right / middle（默认 left）
            screenshot  : 可选，base64 data URL；不提供时自动截图
        """
        from automation import computer_use, osatlas_grounding
        params = action.get("params") or {}
        query = str(params.get("query") or "").strip()
        if not query:
            return {"success": False, "message": "missing_query", "stderr": "computer_click_element requires params.query"}

        # 获取截图
        data_url = str(params.get("screenshot") or "").strip()
        if not data_url:
            shot = computer_use.capture_jpeg_data_url()
            data_url = shot

        # OS-Atlas 定位
        pos = osatlas_grounding.find_element_from_data_url(query, data_url)
        if pos is None:
            return {
                "success": False,
                "message": "osatlas_grounding_failed",
                "stderr": f"OS-Atlas could not locate element: {query!r}",
            }

        x, y = pos
        button = str(params.get("button") or "left")
        click_params = {"x": x, "y": y, "coord_space": "absolute"}
        result = computer_use.run_click(click_params, button=button, clicks=1)
        result["grounding_query"] = query
        result["grounded_at"] = f"({x},{y})"
        return result

    def _exec_computer_double_click(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        button = str(params.get("button") or "left")
        return computer_use.run_click(params, button=button, clicks=2)

    def _exec_computer_move(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        return computer_use.run_move(params)

    def _exec_computer_drag(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        return computer_use.run_drag(params)

    def _exec_computer_scroll(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        return computer_use.run_scroll(params)

    def _exec_computer_key(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        return computer_use.run_key(params)

    def _exec_computer_type(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        return computer_use.run_type_text(params)

    def _exec_computer_wait(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        from automation import computer_use
        params = action.get("params") or {}
        return computer_use.run_wait(params)

    def _exec_window_activate(
        self, action: dict[str, Any], conversation_id: str, methodology_manager: Any, conversation_manager: Any
    ) -> dict[str, Any]:
        """
        按标题子串将指定窗口激活到前台（Windows 专用）。

        params:
            title: 窗口标题关键词（子串匹配，必填）。例如 "微信" 可匹配 "微信" 窗口。
        """
        from automation import computer_use
        params = action.get("params") or {}
        return computer_use.run_window_activate(params)

    def generate_small_talk_reply(self, user_input: str) -> str:
        text = (user_input or "").strip()
        user_line = text or ("请根据图片简单打个招呼或说明你能看到的内容。" if self._turn_vision_data_urls else "你好")
        messages = [
            {
                "role": "system",
                "content": (
                    "你是ARIA助手。请对用户寒暄做简短友好回复。"
                    "要求：1-2句话、总长度不超过50字、不输出JSON、不输出步骤。"
                ),
            },
            {"role": "user", "content": self._user_content_with_optional_vision(user_line)},
        ]
        llm_text = self._call_llm_fast(messages, fallback_text="")
        cleaned = (llm_text or "").strip()
        if cleaned:
            cleaned = cleaned.replace("```", "").strip()
            if len(cleaned) <= 120 and "{" not in cleaned:
                return cleaned
        if any(k in text.lower() for k in ["谢谢", "thank"]):
            return "不客气，我在这儿，随时可以继续帮你。"
        return "你好，我在。告诉我你想解决什么问题，我马上开始。"

    def generate_direct_qa_reply(self, user_input: str, dialogue_context: str = "") -> str:
        """规划器已判定 task_form=qa_only：单轮对话作答（可走 LLM 流式）。"""
        text = (user_input or "").strip()
        ctx = (dialogue_context or "").strip()
        user_blob = (
            f"【本会话近期对话】\n{ctx}\n\n【本轮问题】\n{text}"
            if ctx
            else f"【本轮问题】\n{text or '（见多模态附件）'}"
        )
        messages = [
            {
                "role": "system",
                "content": (
                    "你是 ARIA 助手。根据用户问题直接作答：条理清晰、用语准确；需要时分点或短段落。"
                    "不要输出 JSON、不要编造未提供的链接或执行结果；未给出的信息如实说明不确定。"
                    + _MATH_NOTATION_FOR_CHAT
                ),
            },
            {"role": "user", "content": self._user_content_with_optional_vision(user_blob)},
        ]
        return (
            self._call_llm_fast(
                messages,
                fallback_text="（模型暂不可用，请稍后重试。）",
            )
            or ""
        ).strip()

    # 1. 解析用户问题
    def parse_task(self, user_input: str, dialogue_context: str = "", reuse_task_id: str | None = None) -> dict:
        self.check_cancelled("task_parse_start")
        self.current_task_id = self._resolve_task_id_for_turn(reuse_task_id)
        self.push_event("task_parse", "running", "TaskParser", "PM 正在解析用户需求")
        self.push_log("TaskParser", "正在分析你的问题", "running")
        # 记录模型思考过程
        self.record_model_thought("TaskParser", f"收到用户输入: {user_input}")
        self.record_model_thought("TaskParser", "分析用户意图和任务类型")

        # LLM 任务解析（失败则回退到简单规则解析）
        fallback_task_info = {
            "task_id": self.current_task_id,
            "user_input": user_input,
            "task_type": "text",
            "intent": "general",
            "keywords": user_input.split()[:5],
            "temporal_risk": self._infer_temporal_risk(user_input),
            "execution_surface": "conversation",
            "timestamp": time.time(),
        }
        messages = [
            {
                "role": "system",
                "content": (
                    "你是ARIA任务解析器。请根据用户输入提取任务类型(task_type)、意图(intent)、关键词(keywords)、时效风险(temporal_risk)以及执行面(execution_surface)。只输出严格JSON，不要多余文本。"
                    "keywords为数组，元素为短关键词（3-8个字/词）。"
                    "temporal_risk 只能为 high 或 low：high 表示结论强依赖当下数据（天气、股价、赛果、实时路况等），旧答案不可照抄；low 表示以稳定知识或流程为主。"
                    "execution_surface 只能为下列之一："
                    "local_desktop（用户要操作本机程序/自动化：微信/企微发消息、打开桌面软件、终端命令、本地文件读写、Playwright/桌面 UIA 等，不依赖「上网找教程」完成）；"
                    "web_research（用户明确要上网查资料、读新闻、抓网页/链接摘要、搜教程等）；"
                    "conversation（仅需对话解释/写作/推理，无上述强约束）。"
                    "判定规则：用户点名具体 App 并要求「发送/打开/点击/输入/运行」等时，一律 local_desktop；仅当出现明确检索/网页意图时才用 web_research；犹豫时优先 local_desktop。"
                    "注意：后续多 Agent 链路不会自动在磁盘生成 Word 文件；若用户要可下载文档，通常需要走动作执行(file_write)而非仅靠文本答复。"
                    "若提供了「本会话近期对话」，请结合上下文理解用户本轮补充或修改是否与上文同一任务相关。"
                    "重要：如果用户输入包含多个独立任务（如「搜索新闻并保存到 Word」包含搜索 + 文件保存两个任务），请在 multi_tasks 字段中列出每个子任务的目标描述。"
                    "multi_tasks 为数组，每个元素包含：goal(任务目标字符串), task_type(任务类型), keywords(关键词数组)。"
                ),
            },
            {
                "role": "user",
                "content": self._user_content_with_optional_vision(
                    
                        f"【本会话近期对话】\n{(dialogue_context or '').strip()}\n\n用户输入：{user_input}"
                        if (dialogue_context or "").strip()
                        else f"用户输入：{user_input}"
                    
                ),
            },
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="TaskParser")
        data = self._extract_json_object(llm_text)

        task_info = dict(fallback_task_info)
        if data:
            raw_tr = str(data.get("temporal_risk", "")).strip().lower()
            tr = raw_tr if raw_tr in ("high", "low") else self._infer_temporal_risk(user_input)
            es_raw = str(data.get("execution_surface") or "").strip().lower()
            es = es_raw if es_raw in ("local_desktop", "web_research", "conversation") else "conversation"
            task_info.update(
                {
                    "task_id": fallback_task_info["task_id"],
                    "user_input": user_input,
                    "task_type": str(data.get("task_type") or "text"),
                    "intent": str(data.get("intent") or "general"),
                    "keywords": self._normalize_keywords(data.get("keywords"))[:10],
                    "temporal_risk": tr,
                    "execution_surface": es,
                    "timestamp": time.time(),
                    "multi_tasks": data.get("multi_tasks") or [],
                }
            )
        else:
            task_info["temporal_risk"] = self._infer_temporal_risk(user_input)

        # 写入短期记忆
        self.stm.task_id = task_info["task_id"]
        self.stm.user_input = user_input
        self.stm.temporal_risk = str(task_info.get("temporal_risk") or "low")

        self.record_model_thought("TaskParser", f"任务解析完成，任务类型: {task_info['task_type']}")
        self.push_event(
            "task_parse",
            "success",
            "TaskParser",
            "PM 已完成需求解析",
            {
                "task_type": task_info.get("task_type"),
                "intent": task_info.get("intent"),
                "temporal_risk": task_info.get("temporal_risk"),
                "execution_surface": task_info.get("execution_surface"),
            },
        )
        self.push_log("TaskParser", "问题分析完成", "completed")
        self.check_cancelled("task_parse_end")
        return task_info

    # 2. 匹配本地方法论
    def match_methodology(self, task_info: dict) -> tuple[float, dict]:
        self.check_cancelled("method_match_start")
        self.push_event("method_match", "running", "MethodSearcher", "知识专家正在检索历史方法论")
        self.push_log("MethodSearcher", "正在查找历史解决方案", "running")
        # 记录模型思考过程
        self.record_model_thought("MethodSearcher", f"基于用户输入: {task_info['user_input']} 查找匹配的方法论")
        self.record_model_thought("MethodSearcher", "从长期记忆中查找方法论")
        exact = self.find_exact_methodology(task_info.get("user_input", ""))
        if exact:
            mid = str((exact or {}).get("method_id") or "")
            if mid:
                self.ltm.record_method_hit(mid, retrieval_score=1.0)
            self.record_model_thought("MethodSearcher", "命中同问句精确复用")
            self.push_event(
                "method_match",
                "success",
                "MethodSearcher",
                "命中同问句复用，跳过外网学习",
                {"score": 1.0, "exact_hit": True},
            )
            self.push_log("MethodSearcher", "命中同问句复用", "completed")
            self.check_cancelled("method_match_exact_hit")
            return 1.0, exact

        # 从长期记忆中搜索方法论：传入用户原始输入，尽量包含“场景”语义
        query = (task_info.get("user_input") or "").strip()
        results = self.ltm.search_methodology(query)

        best_match = None
        best_score = 0.0
        ab_meta: dict[str, Any] = {}
        if results:
            best_score, best_match, ab_meta = self._select_methodology_candidate(task_info, results)
            chosen_mid = str((best_match or {}).get("method_id") or "")
            if chosen_mid:
                self.ltm.record_method_hit(chosen_mid, retrieval_score=float(best_score or 0.0))
            task_id = str(task_info.get("task_id") or "")
            if task_id and isinstance(ab_meta, dict) and str(ab_meta.get("mode") or "") == "ab_bandit":
                self._ab_context_by_task[task_id] = dict(ab_meta)
            self.record_model_thought("MethodSearcher", f"成功找到匹配方法论，相似度：{best_score:.2f}")
        else:
            self.record_model_thought("MethodSearcher", "未找到匹配的方法论")

        self.record_model_thought("MethodSearcher", f"匹配完成，最佳匹配相似度：{best_score:.2f}")
        self.push_event(
            "method_match",
            "success",
            "MethodSearcher",
            f"方法论匹配完成，相似度 {best_score:.2f}",
            {"score": best_score, "ab": ab_meta},
        )
        self.push_log("MethodSearcher", f"找到匹配方案，相似度：{best_score:.2f}", "completed")
        self.check_cancelled("method_match_end")
        return best_score, best_match

    def _select_methodology_candidate(
        self,
        task_info: dict[str, Any],
        results: list[tuple[float, dict[str, Any]]],
    ) -> tuple[float, dict[str, Any] | None, dict[str, Any]]:
        """
        A/B 候选选择：
        - 默认选加权分最高的 A
        - 当 A 与 B 接近时，按“分数 + 探索奖励”做小概率探索，避免陷入单一版本
        """
        if not results:
            return 0.0, None, {}
        if len(results) == 1:
            s, m = results[0]
            return float(s), m, {"mode": "single", "chosen": "A"}

        (a_score, a_method), (b_score, b_method) = results[0], results[1]
        gap = float(a_score) - float(b_score)
        if gap >= 0.12:
            return float(a_score), a_method, {"mode": "greedy", "chosen": "A", "gap": round(gap, 4)}

        eps_raw = (os.getenv("ARIA_METHOD_AB_EXPLORATION") or "0.2").strip()
        try:
            epsilon = max(0.0, min(0.5, float(eps_raw)))
        except Exception:
            epsilon = 0.2
        epsilon = float(self.ltm.get_adaptive_ab_epsilon(epsilon))

        a_usage = int((a_method or {}).get("usage_count", 0) or 0)
        b_usage = int((b_method or {}).get("usage_count", 0) or 0)
        # 使用次数越低探索奖励越高
        a_bonus = 0.06 / (1 + a_usage)
        b_bonus = 0.06 / (1 + b_usage)
        a_bandit = float(a_score) + a_bonus
        b_bandit = float(b_score) + b_bonus

        force_explore = random.random() < epsilon
        choose_b = force_explore or (b_bandit > a_bandit)
        chosen_arm = "B" if choose_b else "A"
        chosen_score = float(b_score if choose_b else a_score)
        chosen_method = b_method if choose_b else a_method
        meta = {
            "mode": "ab_bandit",
            "chosen": chosen_arm,
            "gap": round(gap, 4),
            "epsilon": epsilon,
            "arms": {
                "A": {
                    "method_id": str((a_method or {}).get("method_id") or ""),
                    "score": round(float(a_score), 4),
                    "usage": a_usage,
                    "bandit": round(a_bandit, 4),
                },
                "B": {
                    "method_id": str((b_method or {}).get("method_id") or ""),
                    "score": round(float(b_score), 4),
                    "usage": b_usage,
                    "bandit": round(b_bandit, 4),
                },
            },
            "task_id": str(task_info.get("task_id") or ""),
        }
        return chosen_score, chosen_method, meta

    # 3. 无方案 → 调用外网大模型学习
    def learn_from_external(self, task_info: dict) -> dict:
        self.check_cancelled("method_learn_start")
        self.push_event("method_learn", "running", "SolutionLearner", "执行专家正在学习新的解决方案")
        self.push_log("SolutionLearner", "正在外网获取解决方案", "running")
        # 记录模型思考过程
        self.record_model_thought("SolutionLearner", f"调用大模型学习解决方案，用户输入: {task_info['user_input']}")
        self.record_model_thought("SolutionLearner", "正在分析问题并生成解决方案")
        fallback_solution = {
            "scene": task_info["user_input"][:50],
            "keywords": self._normalize_keywords(task_info.get("keywords") or task_info["user_input"].split()[:5]),
            "solve_steps": ["分析问题", "查找资料", "生成解决方案", "验证结果"],
            "applicable_range": "通用",
        }
        temporal = str(task_info.get("temporal_risk") or "low").strip().lower()
        messages = [
            {
                "role": "system",
                "content": (
                    "你是ARIA方案学习器。请把用户需求抽象成一个可复用的方法论(methodology)。只输出严格JSON，不要多余文本。"
                    "JSON字段: scene(字符串), keywords(字符串数组), solve_steps(字符串数组), applicable_range(字符串，可选), outcome_type(字符串，可选: stable|time_bound|pure_procedure)。"
                    "solve_steps 只写「做什么、用什么渠道/工具、如何校验」，用占位符表示地点/日期/品种等变量；禁止把具体实时数值（如今日气温、当前股价）写进步骤当作可复用事实。"
                    "若任务强依赖当下数据（天气、赛果、股价等），outcome_type 须为 time_bound，scene 与 keywords 尽量泛化（如「查询某城市天气」而非绑定单一城市名），便于同类问题复用流程。"
                    "非时效类可省略 outcome_type 或填 stable；纯流程无结论可填 pure_procedure。"
                    "solve_steps 中应体现可替换工具与本地文件落盘思路（如多种编辑器或格式），避免默认绑定单一专有软件。"
                ),
            },
            {
                "role": "user",
                "content": self._user_content_with_optional_vision(
                    f"任务时效(temporal_risk)：{temporal}（high 表示结论会过期，步骤必须是可重复的获取与校验流程）。\n\n"
                    f"用户需求：{task_info.get('user_input', '')}\n\n"
                    "请给出：1) scene，2) keywords，3) solve_steps（4-8条），4) applicable_range，5) 必要时 outcome_type。"
                ),
            },
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="SolutionLearner")
        data = self._extract_json_object(llm_text)
        solution = data if data else fallback_solution
        solution["is_success"] = False
        ot = str(solution.get("outcome_type") or "").strip().lower()
        if ot not in ("stable", "time_bound", "pure_procedure"):
            solution.pop("outcome_type", None)
        if temporal == "high":
            solution.setdefault("outcome_type", "time_bound")
        self.push_event(
            "method_learn",
            "success",
            "SolutionLearner",
            "执行专家已输出候选方法论",
            {"steps_count": len(solution.get("solve_steps", []))},
        )
        self.record_model_thought("SolutionLearner", f"学习完成，生成解决方案: {solution['solve_steps']}")
        self.push_log("SolutionLearner", "学习完成，生成解决方案", "completed")
        self.check_cancelled("method_learn_end")
        return solution

    # 4. 拆分子任务
    def split_sub_tasks(self, task_info: dict, method: dict) -> list:
        self.check_cancelled("task_split_start")
        self.push_event("task_split", "running", "TaskSplitter", "PM 正在拆分子任务")
        self.push_log("TaskSplitter", "正在拆分任务", "running")

        # 多任务拆解：如果检测到多个独立任务，为每个任务生成执行步骤
        multi_tasks = task_info.get("multi_tasks") or []
        if len(multi_tasks) > 1:
            all_sub_tasks = []
            for idx, mt in enumerate(multi_tasks, 1):
                goal = mt.get("goal", "")
                mt_type = mt.get("task_type", "text")
                mt_keywords = mt.get("keywords", [])
                self.push_log("TaskSplitter", f"任务{idx}/{len(multi_tasks)}: {goal}", "running")
                # 为每个子任务调用 LLM 生成执行步骤
                step_messages = [
                    {
                        "role": "system",
                        "content": (
                            "你是 ARIA 任务拆分器。请根据任务目标生成 2-4 个可执行的子任务步骤。只输出严格 JSON。"
                            'JSON 格式：{"sub_tasks": [{"step": "步骤 1", "description": "描述", "agent_type": "TextExecAgent"}, ...]}'
                            "agent_type 只能是 TextExecAgent / VisionExecAgent / SpeechExecAgent。"
                        ),
                    },
                    {
                        "role": "user",
                        "content": f"任务目标：{goal}\n任务类型：{mt_type}\n关键词：{', '.join(mt_keywords)}",
                    },
                ]
                step_llm = self._call_llm(step_messages, fallback_text="", agent_code="TaskSplitter")
                step_data = self._extract_json_object(step_llm)
                steps_list = step_data.get("sub_tasks") if isinstance(step_data.get("sub_tasks"), list) else None
                if not steps_list:
                    steps_list = [{"step": f"执行{goal}", "description": goal, "agent_type": "TextExecAgent"}]
                for step_item in steps_list:
                    step = str(step_item.get("step") or "")
                    atype = str(step_item.get("agent_type") or "TextExecAgent")
                    desc = str(step_item.get("description") or f"执行{step}")
                    pb = str(step_item.get("persona_brief") or "").strip() or self._default_persona_brief(
                        atype, step, desc
                    )
                    all_sub_tasks.append(
                        {
                            "sub_task_id": str(uuid.uuid4()),
                            "task_id": task_info["task_id"],
                            "task_index": idx,
                            "step": step,
                            "description": desc,
                            "agent_type": atype,
                            "persona_brief": pb,
                            "goal": goal,
                        }
                    )
            self.record_model_thought(
                "TaskSplitter", f"多任务拆解完成，共{len(multi_tasks)}个任务，{len(all_sub_tasks)}个子任务"
            )
            self.push_event(
                "task_split",
                "success",
                "TaskSplitter",
                f"任务拆分完成，共{len(all_sub_tasks)}个子任务",
                {"sub_tasks": all_sub_tasks},
            )
            self.push_log("TaskSplitter", f"任务拆分完成，共{len(all_sub_tasks)}个子任务", "completed")
            self.check_cancelled("task_split_end")
            return all_sub_tasks

        # 记录模型思考过程
        self.record_model_thought("TaskSplitter", f"基于方法论拆分子任务，用户输入: {task_info['user_input']}")
        self.record_model_thought("TaskSplitter", f"方法论步骤: {method.get('solve_steps', [])}")
        steps = method.get("solve_steps", []) or []
        fallback_sub_tasks: list[dict[str, Any]] = []
        for step in steps:
            st = str(step)
            desc_fb = f"执行{st}"
            at_fb = "TextExecAgent"
            fallback_sub_tasks.append(
                {
                    "sub_task_id": str(uuid.uuid4()),
                    "task_id": task_info["task_id"],
                    "step": st,
                    "description": desc_fb,
                    "agent_type": at_fb,
                    "persona_brief": self._default_persona_brief(at_fb, st, desc_fb),
                }
            )

        messages = [
            {
                "role": "system",
                "content": (
                    "你是ARIA总指挥（项目经理），负责拆分子任务并为每个执行位写清「个性化人设与交付要求」。"
                    "把 solve_steps 拆成可执行子任务。只输出严格JSON，不要多余文本。"
                    "字段：sub_tasks(数组)。每项必须包含："
                    "step(字符串), description(字符串), "
                    "agent_type(只能是 TextExecAgent / VisionExecAgent / SpeechExecAgent), "
                    "persona_brief(字符串，2-6句中文)：针对该子任务写清角色人设、专业侧重、禁忌与可验收的交付标准，"
                    "须与 agent_type 一致（文本/视觉/语音侧重不同）。"
                ),
            },
            {
                "role": "user",
                "content": self._user_content_with_optional_vision(
                    f"任务原始输入：{task_info.get('user_input', '')}\n\n方法论场景(scene)：{method.get('scene', '') or method.get('scenario', '')}\n\nsolve_steps：{steps}"
                ),
            },
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="TaskSplitter")
        data = self._extract_json_object(llm_text)
        sub_tasks_data = data.get("sub_tasks") if isinstance(data.get("sub_tasks"), list) else None

        if not sub_tasks_data:
            sub_tasks = fallback_sub_tasks
        else:
            sub_tasks = []
            for item in sub_tasks_data:
                step = str(item.get("step") or "")
                atype = str(item.get("agent_type") or "TextExecAgent")
                desc = str(item.get("description") or f"执行{step}")
                pb = str(item.get("persona_brief") or "").strip() or self._default_persona_brief(atype, step, desc)
                sub_tasks.append(
                    {
                        "sub_task_id": str(uuid.uuid4()),
                        "task_id": task_info["task_id"],
                        "step": step,
                        "description": desc,
                        "agent_type": atype,
                        "persona_brief": pb,
                    }
                )

        # 写入短期记忆
        self.stm.sub_tasks = sub_tasks

        self.record_model_thought("TaskSplitter", f"任务拆分完成，共{len(sub_tasks)}个子任务")
        self.push_event(
            "task_split",
            "success",
            "TaskSplitter",
            f"任务拆分完成，共 {len(sub_tasks)} 个子任务",
            {"sub_tasks": sub_tasks},
        )
        self.push_log("TaskSplitter", f"拆分完成，共{len(sub_tasks)}个子任务", "completed")
        self.check_cancelled("task_split_end")
        return sub_tasks

    # 5. 动态生成Agent
    def choose_collaboration_topology(self, task_info: dict, sub_tasks: list) -> str:
        """根据任务复杂度和子任务数量决定协作拓扑：pipeline 或 parallel_with_merge。"""
        complexity = int((task_info or {}).get("complexity_score") or 0)
        if len(sub_tasks) >= 3 or complexity >= 5:
            return "parallel_with_merge"
        return "pipeline"

    def create_agents(self, sub_tasks: list) -> dict:
        self.check_cancelled("agent_create_start")
        self.push_event("agent_create", "running", "TaskSplitter", "PM 正在组建执行小队")
        self.push_log("TaskSplitter", "正在生成执行Agent", "running")
        # 记录模型思考过程
        self.record_model_thought("TaskSplitter", f"开始生成Agent，共{len(sub_tasks)}个子任务")
        time.sleep(0.5)  # 添加延迟，使日志显示更加流畅
        agents = {}
        agent_status = {}
        used_names_by_type: dict[str, set[str]] = {}
        for sub_task in sub_tasks:
            agent_id = str(uuid.uuid4())
            agent_type = sub_task["agent_type"]
            assigned_name = self._pick_exec_agent_name(agent_type, used_names_by_type)
            persona = str(sub_task.get("persona_brief") or "").strip() or self._default_persona_brief(
                agent_type,
                str(sub_task.get("step") or ""),
                str(sub_task.get("description") or ""),
            )
            # 这里应该实例化具体的Agent类，现在只是模拟
            agents[agent_id] = {
                "agent_id": agent_id,
                "agent_type": agent_type,
                "task": sub_task,
                "persona_brief": persona,
                "agent_name": assigned_name,
            }
            agent_status[agent_id] = "created"
            self.record_model_thought(
                "TaskSplitter",
                f"生成Agent: {agent_type}({assigned_name})，统一模型: {self.unified_model}",
            )

        # 写入短期记忆
        self.stm.agent_status = agent_status

        self.record_model_thought("TaskSplitter", f"Agent生成完成，共{len(agents)}个Agent")
        self.push_event(
            "agent_create",
            "success",
            "TaskSplitter",
            f"执行小队组建完成，共 {len(agents)} 名成员",
        )
        self.push_log("TaskSplitter", f"生成完成，共{len(agents)}个Agent", "completed")
        self.check_cancelled("agent_create_end")
        return agents

    # 6. 执行Agent
    def run_agents(
        self,
        agents: dict,
        method: dict[str, Any] | None = None,
        dialogue_context: str = "",
    ) -> list:
        self.check_cancelled("agent_execute_start")
        results: list[dict[str, Any]] = []
        method_ctx = self._methodology_summary_text(method) if method else ""
        for agent_id, agent in agents.items():
            self.check_cancelled("agent_execute_loop")
            agent_type = agent["agent_type"]
            task = agent["task"]
            agent_name = agent.get("agent_name") or self._agent_profile(agent_type)["name"]
            role = self._agent_profile(agent_type)["role"]
            persona = str(agent.get("persona_brief") or "").strip() or self._default_persona_brief(
                agent_type,
                str(task.get("step") or ""),
                str(task.get("description") or ""),
            )
            self.push_event(
                "agent_execute",
                "running",
                agent_type,
                f"{role} 正在执行：{task['description']}",
                {"sub_task_id": task.get("sub_task_id")},
                agent_name_override=agent_name,
            )
            self.push_log(agent_type, f"正在执行子任务：{task['description']}", "running")
            # 记录模型思考过程
            self.record_model_thought(agent_type, f"开始执行任务：{task['description']}")
            self.record_model_thought(agent_type, f"统一模型：{self.unified_model}")

            sys_parts = [
                f"你是ARIA执行专家[{agent_type}]。你将基于给定步骤产出可直接使用的结果。只输出纯文本，不要JSON。",
                _MATH_NOTATION_FOR_CHAT,
            ]
            if str(getattr(self.stm, "temporal_risk", "low")).lower() == "high":
                sys_parts.append(
                    "【时效】本任务结论依赖当下数据：须说明信息时间点或获取渠道，禁止把方法论纲要中的示例数值当作当前事实；"
                    "无法取得实时数据时要明确说明并给出用户可自行核实的方式。"
                )
            sys_parts.extend(
                [
                    "本链路仅为文本推理：你没有调用 file_write、没有访问用户磁盘。严禁声称「已成功创建/保存 .docx」「已写入 我的文档/此电脑>文档」等；若用户要可下载文件，应明确说明须由用户在动作计划中「确认执行」file_write 到工作区，或自行在本机用 Word 另存。",
                    "用户可通过网页回形针上传文件；若子任务涉及已上传文件，正文可能在「原始任务输入」的附件摘要中。不要编造用户未提供的文件内容。",
                    "涉及「创建文档/保存文件」时：不要因未安装 Word 就拒绝；应给出可落地方案——例如建议相对路径如 data/artifacts/xxx.md、记事本或 WPS/LibreOffice/VS Code 等替代、以及可复制粘贴的正文草稿。",
                    "若仍缺关键信息（路径、格式、是否覆盖），在答复末尾用简短编号列出 1～3 个需用户确认的问题。",
                    "【总指挥设定的人设与要求】",
                    persona,
                ]
            )
            user_parts = []
            if (dialogue_context or "").strip():
                user_parts.append(f"【本会话近期对话（与当前任务同一线程）】\n{(dialogue_context or '').strip()}")
            if method_ctx:
                user_parts.append(method_ctx)
            previous_results_text = self._build_exec_context_window(results)
            user_parts.extend(
                [
                    f"原始任务输入：{self.stm.user_input}",
                    f"当前子任务步骤(step)：{task.get('step', '')}",
                    f"子任务描述(description)：{task.get('description', '')}",
                    "【此前各执行者的完整产出（请完整理解，勿遗漏细节；执行链为纯文本传递，无二次解析）】",
                    previous_results_text if previous_results_text else "（尚无）",
                    "",
                    "请仅输出本步骤的最终结果正文。",
                ]
            )
            user_body = "\n".join(user_parts)
            # 多模态体积大：仅首个子任务附带原图，后续步骤依赖前文文本链
            user_content = self._user_content_with_optional_vision(user_body) if not results else user_body
            messages = [
                {"role": "system", "content": "\n".join(sys_parts)},
                {"role": "user", "content": user_content},
            ]
            llm_text = self._call_llm(messages, fallback_text=f"执行完成：{task['description']}", agent_code=agent_type)

            result = {
                "agent_id": agent_id,
                "agent_type": agent_type,
                "agent_name": agent_name,
                "step": task.get("step", ""),
                "description": task.get("description", ""),
                "task_id": task["task_id"],
                "sub_task_id": task["sub_task_id"],
                "result": llm_text.strip(),
                "status": "completed",
                "timestamp": time.time(),
            }
            # 更新短期记忆中的Agent状态
            self.stm.agent_status[agent_id] = "completed"
            self.record_model_thought(agent_type, f"任务执行完成：{result['result']}")
            results.append(result)
            self.push_log(agent_type, "执行完成", "completed")
            self.push_event(
                "agent_execute",
                "success",
                agent_type,
                f"{role} 已完成：{task['step']}",
                {
                    "sub_task_id": task.get("sub_task_id"),
                    "result_preview": (result["result"] or "")[:120],
                },
                agent_name_override=agent_name,
            )

        # 写入短期记忆
        self.stm.results = results

        self.check_cancelled("agent_execute_end")
        return results

    def _build_exec_context_window(self, results: list[dict[str, Any]]) -> str:
        """子 Agent 隔离：只透传最近若干步骤，并限制总字符，避免噪声在链路中累积。"""
        if not results:
            return ""
        try:
            max_steps = max(1, int(os.getenv("ARIA_AGENT_CONTEXT_MAX_STEPS", "4") or "4"))
        except (TypeError, ValueError):
            max_steps = 4
        try:
            max_chars = max(1000, int(os.getenv("ARIA_AGENT_CONTEXT_MAX_CHARS", "5000") or "5000"))
        except (TypeError, ValueError):
            max_chars = 5000
        window = results[-max_steps:]
        text = self._format_exec_results_as_plain_text(window)
        if len(text) <= max_chars:
            return text
        return "…[历史执行上下文已截断]\n" + text[-max_chars:]

    # 7. 校验结果
    def check_result(self, results: list) -> dict:
        self.check_cancelled("quality_check_start")
        self.push_event("quality_check", "running", "QualityChecker", "QA 正在校验并汇总结果")
        self.push_log("QualityChecker", "正在校验结果", "running")
        # 记录模型思考过程
        self.record_model_thought("QualityChecker", f"开始校验结果，共{len(results)}个结果")
        for i, result in enumerate(results):
            self.record_model_thought("QualityChecker", f"校验第{i + 1}个结果：{result['result']}")

        fallback_final = "\n".join([result["result"] for result in results])
        steps_plain = self._format_exec_results_as_plain_text(results if isinstance(results, list) else [])
        qc_sys_parts = [
            "你是ARIA质量校验员。请基于各子步骤产出生成最终结果，并判断整体是否符合要求。只输出严格JSON，不要多余文本。"
            "JSON字段：final_result(字符串), is_success(布尔)。"
            "final_result 须汇总可执行建议；若任务与本地文件/软件相关，应包含多种替代路径（如 .md/.txt、常见编辑器），避免单独一句「无法完成」。"
            "若子步骤声称已生成 .docx 或已保存到「文档」文件夹等，但整个流程未实际执行 file_write 动作，你必须在 final_result 中删除或纠正此类虚假陈述，并说明：真实文件需用户确认 file_write 计划，或自行在本机创建。"
            "Web 对话不能发附件；不要在 final_result 中让用户「通过聊天窗发送文件」。"
            "若信息仍不足，在 final_result 中明确写出 1～3 条需用户补充或选择的问题。",
            _MATH_NOTATION_FOR_CHAT,
        ]
        if str(getattr(self.stm, "temporal_risk", "low")).lower() == "high":
            qc_sys_parts.append(
                "本任务强时效：final_result 不得把过时或示例数值当作当前事实；应写清数据时间点或引导用户通过权威渠道核实最新信息。"
            )
        messages = [
            {
                "role": "system",
                "content": "\n".join(qc_sys_parts),
            },
            {
                "role": "user",
                "content": (
                    f"原始任务输入：{self.stm.user_input}\n\n"
                    "下列为各执行者完整产出（纯文本分块，与执行链传递一致，请据此汇总勿臆造）：\n"
                    f"{steps_plain or '（无）'}"
                ),
            },
        ]
        llm_text = self._call_llm(messages, fallback_text="", agent_code="QualityChecker")
        data = self._extract_json_object(llm_text)
        if not data:
            final_result = fallback_final
            is_success = False
        else:
            final_result = str(data.get("final_result") or fallback_final)
            is_success = bool(data.get("is_success", False))

        self.record_model_thought("QualityChecker", "结果校验完成，生成最终结果")
        self.push_event(
            "quality_check",
            "success",
            "QualityChecker",
            "QA 已完成校验与汇总",
            {"is_success": is_success},
        )
        self.push_log("QualityChecker", "结果校验完成", "completed")
        self.check_cancelled("quality_check_end")
        return {"final_result": final_result, "is_success": is_success}

    # 8. 保存方法论
    def _record_method_feedback(self, task_info: dict, method: dict, result_payload: Any) -> None:
        task_id = str(task_info.get("task_id") or "")
        method_id = str((method or {}).get("method_id") or "")
        is_success = bool(result_payload.get("is_success", False)) if isinstance(result_payload, dict) else False
        quality_score = 1.0 if is_success else 0.0
        if isinstance(result_payload, dict):
            try:
                quality_score = float(result_payload.get("score", quality_score) or quality_score)
            except Exception:
                quality_score = 1.0 if is_success else 0.0
        if method_id:
            self.ltm.record_method_outcome(method_id, is_success=is_success, quality_score=quality_score)
        ab_meta = self._ab_context_by_task.pop(task_id, None) if task_id else None
        if isinstance(ab_meta, dict):
            self.ltm.record_ab_outcome(ab_meta, is_success=is_success, quality_score=quality_score)

    def save_methodology(self, task_info: dict, method: dict, result_payload: Any):
        self.check_cancelled("method_save_start")
        should_save, skip_reason, judge_source = self._should_save_methodology(task_info, method or {}, result_payload)
        if not should_save:
            self._record_method_feedback(task_info, method or {}, result_payload)
            self.record_model_thought("MethodSaver", f"知识沉淀判定为跳过，reason={skip_reason}, source={judge_source}")
            self.push_event(
                "method_save",
                "success",
                "MethodSaver",
                "本轮对话暂不满足沉淀条件，已跳过方法论入库",
                {"skipped": True, "reason": skip_reason, "judge_source": judge_source},
            )
            self.push_log("MethodSaver", f"已跳过方法论保存（{skip_reason}）", "completed")
            self.check_cancelled("method_save_skipped")
            return

        self.push_event("method_save", "running", "MethodSaver", "知识专家正在沉淀方法论")
        self.push_log("MethodSaver", "正在保存方案", "running")
        # 记录模型思考过程
        self.record_model_thought("MethodSaver", f"开始保存方法论，用户输入: {task_info['user_input']}")
        self.record_model_thought("MethodSaver", "将方法论保存到长期记忆")

        is_success = False
        if isinstance(result_payload, dict):
            is_success = bool(result_payload.get("is_success", False))
        else:
            is_success = False

        normalized_method = self._normalize_methodology(method or {}, task_info)
        normalized_method["success_count"] = 1 if is_success else 0
        normalized_method["is_success"] = is_success
        if isinstance(result_payload, dict):
            quality_metrics = result_payload.get("quality_metrics")
            if isinstance(quality_metrics, dict):
                normalized_method["quality_metrics"] = quality_metrics
            try:
                normalized_method["score"] = float(
                    result_payload.get("score", normalized_method.get("score", 0.0)) or 0.0
                )
            except Exception:
                normalized_method["score"] = float(normalized_method.get("score", 0.0) or 0.0)
        normalized_method.setdefault("evidence_refs", [])
        normalized_method["evidence_refs"] = list(normalized_method.get("evidence_refs") or []) + [
            f"task_id:{task_info.get('task_id', '')}".strip(":"),
        ]

        self.ltm.add_methodology(normalized_method)
        self._record_method_feedback(task_info, normalized_method, result_payload)

        # 保存到中期记忆作为模板
        task_template = {
            "template_id": str(uuid.uuid4()),
            "task_type": task_info["task_type"],
            "solve_steps": normalized_method["solve_steps"],
            "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        }
        self.mtm.task_templates.append(task_template)
        self.mtm.save()

        self.record_model_thought("MethodSaver", "方法论保存完成")
        self.push_event("method_save", "success", "MethodSaver", "方法论沉淀完成")
        self.push_log("MethodSaver", "方案保存完成", "completed")
        self.check_cancelled("method_save_end")

    # 9. 销毁所有Agent
    def destroy_agents(self, agents: dict):
        self.push_log("TaskSplitter", "正在销毁Agent", "running")
        # 记录模型思考过程
        self.record_model_thought("TaskSplitter", f"开始销毁Agent，共{len(agents)}个Agent")
        time.sleep(0.3)  # 添加延迟，使日志显示更加流畅
        # 实际应该调用Agent的销毁方法，释放资源
        agents.clear()
        # 任务结束，清空短期记忆
        self.stm.clear()
        self.record_model_thought("TaskSplitter", "Agent销毁完成，短期记忆已清空")
        self.push_event("agent_destroy", "success", "TaskSplitter", "小队任务结束，已释放资源")
        self.push_log("TaskSplitter", "Agent销毁完成", "completed")

    # 10. 推送日志到UI
    def push_log(self, agent_name: str, content: str, status="running"):
        log_entry = {"agent": agent_name, "content": content, "status": status, "timestamp": time.time()}
        self.execution_log.append(log_entry)
        # 同时写入短期记忆
        self.stm.logs.append(log_entry)

    # 获取执行日志
    def get_execution_log(self):
        return self.execution_log

    # 清空执行日志
    def clear_execution_log(self):
        self.execution_log = []

    # ------------------------------------------------------------------ #
    # runtime/ 依赖的桩方法                                                  #
    # ------------------------------------------------------------------ #

    def _math_notation_hint(self) -> str:
        return _MATH_NOTATION_FOR_CHAT

    def _memory_system_prompt_fragment(self, task_text: str = "") -> str:
        try:
            auto_memory = getattr(self, "auto_memory", None)
            if auto_memory and hasattr(auto_memory, "get_system_prompt_fragment"):
                return auto_memory.get_system_prompt_fragment(task_text)
        except Exception:
            pass
        return ""

    SAFE_ACTION_TYPES = SAFE_ACTION_TYPES  # single source: runtime.permissions

    def _build_taor_task_info(self, user_input: str, pa_result: dict) -> dict:
        """为 TAOR 模式构造轻量 task_info，复用 pa_result 已有字段，避免额外 LLM 调用。"""
        task_id = self._resolve_task_id_for_turn(None)
        self.current_task_id = task_id
        return {
            "task_id": task_id,
            "user_input": user_input,
            "task_type": str(pa_result.get("task_form") or "text"),
            "intent": str(pa_result.get("summary") or "general"),
            "keywords": user_input.split()[:5],
            "temporal_risk": str(
                pa_result.get("temporal_risk") or self._infer_temporal_risk(user_input)
            ),
            "outcome_type": str(pa_result.get("outcome_type") or "stable"),
            "complexity_score": int(pa_result.get("complexity_score") or 3),
            "execution_surface": str(pa_result.get("task_form") or "conversation"),
            "timestamp": time.time(),
        }

    def _build_taor_method_from_result(
        self,
        matched_method: dict | None,
        taor_result: dict,
        task_info: dict,
    ) -> dict:
        """从 TAOR 执行结果构造待保存的方法论 dict。"""
        tool_trace = taor_result.get("tool_trace") or []
        solve_steps = [
            f"{i + 1}. [{e['action'].get('type', '')}] {e['action'].get('reason', '')[:200]}"
            for i, e in enumerate(tool_trace)
            if isinstance(e.get("action"), dict)
        ] or [taor_result.get("final_result", "")[:300]]

        base = dict(matched_method) if matched_method else {}
        base.update(
            {
                "scene": base.get("scene") or task_info.get("user_input", "")[:80],
                "solve_steps": solve_steps,
                "keywords": base.get("keywords") or task_info.get("keywords") or [],
                "is_success": taor_result.get("is_success", False),
                "score": 1.0 if taor_result.get("is_success") else 0.0,
                "quality_metrics": {"total_steps": len(tool_trace)},
            }
        )
        return base

    def run_taor_pipeline(
        self,
        user_input: str,
        dialogue_context: str = "",
        conversation_id: str = "",
    ) -> dict:
        from runtime.taor_loop import TAORLoop
        from runtime.hybrid_planner import HybridPlanner

        self.current_conversation_id = conversation_id

        plan_context: str = ""
        tracker = None
        pa_result: dict = {}

        if os.getenv("ARIA_HYBRID_PLAN", "0").strip().lower() in ("1", "true", "yes"):
            planner = HybridPlanner(self)
            pa_result = self.plan_actions(user_input, dialogue_context)
            if planner.should_plan(user_input, pa_result):
                self.push_event(
                    "hybrid_plan",
                    "running",
                    "HybridPlanner",
                    "生成执行路线图",
                    {"complexity_score": pa_result.get("complexity_score")},
                )
                hybrid_plan = planner.build_plan(user_input, dialogue_context, pa_result)
                plan_context = planner.format_plan_for_system_prompt(hybrid_plan)
                tracker = planner.make_tracker(hybrid_plan)
                self.push_event(
                    "hybrid_plan",
                    "success",
                    "HybridPlanner",
                    f"路线图生成完成，{len(hybrid_plan.get('sub_goals', []))} 个子目标",
                    {"plan": hybrid_plan},
                )

        # ── 方法论查找：执行前注入历史经验 ──────────────────────────────────
        matched_method: dict | None = None
        task_info: dict = {}
        _taor_method_enabled = (
            os.getenv("ARIA_TAOR_METHODOLOGY", "1").strip().lower()
            not in ("0", "false", "no")
        )
        if _taor_method_enabled:
            task_info = self._build_taor_task_info(user_input, pa_result)
            try:
                _score, matched_method = self.match_methodology(task_info)
                min_score = float(os.getenv("ARIA_TAOR_METHOD_MIN_SCORE", "0.5"))
                if matched_method and _score < min_score:
                    matched_method = None
            except Exception as e:
                if getattr(self, "_cancelled", False):
                    raise
                self.push_event(
                    "method_match",
                    "warning",
                    "MethodSearcher",
                    f"TAOR 方法论匹配失败（已跳过）: {e}",
                    {},
                )
                matched_method = None

        # ── 执行 TAOR 循环 ────────────────────────────────────────────────
        taor_result = TAORLoop(self).run(
            user_input,
            dialogue_context,
            method=matched_method,
            plan_context=plan_context,
            plan_tracker=tracker,
        )

        # ── 方法论保存：执行后沉淀成功经验 ──────────────────────────────────
        if _taor_method_enabled and task_info:
            try:
                method_to_save = self._build_taor_method_from_result(
                    matched_method, taor_result, task_info
                )
                self.save_methodology(task_info, method_to_save, taor_result)
            except Exception as e:
                if getattr(self, "_cancelled", False):
                    raise
                self.push_event(
                    "method_save",
                    "warning",
                    "MethodSaver",
                    f"TAOR 方法论保存失败（已跳过）: {e}",
                    {},
                )

        return taor_result
