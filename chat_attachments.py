"""
用户通过 Web 上传的聊天附件：保存到工作区并抽取可送入 LLM 的文本摘要。
路径均限制在 data/artifacts/uploads/<conversation_id>/ 下。
"""

from __future__ import annotations

import base64
import json
import re
import uuid
from pathlib import Path
from typing import Any, TYPE_CHECKING

if TYPE_CHECKING:
    from aria_manager import ARIAManager

# 上传限制
MAX_UPLOAD_BYTES = 15 * 1024 * 1024
MAX_FILES_PER_MESSAGE = 8
MAX_EXTRACT_TOTAL_CHARS = 18_000
MAX_TEXT_FILE_READ = 120_000

# 送入 chat.completions 的 vision 片段（data URL）；单张过大或过多会跳过以免 413/超时
_VISION_MIME = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
}
MAX_VISION_IMAGES_PER_TURN = 6
MAX_VISION_BYTES_PER_IMAGE = 4 * 1024 * 1024

ALLOWED_EXTENSIONS = frozenset(
    {
        "txt",
        "md",
        "csv",
        "json",
        "log",
        "py",
        "html",
        "htm",
        "xml",
        "png",
        "jpg",
        "jpeg",
        "gif",
        "webp",
        "pdf",
        "docx",
        "xlsx",
        "xlsm",
        "pptx",
    }
)


def _ext(name: str) -> str:
    return Path(name or "").suffix.lower().lstrip(".")


def _safe_filename(name: str) -> str:
    base = Path(name or "file").name
    base = re.sub(r"[^\w\u4e00-\u9fff.\-()+]", "_", base)
    return (base or "file")[:180]


def uploads_dir_for_conversation(conversation_id: str) -> str:
    cid = (conversation_id or "").strip().replace("\\", "/").replace("..", "")
    return f"data/artifacts/uploads/{cid}"


def save_uploaded_file(manager: ARIAManager, conversation_id: str, storage) -> dict[str, Any]:
    """
    storage: werkzeug FileStorage
    返回 { path, name, size, ext }，path 为相对工作区路径；失败抛 ValueError。
    """
    if not storage or not getattr(storage, "filename", None):
        raise ValueError("empty_file")
    ext = _ext(storage.filename)
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"unsupported_type:{ext}")
    cl_raw = getattr(storage, "content_length", None)
    try:
        cl = int(cl_raw) if cl_raw else 0
    except (TypeError, ValueError):
        cl = 0
    if cl and cl > MAX_UPLOAD_BYTES:
        raise ValueError("file_too_large")

    rel_dir = uploads_dir_for_conversation(conversation_id)
    short = uuid.uuid4().hex[:10]
    fname = f"{short}_{_safe_filename(storage.filename)}"
    rel_path = f"{rel_dir}/{fname}".replace("\\", "/")

    full = manager._ensure_safe_path(rel_path)
    full.parent.mkdir(parents=True, exist_ok=True)
    storage.save(str(full))
    size = full.stat().st_size
    if size > MAX_UPLOAD_BYTES:
        try:
            full.unlink()
        except OSError:
            pass
        raise ValueError("file_too_large")
    if size <= 0:
        try:
            full.unlink()
        except OSError:
            pass
        raise ValueError("empty_file")
    return {"path": rel_path, "name": storage.filename, "size": size, "ext": ext}


def validate_attachment_path(manager: ARIAManager, conversation_id: str, rel_path: str) -> str:
    """校验 JSON 回传的 path 必须落在本会话的上传目录下。"""
    p = (rel_path or "").strip().replace("\\", "/")
    prefix = uploads_dir_for_conversation(conversation_id) + "/"
    if not p.startswith(prefix) or ".." in p:
        raise ValueError("invalid_attachment_path")
    full = manager._ensure_safe_path(p)
    if not full.is_file():
        raise ValueError("attachment_not_found")
    return p


def extract_llm_excerpt(manager: ARIAManager, records: list[dict[str, Any]], budget: int = MAX_EXTRACT_TOTAL_CHARS) -> str:
    """根据已保存的附件路径抽取文本，供拼入 user 侧提示。"""
    parts: list[str] = []
    used = 0
    for rec in records:
        if used >= budget:
            parts.append("\n... [附件摘要总长度已达上限]")
            break
        rel = str(rec.get("path") or "").strip()
        if not rel:
            continue
        name = str(rec.get("name") or Path(rel).name)
        ext = str(rec.get("ext") or _ext(name)).lower()
        try:
            full = manager._ensure_safe_path(rel)
        except Exception:
            parts.append(f"- {name}: （路径无效，已跳过）")
            continue
        if not full.is_file():
            parts.append(f"- {name}: （文件不存在）")
            continue

        chunk = _extract_one_file(full, name, ext)
        line = f"### 文件: {name} ({ext or 'unknown'})\n{chunk}".strip()
        if used + len(line) > budget:
            line = line[: max(0, budget - used - 40)] + "\n... [截断]"
        parts.append(line)
        used += len(line) + 2

    return "\n\n".join(parts).strip()


def image_data_urls_from_attachment_records(
    manager: ARIAManager,
    records: list[dict[str, Any]],
    *,
    max_images: int = MAX_VISION_IMAGES_PER_TURN,
    max_bytes: int = MAX_VISION_BYTES_PER_IMAGE,
) -> list[str]:
    """从已落盘的上传图片生成 OpenAI 兼容的 data:image/...;base64,... URL 列表。"""
    out: list[str] = []
    for rec in records:
        if len(out) >= max_images:
            break
        rel = str(rec.get("path") or "").strip()
        ext = (str(rec.get("ext") or "").lower() or _ext(str(rec.get("name") or rel))).lower()
        if ext not in _VISION_MIME:
            continue
        try:
            full = manager._ensure_safe_path(rel)
        except Exception:
            continue
        if not full.is_file():
            continue
        try:
            raw = full.read_bytes()
        except OSError:
            continue
        if len(raw) > max_bytes:
            continue
        b64 = base64.b64encode(raw).decode("ascii")
        out.append(f"data:{_VISION_MIME[ext]};base64,{b64}")
    return out


def _extract_one_file(full: Path, display_name: str, ext: str) -> str:
    ext = ext.lower()
    if ext in ("png", "jpg", "jpeg", "gif", "webp"):
        try:
            sz = full.stat().st_size
        except OSError:
            sz = 0
        return (
            f"（图片文件，约 {sz // 1024} KB；若接入点多模态且未超大小限制，同轮会一并传入模型。"
            f"相对路径: {full.name}）"
        )

    if ext == "pdf":
        return _try_pdf_text(full)

    if ext == "docx":
        return _try_docx_text(full)

    if ext in ("xlsx", "xlsm"):
        return _try_xlsx_text(full)

    if ext == "pptx":
        return _try_pptx_text(full)

    if ext in ("txt", "md", "csv", "json", "log", "py", "html", "htm", "xml"):
        return _read_plain_text(full)

    return "（未实现该类型的文本抽取）"


def _read_plain_text(full: Path) -> str:
    raw = full.read_bytes()[:MAX_TEXT_FILE_READ]
    for enc in ("utf-8", "utf-8-sig", "gbk", "latin-1"):
        try:
            t = raw.decode(enc)
            break
        except UnicodeDecodeError:
            t = ""
    else:
        t = raw.decode("utf-8", errors="replace")
    t = t.strip()
    if len(t) > 14_000:
        t = t[:14_000] + "\n... [文本已截断]"
    return t or "（空文件）"


def _try_pdf_text(full: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        return "（未安装 pypdf，无法抽取 PDF 文本。可执行: pip install pypdf）"
    try:
        reader = PdfReader(str(full))
        chunks: list[str] = []
        for i, page in enumerate(reader.pages[:30]):
            try:
                chunks.append(page.extract_text() or "")
            except Exception:
                pass
        t = "\n".join(chunks).strip()
        if len(t) > 14_000:
            t = t[:14_000] + "\n... [PDF 文本已截断]"
        return t or "（PDF 中未解析到可读文本，可能为扫描件）"
    except Exception as e:
        return f"（PDF 读取失败: {e}）"


def _try_docx_text(full: Path) -> str:
    try:
        from docx import Document  # type: ignore
    except Exception:
        return "（未安装 python-docx）"
    try:
        doc = Document(str(full))
        paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        t = "\n".join(paras).strip()
        if len(t) > 14_000:
            t = t[:14_000] + "\n... [docx 已截断]"
        return t or "（docx 中无段落文本）"
    except Exception as e:
        return f"（docx 读取失败: {e}）"


def _try_xlsx_text(full: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception:
        return "（未安装 openpyxl）"
    try:
        wb = load_workbook(str(full), read_only=True, data_only=True)
        ws = wb.active
        rows_out: list[str] = []
        for i, row in enumerate(ws.iter_rows(max_row=80, max_col=24, values_only=True)):
            cells = [("" if c is None else str(c)) for c in row]
            if any(x.strip() for x in cells):
                rows_out.append("\t".join(cells))
        wb.close()
        t = "\n".join(rows_out).strip()
        if len(t) > 14_000:
            t = t[:14_000] + "\n... [表格已截断]"
        return t or "（表格为空）"
    except Exception as e:
        return f"（xlsx 读取失败: {e}）"


def _try_pptx_text(full: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return "（未安装 python-pptx）"
    try:
        prs = Presentation(str(full))
        chunks: list[str] = []
        for slide in prs.slides[:12]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text.strip())
        t = "\n".join(chunks).strip()
        if len(t) > 14_000:
            t = t[:14_000] + "\n... [pptx 已截断]"
        return t or "（幻灯片中未解析到文本）"
    except Exception as e:
        return f"（pptx 读取失败: {e}）"


def merge_json_attachments(
    manager: ARIAManager,
    conversation_id: str,
    items: Any,
) -> list[dict[str, Any]]:
    """解析前端 JSON 中的 attachments: [{path, name?}]"""
    if not items:
        return []
    if not isinstance(items, list):
        return []
    out: list[dict[str, Any]] = []
    for it in items[:MAX_FILES_PER_MESSAGE]:
        if not isinstance(it, dict):
            continue
        p = validate_attachment_path(manager, conversation_id, str(it.get("path") or ""))
        ext = _ext(p)
        out.append(
            {
                "path": p,
                "name": str(it.get("name") or Path(p).name),
                "ext": ext,
                "size": Path(manager._ensure_safe_path(p)).stat().st_size,
            }
        )
    return out
